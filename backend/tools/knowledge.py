"""Knowledge Base Tool – Multi-Folder RAG mit Vektor-Suche (ChromaDB) und TF-IDF Fallback."""

import asyncio
import json
import logging
import math
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
from collections import Counter
from pathlib import Path

from backend.tools.base import BaseTool
from backend.config import config

PROJECT_ROOT = Path(__file__).parent.parent.parent
INDEX_CACHE_PATH = PROJECT_ROOT / "data" / "knowledge_index.json"
DEFAULT_FOLDER = "data/knowledge"
# Groessenlimit je Datei. Von 50 auf 150 angehoben (2026-08-01), nachdem auf
# ECHT ein 130-MB-PDF (KBV-Bewertungsmassstab) als einziger ECHTER Indexfehler
# uebrig blieb.
#
# Das Anheben war erst vertretbar, NACHDEM die Speicherspitze der PDF-Extraktion
# von 6000 MB auf 306 MB gesenkt wurde (siehe _extract_text_raw). Mit dem alten
# Stand haette dieses Limit den OOM-Kill zurueckgeholt, an dem der Reindex
# historisch dreimal gestorben ist – die VM hat 16 GB und im Betrieb rund 8 GB
# frei. Wer das Limit weiter anhebt, misst vorher die Spitze am groessten
# tatsaechlich vorhandenen Dokument.
DEFAULT_MAX_SIZE_MB = 150

# Maximale Zeichen pro Treffer-Chunk in der Tool-Ausgabe.
# MUSS groesser sein als ein vollstaendiger Chunk (_chunk_text: 200 Woerter,
# ~1600 Zeichen) – sonst wird der gefundene Treffer mitten im Text abgeschnitten
# und das LLM antwortet auf einem Ausschnitt, der die Antwort gar nicht enthaelt.
CHUNK_OUTPUT_LIMIT = 3000

EXTENSIONS_TEXT = {
    ".txt", ".md", ".json", ".csv", ".log", ".py", ".sh",
    ".yaml", ".yml", ".cfg", ".conf", ".ini",
}
EXTENSIONS_PDF = {".pdf"}
EXTENSIONS_DOCX = {".docx", ".doc"}
EXTENSIONS_XLSX = {".xlsx", ".xls"}
EXTENSIONS_PPTX = {".pptx"}
EXTENSIONS_VIDEO = {".mp4", ".mkv", ".avi", ".webm", ".mov", ".m4v", ".flv", ".wmv"}
EXTENSIONS_AUDIO = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac", ".wma", ".opus"}
EXTENSIONS_IMAGE = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp"}
# OneNote-ABSCHNITTSDATEIEN. Bewusst OHNE ".onetoc2": das ist der
# Notizbuch-Index ohne eigenen Inhalt, und der Parser (Apache Tika) hat dafuer
# keinen Handler. Siehe backend/tools/onenote.py.
EXTENSIONS_ONENOTE = {".one"}


def alle_endungen() -> set:
    """DIE EINE Quelle fuer "welche Datei ist indizierbar".

    Bis 2026-09-04 stand diese Vereinigung an SECHS Stellen als Handarbeit
    (zweimal hier, dreimal in main.py, einmal in knowledge_sync.py). Wer ein
    Format ergaenzte, musste alle sechs finden – und die vergessene Stelle
    faellt nicht auf, sie laesst die Datei nur still liegen: nicht indiziert,
    nicht hochladbar oder nicht uebertragbar, je nachdem welche es war.
    ``tests/test_onenote_import.py`` haelt fest, dass niemand die Vereinigung
    wieder von Hand zusammensetzt.
    """
    return (EXTENSIONS_TEXT | EXTENSIONS_PDF | EXTENSIONS_DOCX | EXTENSIONS_XLSX
            | EXTENSIONS_PPTX | EXTENSIONS_VIDEO | EXTENSIONS_AUDIO
            | EXTENSIONS_IMAGE | EXTENSIONS_ONENOTE)

_cache_lock = threading.Lock()
_log = logging.getLogger("jarvis.knowledge")

# ─── Indizierungs-Fortschritt (thread-sicher) ────────────────────────────────
# started_at/finished_at sind Unix-Zeitstempel (float) – die Oberflaeche zeigt
# damit Startzeit und Laufdauer an.
_index_progress: dict = {"running": False, "phase": "", "done": 0, "total": 0,
                         "vector_done": 0, "vector_total": 0, "vector_base": 0,
                         "chunks": 0, "error": "", "current_file": "",
                         "started_at": 0.0, "finished_at": 0.0, "cancelled": False,
                         "failed": 0,
                         # WELCHE Dateien gescheitert sind, mit Grund. Bis
                         # 2026-07-31 gab es nur die ZAHL, und die Oberflaeche
                         # verwies auf „siehe Journal" – wo die Namen gar nicht
                         # standen: der haeufigste Zweig protokolliert mit
                         # _log.info, und ohne Logging-Konfiguration verwirft
                         # Python alles unterhalb von WARNING. Der Hinweis zeigte
                         # damit auf Information, die nie geschrieben wurde.
                         "failed_list": []}
# Deckel fuer failed_list: bei einem unerreichbaren Netzlaufwerk koennen tausende
# Dateien scheitern. Die Liste geht in JEDE Fortschritts-Antwort (das Frontend
# pollt im Sekundentakt) und in last_index.json – unbegrenzt waere sie eine
# Bremse und ein Speicherleck zugleich.
MAX_FAILED_LIST = 50
_progress_lock = threading.Lock()
# Alle wieviel Dateien der Bulk-Reindex auf Platte sichert + Speicher ans OS
# zurueckgibt. Kompromiss: haeufiger = weniger Verlust bei Absturz, aber mehr
# I/O; 25 verliert im schlimmsten Fall 25 Dateien, die die Wiederaufnahme ohnehin
# nachholt.
CHECKPOINT_EVERY = 25
# Metadaten des gerade laufenden Laufs (fuer die Platten-Checkpoints).
_current_run: dict = {}
# Verhindert PARALLELE Reindex-Laeufe – sonst teilen sie sich _index_progress und
# die Zaehler ueberschreiben sich (z.B. vector_done=48 / vector_total=10 -> 480%).
_reindex_lock = threading.Lock()
# Kam waehrend eines laufenden Reindex eine weitere Anfrage, wird GENAU EINMAL
# nachgeholt (coalesced) – so gehen frisch hinzugefuegte Dateien nicht verloren.
_reindex_rerun = threading.Event()
# Abbruchwunsch des Benutzers. Wird nur zwischen zwei Dateien geprueft – eine
# laufende Einbettung wird nicht mitten drin abgeschossen.
_reindex_cancel = threading.Event()

def get_index_progress() -> dict:
    with _progress_lock:
        return dict(_index_progress)

def _set_progress(**kwargs):
    with _progress_lock:
        _index_progress.update(kwargs)


def _note_failed(path: str, grund: str) -> None:
    """Haelt eine gescheiterte Datei MIT GRUND fest – fuer die Oberflaeche.

    Der Grund ist so formuliert, dass er ohne Vorwissen handlungsleitend ist
    („zu gross", „kein Parser") – nicht als Ausnahmetext, mit dem ein Admin
    nichts anfangen kann.
    """
    try:
        rel = str(Path(path).relative_to(PROJECT_ROOT))
    except (ValueError, TypeError):
        rel = str(path)
    with _progress_lock:
        liste = _index_progress.setdefault("failed_list", [])
        if len(liste) < MAX_FAILED_LIST:
            liste.append({"file": rel, "reason": grund})

# ─── Vector Store (optional, Fallback auf TF-IDF) ────────────────
_vector_store = None
_vector_store_checked = False

# ─── Gecachte Stats (Format-Support ändert sich nie zur Laufzeit) ─
_stats_cache: dict | None = None
_stats_cache_lock = threading.Lock()


# Nach einem FEHLER (nicht: fehlender Abhaengigkeit) wird die Initialisierung
# spaeter erneut versucht. Frueher galt ein einmaliger Fehler bis zum Neustart –
# der Dienst lief dann dauerhaft ohne semantische Suche weiter, mit einer
# einzigen Warnzeile im Journal.
_VS_RETRY_SEC = 60.0
_vector_store_retry_after = 0.0


def _get_vector_store():
    """Gibt VectorStore-Singleton zurueck oder None wenn Dependencies fehlen."""
    global _vector_store, _vector_store_checked, _vector_store_retry_after
    if _vector_store is not None:
        return _vector_store
    if _vector_store_checked and time.time() < _vector_store_retry_after:
        return None
    _vector_store_checked = True
    try:
        from backend.tools.vector_store import VectorStore
        vs = VectorStore(PROJECT_ROOT / "data" / "vector_store")
        _vector_store = vs
        _log.info("VectorStore verfuegbar – semantische Suche aktiv")
        return vs
    except ImportError as e:
        # Fehlende Abhaengigkeit aendert sich zur Laufzeit nicht -> kein Neuversuch.
        _vector_store_retry_after = float("inf")
        _log.info(f"VectorStore nicht verfuegbar (faiss-cpu/sentence-transformers fehlt): {e}")
        return None
    except Exception as e:
        # Voruebergehend (Rechte, Platte voll, beschaedigter Index) -> spaeter erneut.
        _vector_store_retry_after = time.time() + _VS_RETRY_SEC
        _log.warning(f"VectorStore Initialisierung fehlgeschlagen (Neuversuch in "
                     f"{int(_VS_RETRY_SEC)}s): {e}")
        return None


def vector_store_status() -> dict:
    """Zustand der Vektorsuche fuer die Statistik-Anzeige."""
    # DREI Zustaende, nicht zwei: _vector_store entsteht erst beim ersten
    # Zugriff. "noch nicht initialisiert" darf NICHT als "nicht verfuegbar"
    # gemeldet werden – sonst zeigt die Oberflaeche nach jedem Neustart eine
    # Stoerung an, die es nicht gibt (gleiche Lazy-Falle wie beim Hauptagenten).
    if _vector_store is not None:
        state = "ok"
    elif not _vector_store_checked:
        state = "unbenutzt"           # noch nie gebraucht – keine Aussage moeglich
    elif _vector_store_retry_after == float("inf"):
        state = "nicht_installiert"   # faiss/sentence-transformers fehlen
    else:
        state = "fehler"              # Aufbau gescheitert, Neuversuch laeuft
    return {
        "state": state,
        "available": _vector_store is not None,
        "permanently_off": state == "nicht_installiert",
        "retry_in": max(0, int(_vector_store_retry_after - time.time()))
        if state == "fehler" else 0,
    }


def preload_embedding_model():
    """Lädt das Embedding-Modell im Hintergrund vor (vermeidet Kaltstart bei erster Suche)."""
    vs = _get_vector_store()
    if vs is None:
        print("[knowledge] Embedding-Preload übersprungen (kein VectorStore)", flush=True)
        return
    try:
        from backend.tools.vector_store import _get_embedding_model
        print("[knowledge] Lade Embedding-Modell vor...", flush=True)
        _get_embedding_model()
        print("[knowledge] Embedding-Modell vorgeladen ✓", flush=True)
    except Exception as e:
        print(f"[knowledge] Embedding-Modell Preload fehlgeschlagen: {e}", flush=True)


def _rebuild_vector_index(folders: list[Path], max_bytes: int, force: bool = False) -> bool:
    """Inkrementeller Vektor-Index Aufbau. Gibt True zurueck wenn Index Inhalt hat.

    force=False (Suchpfad): Kein Bulk-Aufbau bei leerem Index, max. INLINE_LIMIT
    Dateien, und AUSDRUECKLICH KEIN Entfernen verwaister Eintraege.
    force=True  (Neu-Indizieren): Alle Dateien, inkl. Aufraeumen.

    Warum das Aufraeumen nicht mehr im Suchpfad passiert: Es hielt "Datei nicht
    auffindbar" fuer "Datei geloescht". Antwortete ein Netzlaufwerk kurz nicht,
    ueberging ``_all_files`` den ganzen Ordner – und die naechstbeste Suche loeschte
    saemtliche Chunks dieses Shares aus dem Index. Aufraeumen ist eine
    Wartungsaufgabe und gehoert in den ausdruecklichen Neuaufbau.
    """
    vs = _get_vector_store()
    if vs is None:
        return False

    indexed = vs.get_indexed_files()

    if not force:
        # Leerer Index: kein Inline-Bulk-Indexing
        if not indexed:
            _log.debug("Vektor-Index leer – bitte Neu-Indizieren ausfuehren")
            return False

    # Nur erreichbare Ordner betrachten. Die Liste wird VOR dem Scan bestimmt,
    # damit Scan und Aufraeumen garantiert denselben Stand sehen.
    alive = [f for f in folders if _safe_exists(f)]
    if len(alive) != len(folders):
        missing = [str(f) for f in folders if f not in alive]
        _log.warning("Nicht erreichbare Wissensordner werden uebersprungen "
                     "(Index bleibt unangetastet): %s", ", ".join(missing))
    # Suchpfad: kurz gecachte Liste (der Ordner-Scan lief bisher bei JEDER Suche
    # ueber alle Wissensordner inkl. Netzlaufwerken). Der Neuaufbau nimmt den
    # echten Stand.
    files = _all_files(alive) if force else _all_files_cached(alive)
    current_paths = {str(f) for f in files}

    # Verwaiste Eintraege entfernen – nur beim ausdruecklichen Neuaufbau und nur
    # fuer Dateien aus ERREICHBAREN Ordnern. In EINEM Neuaufbau statt N.
    if force:
        alive_prefixes = tuple(str(r).rstrip(os.sep) + os.sep for r in alive)
        stale = [p for p in indexed
                 if p not in current_paths and p.startswith(alive_prefixes)]
        if stale:
            removed = vs.remove_files(stale)
            _log.info(f"{len(stale)} verwaiste Datei(en) aus dem Index entfernt "
                      f"({removed} Chunks)")

    # Neue/geaenderte Dateien ermitteln
    to_index = []
    for filepath in files:
        path_str = str(filepath)
        try:
            mtime = filepath.stat().st_mtime
        except Exception:
            continue
        if indexed.get(path_str) != mtime:
            to_index.append(filepath)

    if not force and len(to_index) > INLINE_LIMIT:
        _log.info(f"{len(to_index)} neue/geaenderte Dateien – nur {INLINE_LIMIT} inline, Rest via Neu-Indizieren")
        to_index = to_index[:INLINE_LIMIT]

    # Bereits indizierte Dateien zaehlen mit (der Voll-Reindex ueberspringt
    # unveraenderte Dateien) – so zeigt der Balken bei einer WIEDERAUFNAHME nach
    # Absturz den echten Gesamtstand, nicht nur die Rest-Dateien.
    already = max(0, len(indexed) - len([f for f in to_index if str(f) in indexed]))
    total = len(to_index)
    _set_progress(phase="Vektor", vector_done=0, vector_total=total, vector_base=already)

    changed = 0
    failed = 0
    cancelled = False
    for i, filepath in enumerate(to_index):
        # Abbruch nur ZWISCHEN zwei Dateien – die bereits geschriebenen Chunks
        # bleiben gueltig, der Index ist danach lediglich unvollstaendig.
        if _reindex_cancel.is_set():
            cancelled = True
            _log.info(f"Vektor-Index: Abbruch nach {i}/{total} Dateien")
            break
        path_str = str(filepath)
        _set_progress(vector_done=i + 1, phase=f"Vektor: {filepath.name[:40]}",
                      current_file=filepath.name)
        try:
            mtime = filepath.stat().st_mtime
            text = _extract_text(filepath, max_bytes)
            if text and text.strip():
                chunks = _chunk_text(text)
                # save=False: nicht bei jeder Datei den ganzen Index schreiben.
                vs.add_chunks(path_str, chunks, mtime, save=False)
                changed += 1
            elif text is None:
                # NICHT lesbar (zu gross, Parser fehlt, defekt) – der bisherige
                # Indexstand BLEIBT. Frueher wurde hier entfernt: eine wachsende
                # Datei verlor beim Ueberschreiten des Groessenlimits still ihre
                # Chunks, ohne dass irgendwo ein Fehler auftauchte.
                failed += 1
                grund = _unlesbar_grund(filepath, max_bytes)
                if path_str in indexed:
                    _note_failed(path_str, f"{grund} – bisheriger Indexstand bleibt erhalten")
                    _log.warning("Nicht lesbar (%s), bisheriger Indexstand bleibt: %s",
                                 grund, path_str)
                else:
                    _note_failed(path_str, grund)
                    _log.warning("Nicht lesbar (%s), wird nicht indiziert: %s",
                                 grund, path_str)
            else:
                # Lesbar, aber leer -> Eintrag ist gegenstandslos.
                vs.remove_file(path_str)
        except Exception as e:
            failed += 1
            _note_failed(path_str, f"{type(e).__name__}: {e}"[:200])
            _log.warning("Indizierung fehlgeschlagen fuer %s: %s", path_str, e)

        # Laufende Chunk-Zahl mitfuehren, damit ALLE offenen Clients dieselbe
        # Live-Zahl sehen (sonst zeigt ein Browser, der die Kachel vor dem Lauf
        # geladen hat, dauerhaft den alten Stand – z.B. 16453 statt 9715).
        try:
            _set_progress(chunks=vs.chunk_count())
        except Exception:
            pass

        # Checkpoint: alle CHECKPOINT_EVERY Dateien auf Platte sichern, Speicher
        # ans OS zurueckgeben (verhindert Heap-Wachstum → OOM) und den
        # Fortschritt persistent festhalten (welche Datei, wie weit) – so ist
        # nach einem Absturz sichtbar, wo es endete, und die Wiederaufnahme
        # setzt genau dort fort.
        if (i + 1) % CHECKPOINT_EVERY == 0:
            try:
                vs.save()
                from backend.tools.vector_store import release_memory_to_os
                release_memory_to_os()
            except Exception as e:
                _log.warning(f"Checkpoint fehlgeschlagen: {e}")
            _write_run_checkpoint(done=i + 1, total=total,
                                  current_file=filepath.name,
                                  indexed_files=already + changed)

    # Rest sichern (der letzte, unvollstaendige Checkpoint-Block).
    try:
        vs.save()
    except Exception as e:
        _log.warning(f"Abschluss-Speichern fehlgeschlagen: {e}")

    if not cancelled:
        _set_progress(vector_done=total, vector_total=total)
    # Fehlerzahl mitfuehren: ein Lauf, der die Haelfte der Dateien verschluckt,
    # sah bisher wie ein Erfolg aus (die Ausnahmen wurden still geschluckt).
    _set_progress(failed=failed)
    if changed or failed:
        _log.info(f"Vektor-Index aktualisiert: {changed} Datei(en), {failed} fehlgeschlagen")
    return vs.chunk_count() > 0


# Gelernte Konversationen (learned/conv_*.md) tragen die urspruengliche
# Benutzerfrage als Ueberschrift. Dadurch sind sie fuer genau diese Frage der
# perfekte semantische Treffer – unabhaengig davon, ob ihr Inhalt zur Frage
# passt – und verdraengen die Primaerdokumentation vom ersten Platz. Das ist
# eine selbstverstaerkende Schleife: eine falsche Antwort wird gelernt und beim
# naechsten Mal bevorzugt wieder ausgeliefert. Deshalb im Ranking abwerten.
LEARNED_PENALTY = 0.6


def _is_learned_note(path_str: str) -> bool:
    p = path_str.replace("\\", "/")
    return "/knowledge/learned/" in p or "/knowledge/pending/" in p


def _learned_weight(file_path: str) -> float:
    """Herkunfts-Gewicht fuer das Ranking (1.0 = normal)."""
    return LEARNED_PENALTY if _is_learned_note(file_path) else 1.0


class _TrefferListe(list):
    """Trefferliste mit einem Zusatzmerkmal am Ergebnis selbst.

    Warum kein Modul-Merker: Suchen laufen ueber ``asyncio.to_thread``, mehrere
    Benutzer koennen gleichzeitig suchen. Ein Modul-Dict wuerde die Kennzeichnung
    des einen Laufs an die Antwort des anderen haengen – dieselbe Klasse Fehler
    wie beim Actor-Kontext (siehe CLAUDE.md). Ein ContextVar hilft hier NICHT:
    ``to_thread`` kopiert den Kontext IN den Thread, Aenderungen darin kommen
    nicht zurueck. Also haengt das Merkmal am zurueckgegebenen Objekt.

    ``kein_anker``: kein einziges Wort der Anfrage kommt im Bestand vor.
    """
    kein_anker: bool = False


def _vector_search(query: str, max_results: int,
                   allow_paths: set | None = None) -> list[tuple[float, str, str]] | None:
    """Hybride Suche (semantisch + BM25) via VectorStore.

    Gibt None zurueck wenn kein VectorStore verfuegbar ist.

    ``allow_paths`` schraenkt auf die Wissensgruppen des Fragenden ein und wird
    IN die Suche gereicht – nicht nachtraeglich gefiltert (Begruendung siehe
    ``VectorStore.search_hybrid``).

    Die Abwertung gelernter Notizen wird als ``weight_fn`` IN die Suche gereicht,
    damit sie VOR dem relativen Cut greift. Vorher lief sie hier nachtraeglich:
    Stand eine gelernte Notiz auf Platz 1, wurde der Cut an ihrem unabgewerteten
    Score gemessen – Primaerdokumente, die nach der Abwertung vorn gelegen
    haetten, waren da schon verworfen.
    """
    vs = _get_vector_store()
    if vs is None:
        return None
    # search_hybrid_ex statt search_hybrid + has_lexical_anchor: letzteres hat
    # denselben BM25-Durchlauf ein ZWEITES Mal gerechnet, obwohl die Hybridsuche
    # das Ergebnis Millisekunden vorher schon hatte und wegwarf.
    results, ohne_anker = vs.search_hybrid_ex(query, max_results,
                                              weight_fn=_learned_weight,
                                              allow_paths=allow_paths)
    if not results:
        return None
    kein_anker = ohne_anker is True

    converted = _TrefferListe()
    for score, file_path, chunk in results:
        try:
            rel = str(Path(file_path).relative_to(PROJECT_ROOT))
        except ValueError:
            rel = file_path
        converted.append((score, rel, chunk))
    converted.kein_anker = kein_anker
    return converted


# ─── Aehnlichkeitspruefung fuer Entwuerfe (Dubletten/Widersprueche) ──────────
# Ab welcher Cosine-Aehnlichkeit ein vorhandener Chunk dem Menschen gezeigt wird.
# Nicht geraten: am Echt-Index (9207 Chunks) liegen echte Fachfragen gegen ihren
# passenden Chunk bei 0.84-0.90, thematisch fremde Texte bei 0.79-0.83. 0.86
# liegt im oberen Bereich – die Liste soll kurz und einschlaegig sein, nicht
# vollstaendig. Wer alles sehen will, sucht in /wissen.
AEHNLICH_AB = 0.86
# Je Abfrage (Titel, Zusammenfassung, Fakten, Fragen) hoechstens so viele
# Nachbarn; insgesamt gedeckelt, damit die Pruefansicht nicht zur Trefferliste wird.
AEHNLICH_JE_ABFRAGE = 3
AEHNLICH_GESAMT = 8
# Relativer Schnitt: nur Treffer, die nah am besten liegen. Siehe Begruendung
# in find_similar_existing() – der Abstand zum Spitzenwert trennt, nicht der
# Absolutwert.
AEHNLICH_REL = 0.97


def find_similar_existing(doc: dict) -> dict:
    """Sucht zu einem Entwurf bereits vorhandene, sehr aehnliche Chunks.

    Abgefragt wird der INHALT, nicht nur der Titel: eine Dublette traegt selten
    dieselbe Ueberschrift. Genommen werden Titel, Zusammenfassung, die ersten
    Fakten und die ersten Fragen – jede fuer sich, weil ein zusammengesetzter
    Text den Vektor verwaessert und dann zu allem maessig passt.

    Gelernte Notizen werden AUSGESCHLOSSEN: sie tragen die Benutzerfrage als
    Ueberschrift und waeren fuer eine Frage-Antwort-Abfrage immer der Top-Treffer,
    unabhaengig vom Inhalt (dieselbe Selbstverstaerkung, deretwegen sie in der
    normalen Suche mit LEARNED_PENALTY abgewertet werden).

    Rueckgabe: ``{"items": [{file, score, text, matched}], "checked": n}``
    """
    vs = _get_vector_store()
    if vs is None:
        raise RuntimeError("Vektor-Index nicht verfuegbar")

    abfragen: list[tuple[str, str]] = []
    titel = (doc.get("title") or "").strip()
    if titel:
        abfragen.append(("Titel", titel))
    zus = (doc.get("summary") or "").strip()
    if zus:
        abfragen.append(("Zusammenfassung", zus[:400]))
    for f in (doc.get("facts") or [])[:3]:
        f = str(f).strip()
        if f:
            abfragen.append(("Fakt", f[:300]))
    for qa in (doc.get("qa_pairs") or [])[:3]:
        frage = str((qa or {}).get("question") or "").strip()
        antwort = str((qa or {}).get("answer") or "").strip()
        if frage:
            # Frage UND Antwort: die Frage allein findet oft nur gleichartige
            # Fragen, die Antwort traegt die eigentliche Aussage.
            abfragen.append(("Frage", (frage + " " + antwort)[:300]))

    gesehen: set[str] = set()
    items: list[dict] = []
    for wo, text in abfragen:
        try:
            treffer = vs.search(text, AEHNLICH_JE_ABFRAGE)
        except Exception:
            continue
        for score, pfad, chunk in treffer:
            if float(score) < AEHNLICH_AB:
                continue
            if _is_learned_note(pfad):
                continue
            schluessel = f"{pfad}|{chunk[:80]}"
            if schluessel in gesehen:
                continue
            gesehen.add(schluessel)
            try:
                rel = str(Path(pfad).relative_to(PROJECT_ROOT))
            except (ValueError, TypeError):
                rel = str(pfad)
            items.append({"file": rel, "score": round(float(score), 3),
                          "text": chunk.strip()[:600], "matched": wo})
    items.sort(key=lambda x: x["score"], reverse=True)

    # RELATIVER SCHNITT – ohne ihn ist die Liste ueberwiegend Beifang.
    # e5 komprimiert Cosine auf ~0.83-0.93; oberhalb einer absoluten Schwelle
    # landen deshalb auch voellig unbeteiligte Chunks. Gemessen am Echt-Index:
    # ein woertlich uebernommener Chunk fand sich mit 0.931 – und daneben drei
    # thematisch fremde mit 0.868-0.871. Der Abstand zum Spitzenwert ist das
    # aussagekraeftige Signal, nicht der Absolutwert.
    #
    # Der Grund, warum das wichtig ist, ist derselbe wie bei den
    # "fehlgeschlagenen Dateien": eine Warnung, die zu drei Vierteln aus
    # Rauschen besteht, wird nach dem zweiten Mal weggeklickt.
    if items:
        boden = items[0]["score"] * AEHNLICH_REL
        items = [x for x in items if x["score"] >= boden]
    return {"items": items[:AEHNLICH_GESAMT], "checked": len(abfragen)}


def _get_skill_config() -> dict:
    try:
        return config.get_skill_states().get("knowledge", {}).get("config", {})
    except Exception:
        return {}


def _get_folders() -> list[Path]:
    cfg = _get_skill_config()
    folders_str = cfg.get("folders", DEFAULT_FOLDER)
    paths = []
    for f in folders_str.split(","):
        f = f.strip()
        if not f:
            continue
        p = Path(f)
        if not p.is_absolute():
            p = PROJECT_ROOT / p
        paths.append(p)
    return paths or [PROJECT_ROOT / DEFAULT_FOLDER]


# ─── Nicht-blockierende Verfügbarkeitspruefung fuer (Netz-)Ordner ─────────────
# Ein totes CIFS/NFS-Mount laesst exists()/os.walk() bis zum Kernel-Timeout
# blockieren ("Lädt…" haengt ewig). Wir pruefen exists() daher in einem
# Daemon-Thread mit kurzem Timeout und cachen ein negatives Ergebnis kurz.
_avail_down_until: dict[str, float] = {}
_AVAIL_DOWN_TTL = 30.0   # Sekunden, wie lange ein totes Mount als "weg" gilt


def _safe_exists(path, timeout: float = 2.0) -> bool:
    """exists()-Check, der bei toten Netzlaufwerken NICHT blockiert.

    Laeuft in einem Daemon-Thread; Timeout oder OSError => False. Ein als
    "blockierend/tot" erkanntes Verzeichnis wird kurz gecacht, damit nicht
    jeder Aufruf (z.B. Stats-Polling) erneut ins Timeout laeuft."""
    key = str(path)
    now = time.time()
    until = _avail_down_until.get(key)
    if until and now < until:
        return False

    result = {"ok": False}

    def _check():
        try:
            result["ok"] = os.path.exists(key)
        except OSError:
            result["ok"] = False

    th = threading.Thread(target=_check, daemon=True)
    th.start()
    th.join(timeout)
    if th.is_alive():
        # Haengt am toten Mount -> kurz als "weg" merken und den Thread
        # (Daemon) sich selbst beenden lassen, sobald der Kernel zurueckkehrt.
        _avail_down_until[key] = now + _AVAIL_DOWN_TTL
        _log.warning("Ordner reagiert nicht (Netzlaufwerk tot?), wird übersprungen: %s", key)
        return False
    _avail_down_until.pop(key, None)
    return bool(result["ok"])


def _bounded_call(fn, timeout: float, default):
    """Führt ``fn`` in einem Daemon-Thread aus und bricht nach ``timeout`` ab.

    Gibt bei Timeout ``default`` zurück (der haengende Thread laeuft als Daemon
    im Hintergrund aus). Schützt Hot-Paths (z.B. Stats) vor toten Netzlaufwerken."""
    box = {"val": default}

    def _run():
        try:
            box["val"] = fn()
        except Exception:
            box["val"] = default

    th = threading.Thread(target=_run, daemon=True)
    th.start()
    th.join(timeout)
    return box["val"]


def _get_max_bytes() -> int:
    try:
        mb = float(_get_skill_config().get("max_file_size_mb", DEFAULT_MAX_SIZE_MB))
    except Exception:
        mb = DEFAULT_MAX_SIZE_MB
    return int(mb * 1024 * 1024)


def _transcribe_media(filepath: Path) -> str | None:
    """Transkribiert Audio/Video via ffmpeg + faster-whisper."""
    if not shutil.which("ffmpeg"):
        _log.warning("ffmpeg nicht gefunden – Video/Audio-Support deaktiviert")
        return None

    try:
        from faster_whisper import WhisperModel
    except ImportError:
        _log.warning("faster-whisper nicht installiert – Video/Audio-Support deaktiviert")
        return None

    tmpdir = None
    try:
        # Audio aus Video/Audio extrahieren → WAV (16kHz mono, optimal für Whisper)
        tmpdir = tempfile.mkdtemp(prefix="jarvis_kb_")
        wav_path = os.path.join(tmpdir, "audio.wav")

        cmd = [
            "ffmpeg", "-i", str(filepath),
            "-vn",                    # kein Video
            "-acodec", "pcm_s16le",   # PCM 16-bit
            "-ar", "16000",           # 16kHz
            "-ac", "1",               # Mono
            "-y",                     # Überschreiben
            wav_path
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=300)
        if result.returncode != 0 or not os.path.exists(wav_path):
            _log.error(f"ffmpeg fehlgeschlagen für {filepath}: {result.stderr[:200]}")
            return None

        # Whisper-Modell laden (eigene Instanz, nicht die aus main.py)
        # Nutze "small" als Default – guter Kompromiss aus Geschwindigkeit und Qualität
        cfg = _get_skill_config()
        model_name = cfg.get("whisper_model", "small")
        model = WhisperModel(model_name, device="cpu", compute_type="int8")

        segments, info = model.transcribe(wav_path, language="de")
        text = " ".join([seg.text for seg in segments]).strip()

        if text:
            # Dateiname + erkannte Sprache als Kontext
            header = f"[Transkription: {filepath.name} | Sprache: {info.language}]"
            _log.info(f"Transkription OK für {filepath.name}: {len(text)} Zeichen")
            return f"{header}\n{text}"

        _log.warning(f"Keine Sprache erkannt in {filepath.name}")
        return None

    except subprocess.TimeoutExpired:
        _log.error(f"ffmpeg Timeout für {filepath}")
        return None
    except Exception as e:
        _log.error(f"Transkription fehlgeschlagen für {filepath}: {e}")
        return None
    finally:
        if tmpdir and os.path.exists(tmpdir):
            shutil.rmtree(tmpdir, ignore_errors=True)


def _ocr_image(filepath: Path) -> str | None:
    """OCR auf einem Bild via Tesseract (Deutsch+Englisch). Gibt erkannten Text zurueck.

    Lokal, kein LLM. Voraussetzung: System-Paket 'tesseract-ocr' (+ Sprachpakete)
    und Python-Pakete 'pytesseract' + 'Pillow'. Fehlt etwas, wird None zurueckgegeben
    (das LLM kann dann ggf. noch das Bild selbst auswerten – siehe extract_from_file).
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        _log.warning("pytesseract/Pillow nicht installiert – Bild-OCR deaktiviert")
        return None
    try:
        # Sprachen auf verfuegbare beschraenken (deu/eng), sonst Tesseract-Default
        lang = None
        try:
            avail = set(pytesseract.get_languages(config=""))
            sel = [l for l in ("deu", "eng") if l in avail]
            lang = "+".join(sel) if sel else None
        except Exception:
            lang = "deu+eng"
        with Image.open(str(filepath)) as img:
            img.load()
            text = pytesseract.image_to_string(img, lang=lang) if lang else pytesseract.image_to_string(img)
        text = (text or "").strip()
        return text or None
    except Exception as e:
        _log.warning(f"Bild-OCR fehlgeschlagen ({filepath.name}): {e}")
        return None


def _ocr_sprachen() -> str | None:
    """Verfuegbare Tesseract-Sprachen auf deu/eng eingegrenzt."""
    try:
        import pytesseract
        avail = set(pytesseract.get_languages(config=""))
        return "+".join([l for l in ("deu", "eng") if l in avail]) or None
    except Exception:
        return "deu+eng"


def _ocr_pdf_seiten(pdf_bytes: bytes, erste: int = 1, letzte: int = 20) -> dict[int, str]:
    """OCR einzelner Seiten – Rueckgabe {Seitennummer (1-basiert): Text}.

    SEITENWEISE und nicht als ein Fliesstext, weil der Aufrufer die Fassungen
    Seite fuer Seite gegeneinander stellen muss: bei einem PDF mit 54 Seiten und
    einem OCR-Deckel von 30 duerfen die restlichen 24 Seiten nicht verloren
    gehen. Ein zusammengefuegter String liesse sich dafuer nicht mehr zerlegen.
    """
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError:
        _log.warning("pdf2image/pytesseract fehlt – PDF-OCR deaktiviert")
        return {}
    if letzte < erste:
        return {}
    try:
        images = convert_from_bytes(pdf_bytes, dpi=200, first_page=erste, last_page=letzte)
    except Exception as e:
        _log.warning("PDF->Bild fehlgeschlagen: %s", e)
        return {}
    lang = _ocr_sprachen()
    aus: dict[int, str] = {}
    for versatz, img in enumerate(images):
        try:
            t = pytesseract.image_to_string(img, lang=lang) if lang else pytesseract.image_to_string(img)
            t = (t or "").strip()
            if t:
                aus[erste + versatz] = t
        except Exception:
            continue
    return aus


def _ocr_pdf_bytes(pdf_bytes: bytes, max_pages: int = 20) -> str:
    """OCR-Fallback fuer gescannte/bildbasierte PDFs (ohne Text-Layer).

    Rendert die Seiten via pdf2image/poppler zu Bildern und liest sie per
    Tesseract (deu+eng). Gibt erkannten Text zurueck oder '' wenn nicht moeglich.
    """
    seiten = _ocr_pdf_seiten(pdf_bytes, 1, max_pages)
    return "\n\n".join(f"[Seite {nr} (OCR)]\n{t}" for nr, t in sorted(seiten.items()))


# ===========================================================================
# Textqualitaet: erkennt eine BESCHAEDIGTE Textebene und laesst OCR entscheiden
# ===========================================================================
# ANLASS (Vorfall 2026-08-12 auf ECHT): ein 54-seitiges Anschreiben-PDF lieferte
# ueber pdfplumber 80.586 Zeichen - die alte Schwelle "weniger als 80 Zeichen ->
# OCR" greift also NIE. Der Text war aber teils unbrauchbar, weil die
# Zeichentabelle (cmap) der eingebetteten Schriften beschaedigt ist:
#   "Datum: OL.O7.2026" statt "01.07.2026"   (Buchstabe O statt Ziffer 0)
#   "Lauerstr.'14"      statt "Lauerstr. 14"
#   "ftir"              statt "fuer"
#   "ngirrrsf#s$"       als Logo-Zeile
# Das Modell hat daraufhin 17 Extraktionsskripte gebaut und die Adressen
# trotzdem nicht saubergekriegt.
#
# WAS AN DIESEM PROBLEM UEBERRASCHT (auf ECHT gemessen, nicht geschaetzt):
# die naheliegenden Kennzahlen sehen den Schaden NICHT. Stoppwortanteil,
# Vokalanteil und Sonderzeichenquote sind bei der kaputten Fassung so gut wie
# bei der OCR-Fassung - teils sogar besser (Stoppwortanteil 12,9 % gegen 11,4 %).
# Der Schaden besteht aus ZEICHEN-SUBSTITUTIONEN, und die erzeugen weiterhin
# aussprechbare Woerter. Wer hier eine Wortlisten- oder Entropiepruefung baut,
# misst am Problem vorbei.
#
# DESHALB ZWEI STUFEN, und die zweite ist die eigentliche Entscheidung:
#   1. Vorfilter (Millisekunden, reine Regex) - er stellt nur einen VERDACHT
#      fest und darf grosszuegig sein, weil ein Fehlalarm nur die Stichprobe
#      kostet.
#   2. OCR-STICHPROBE auf zwei Seiten + Vergleich beider Fassungen. Erst wenn
#      OCR messbar mehr liefert, wird das ganze Dokument neu gelesen.
#
# WARUM NICHT DER VORFILTER ALLEIN: er ist nachweislich fragil. An 753 echten
# Fachdokumenten gemessen schlug eine erste, breitere Fassung bei ICD-10-Codes
# ("O61.0"), PPR-Pflegekategorien ("A4S1") und GUIDs ("43B3B851") an - alles
# voellig korrekter Text, den OCR zeichengleich liefert. In einer Klinik-
# umgebung waere das ein Dauerfehlalarm. Die Stichprobe faengt genau das ab:
# sie stellt fest, dass OCR nichts gewinnt, und der Textlayer bleibt stehen.
#
# MESSWERTE (ECHT, 753 Fachdokumente + das gemeldete PDF):
#   Vorfilter:   Verdacht bei 29 von 753 (3,9 %)
#   Stichprobe:  gemeldetes PDF  Strukturtreffer 21 -> 37, Wortquote 58,6 -> 61,4  => OCR
#                gesunde Faelle  Strukturtreffer gleich, Wortquote 76,5 -> 59,2    => Textebene
#   Kosten:      Stichprobe 3,8-7,0 s, volles OCR rund 1,9 s je Seite

def _qs_zahl(name: str, vorgabe: int | float, klein, gross):
    """Stellschraube aus der Umgebung, auf einen sinnvollen Bereich begrenzt."""
    try:
        w = type(vorgabe)(os.environ.get(name, vorgabe))
    except (TypeError, ValueError):
        return vorgabe
    return max(klein, min(gross, w))


# 0 schaltet die Pruefung ganz ab (dann bleibt es beim alten Verhalten:
# OCR nur, wenn praktisch kein Text da ist).
_PDF_QS_AKTIV = os.environ.get("JARVIS_PDF_QS", "1").strip().lower() not in ("0", "aus", "off", "false")
# Unterhalb dieser Laenge ist jede Quote Zufall.
_PDF_QS_MIN_ZEICHEN = 200
# Seiten fuer die Stichprobe. Zwei reichten im Feldtest zur klaren Trennung;
# jede weitere kostet rund zwei Sekunden.
_PDF_QS_PROBE_SEITEN = _qs_zahl("JARVIS_PDF_QS_PROBE", 2, 1, 10)
# Obergrenze fuer das anschliessende volle OCR. Rund 1,9 s je Seite - der Wert
# ist eine Zeitentscheidung, keine technische Grenze. Was darueber liegt, bleibt
# bei der Textebene und wird im Ergebnis AUSGEWIESEN (kein stiller Verlust).
_PDF_QS_MAX_SEITEN = _qs_zahl("JARVIS_PDF_QS_MAX_SEITEN", 30, 1, 400)
# Wieviele Strukturtreffer muss OCR mindestens zusaetzlich finden.
_PDF_QS_STRUKTUR_PLUS = _qs_zahl("JARVIS_PDF_QS_STRUKTUR_PLUS", 2, 1, 100)
# Um wieviele Prozentpunkte muss die Wortquote steigen.
_PDF_QS_WORT_PLUS = _qs_zahl("JARVIS_PDF_QS_WORT_PLUS", 4.0, 0.5, 50.0)

# Verstuemmeltes Datum: mindestens ein O/I/l in einem sonst gueltigen Datum.
_TQ_DATUM = re.compile(
    r"\b(?=[^\s]*[OIl])[0-9OIl]{1,2}[.,][0-9OIl]{1,2}[.,](?:19|20)[0-9OIl]{2}\b")
# Buchstabe ZWISCHEN zwei Ziffern. Die Einklammerung ist entscheidend: ein
# fuehrendes O vor Ziffern waere ein ICD-10-Code (O61.0), kein Schaden.
# S und B sind bewusst NICHT dabei - "A4S1" und "43B3B851" sind echte Kennungen.
_TQ_INNEN = re.compile(r"\d[OIl]\d")
# Apostroph zwischen Wort/Punkt und Ziffer: "Lauerstr.'14"
_TQ_KLEB = re.compile(r"[A-Za-zÄÖÜäöüß.]['`´]\d")
_TQ_VOKAL = re.compile(r"[aeiouäöüyAEIOUÄÖÜY]")

# Strukturdaten - der Massstab im Vergleich. Bewusst genau die Angaben, wegen
# derer man ein solches Dokument ueberhaupt auswertet.
_TQ_MAIL = re.compile(r"\b[\w.+-]+@[\w-]+\.[\w.]{2,}\b")
_TQ_TEL = re.compile(r"(?:\+\d{2}[\s./-]?|\b0)\d{2,5}[\s./-]?\d{3,}")
_TQ_PLZ = re.compile(r"\b\d{5}\b")
_TQ_DAT_OK = re.compile(r"\b\d{1,2}\.\d{1,2}\.(?:19|20)\d{2}\b")
_TQ_IBAN = re.compile(r"\b[A-Z]{2}\d{2}[A-Z0-9]{10,28}\b")
_TQ_WORT = re.compile(r"[A-Za-zÄÖÜäöüß]{4,}")

# Wortlisten des Systems. OPTIONAL - fehlt sie, entscheidet allein die Zahl der
# Strukturtreffer (dann strenger, siehe _ocr_gewinnt).
_WORTLISTEN = ("/usr/share/dict/ngerman", "/usr/share/hunspell/de_DE.dic",
               "/usr/share/dict/german", "/usr/share/myspell/de_DE.dic")
_wortliste_cache: set[str] | None = None


def _wortliste() -> set[str]:
    """Deutsche Wortliste, einmal geladen. Leer, wenn keine vorhanden ist."""
    global _wortliste_cache
    if _wortliste_cache is not None:
        return _wortliste_cache
    _wortliste_cache = set()
    for p in _WORTLISTEN:
        f = Path(p)
        if not f.exists():
            continue
        try:
            woerter = set()
            for z in f.read_text(encoding="utf-8", errors="ignore").split("\n"):
                # hunspell haengt Flags mit "/" an: "Haus/N"
                w = z.split("/")[0].strip().lower()
                if len(w) >= 4:
                    woerter.add(w)
            if len(woerter) > 5000:
                _wortliste_cache = woerter
                _log.info("Wortliste geladen: %s (%d Eintraege)", p, len(woerter))
                break
        except Exception as e:
            _log.debug("Wortliste %s nicht lesbar: %s", p, e)
    if not _wortliste_cache:
        _log.info("Keine Wortliste gefunden – PDF-Qualitaetspruefung nutzt nur Strukturdaten")
    return _wortliste_cache


def pdf_text_verdacht(text: str) -> dict:
    """Stufe 1: Sieht die Textebene beschaedigt aus? Rueckgabe = Verdachtsgruende.

    Leeres dict heisst "kein Verdacht". Die Funktion entscheidet NICHTS - sie
    sagt nur, ob sich die Stichprobe lohnt.
    """
    if not text or len(text) < _PDF_QS_MIN_ZEICHEN:
        return {}
    kilo = len(text) / 1000
    gruende: dict[str, float] = {}

    n = len(_TQ_DATUM.findall(text))
    if n:
        gruende["datum_verstuemmelt"] = n

    n = len(_TQ_INNEN.findall(text))
    if n / kilo >= 0.15:
        gruende["ziffer_buchstabe_je_1k"] = round(n / kilo, 3)

    n = len(_TQ_KLEB.findall(text))
    if n / kilo >= 0.10:
        gruende["apostroph_vor_zahl_je_1k"] = round(n / kilo, 3)

    # Zeilen ohne jeden Vokal, aber mit Buchstaben UND Sonderzeichen: das ist
    # das Bild einer Logo- oder Symbolschrift ohne brauchbare Zeichentabelle.
    muell = 0
    for z in text.split("\n"):
        z = z.strip()
        if not 6 <= len(z) <= 40 or _TQ_VOKAL.search(z):
            continue
        if sum(c.isalpha() for c in z) >= 4 and any(
                not c.isalnum() and not c.isspace() for c in z):
            muell += 1
    if muell:
        gruende["muellzeilen"] = muell
    return gruende


def text_guete(text: str) -> dict:
    """Wie brauchbar ist dieser Text? Strukturtreffer + Wortquote in Prozent.

    Absichtlich KEINE Bewertung der Schadensmuster aus ``pdf_text_verdacht``:
    mit demselben Massstab zu messen, der den Verdacht ausgeloest hat, waere ein
    Zirkelschluss. Gezaehlt wird, was am Ende gebraucht wird.
    """
    if not text or not text.strip():
        return {"struktur": 0, "wortquote": 0.0}
    struktur = (len(_TQ_MAIL.findall(text)) + len(_TQ_TEL.findall(text))
                + len(_TQ_PLZ.findall(text)) + len(_TQ_DAT_OK.findall(text))
                + len(_TQ_IBAN.findall(text)))
    liste = _wortliste()
    woerter = [w.lower() for w in _TQ_WORT.findall(text)]
    quote = (100 * sum(1 for w in woerter if w in liste) / len(woerter)
             if woerter and liste else 0.0)
    return {"struktur": struktur, "wortquote": round(quote, 1)}


def _ocr_gewinnt(alt: dict, neu: dict, hat_wortliste: bool) -> bool:
    """Ist die OCR-Fassung messbar besser als die Textebene?

    Fail-closed: im Zweifel bleibt die Textebene. Ein unnoetiger OCR-Lauf kostet
    Minuten, und eine schlechtere Fassung zu uebernehmen waere Verschlimmbesserung
    - auf gesunden Dokumenten faellt die Wortquote durch OCR nachweislich
    (76,5 % -> 59,2 % im Feldtest).
    """
    mehr_struktur = (neu["struktur"] >= alt["struktur"] + _PDF_QS_STRUKTUR_PLUS
                     and neu["struktur"] > alt["struktur"] * 1.25)
    if not hat_wortliste:
        # Ohne Wortliste ist die Zahl der Strukturtreffer der einzige Massstab.
        return mehr_struktur
    bessere_woerter = neu["wortquote"] > alt["wortquote"] + _PDF_QS_WORT_PLUS
    # Eine deutlich SCHLECHTERE Wortquote widerlegt auch mehr Strukturtreffer:
    # verrauschtes OCR erzeugt Zahlenfolgen, die wie Telefonnummern aussehen.
    if neu["wortquote"] + _PDF_QS_WORT_PLUS < alt["wortquote"]:
        return False
    return mehr_struktur or bessere_woerter


def pdf_qualitaet_sichern(pdf_bytes: bytes, seiten: list[str]) -> tuple[list[str], dict]:
    """Prueft die Textebene und ersetzt sie SEITENWEISE durch OCR, wo es besser ist.

    ``seiten`` ist der Text je Seite aus pdfplumber (Index 0 = Seite 1).
    Rueckgabe: (Seiten nach der Pruefung, Bericht).

    Der Bericht ist kein Beiwerk - er geht in den Hinweis an das Modell und ins
    Journal. Ohne ihn waere von aussen nicht erkennbar, ob eine Antwort auf der
    Textebene oder auf OCR beruht, und der naechste Fehlerbericht begaenne wieder
    bei null.
    """
    beginn = time.time()
    bericht = {"geprueft": False, "verdacht": {}, "ocr": False, "seiten_ocr": 0,
               "seiten_gesamt": len(seiten), "grund": ""}
    if not _PDF_QS_AKTIV or not seiten:
        bericht["grund"] = "abgeschaltet" if not _PDF_QS_AKTIV else "keine Seiten"
        return seiten, bericht

    voll = "\n\n".join(s for s in seiten if s)
    bericht["geprueft"] = True
    verdacht = pdf_text_verdacht(voll)
    bericht["verdacht"] = verdacht
    if not verdacht:
        bericht["grund"] = "Textebene unauffaellig"
        return seiten, bericht

    # --- Stufe 2: Stichprobe auf den textreichsten Seiten --------------------
    # Die textreichsten, weil eine fast leere Seite in beiden Fassungen gleich
    # aussieht und nichts entscheidet.
    kandidaten = sorted(range(len(seiten)), key=lambda i: -len(seiten[i]))
    probe = sorted(i for i in kandidaten[:_PDF_QS_PROBE_SEITEN] if seiten[i].strip())
    if not probe:
        bericht["grund"] = "keine Seite mit Text fuer die Stichprobe"
        return seiten, bericht

    alt = {"struktur": 0, "wortquote": 0.0}
    neu = {"struktur": 0, "wortquote": 0.0}
    getroffen = 0
    for i in probe:
        ocr = _ocr_pdf_seiten(pdf_bytes, i + 1, i + 1).get(i + 1, "")
        if not ocr:
            continue
        getroffen += 1
        a, b = text_guete(seiten[i]), text_guete(ocr)
        alt["struktur"] += a["struktur"]; alt["wortquote"] += a["wortquote"]
        neu["struktur"] += b["struktur"]; neu["wortquote"] += b["wortquote"]
    if not getroffen:
        # Kein OCR moeglich (Pakete fehlen, poppler kaputt) - das ist KEIN
        # Qualitaetsurteil, also bleibt alles wie es ist.
        bericht["grund"] = "OCR nicht verfuegbar"
        return seiten, bericht
    alt["wortquote"] = round(alt["wortquote"] / getroffen, 1)
    neu["wortquote"] = round(neu["wortquote"] / getroffen, 1)
    bericht["probe"] = {"seiten": [i + 1 for i in probe], "textebene": alt, "ocr": neu}

    if not _ocr_gewinnt(alt, neu, bool(_wortliste())):
        bericht["grund"] = ("Stichprobe: OCR bringt nichts "
                            f"(Struktur {alt['struktur']}->{neu['struktur']}, "
                            f"Wortquote {alt['wortquote']}->{neu['wortquote']})")
        return seiten, bericht

    # --- Stufe 3: volles OCR, seitenweise gemischt --------------------------
    grenze = min(len(seiten), _PDF_QS_MAX_SEITEN)
    ocr_seiten = _ocr_pdf_seiten(pdf_bytes, 1, grenze)
    ergebnis = list(seiten)
    ersetzt = 0
    for nr, txt in ocr_seiten.items():
        i = nr - 1
        if not (0 <= i < len(ergebnis)):
            continue
        # Auch hier je Seite pruefen: in einem Dokument koennen sich saubere und
        # beschaedigte Schriften mischen, und eine Seite zu verschlechtern waere
        # genau der Fehler, den diese Funktion verhindern soll.
        if _ocr_gewinnt(text_guete(ergebnis[i]), text_guete(txt), bool(_wortliste())) \
                or len(ergebnis[i].strip()) < 40:
            ergebnis[i] = txt
            ersetzt += 1
    bericht["ocr"] = ersetzt > 0
    bericht["seiten_ocr"] = ersetzt
    bericht["dauer_s"] = round(time.time() - beginn, 1)
    bericht["grund"] = (f"OCR angewandt auf {ersetzt} von {len(seiten)} Seiten"
                        if ersetzt else "OCR brachte seitenweise doch keinen Gewinn")
    if len(seiten) > grenze:
        bericht["nicht_geprueft_ab"] = grenze + 1
    # Die Dauer gehoert ins Journal: rund zwei Sekunden je Seite sind der Preis
    # dieser Funktion, und wer eine langsame Antwort untersucht, soll sie hier
    # finden statt sie zu suchen.
    _log.info("PDF-Qualitaet: %s in %.1fs | Verdacht %s",
              bericht["grund"], bericht["dauer_s"], verdacht)
    return ergebnis, bericht


def qualitaets_hinweis(bericht: dict) -> str:
    """Ein Satz fuer das Modell – oder nichts, wenn es nichts zu sagen gibt.

    Der Hinweis nennt AUSDRUECKLICH auch die Schwaeche der Texterkennung. Am
    gemeldeten PDF gemessen liefert OCR zwar deutlich mehr brauchbare Adressen,
    aber eigene Lesefehler ("auftrag@ibsvS.de" statt "ibsv3", "1ab@" statt
    "lab@"). Wer nur meldet "wurde per OCR gelesen", suggeriert einen sauberen
    Text – und das Modell uebernimmt solche Adressen dann ungeprueft.
    """
    if not bericht.get("ocr"):
        return ""
    t = (f"[Hinweis zur Textqualitaet: Die Textebene dieses PDFs ist beschaedigt "
         f"(Zeichenfehler wie 'O' statt '0'). {bericht['seiten_ocr']} von "
         f"{bericht['seiten_gesamt']} Seiten wurden deshalb per Texterkennung neu "
         f"gelesen – das ist deutlich besser, kann aber eigene Lesefehler "
         f"enthalten (Ziffern und E-Mail-Adressen).")
    ab = bericht.get("nicht_geprueft_ab")
    if ab:
        t += (f" Ab Seite {ab} stammt der Text weiterhin aus der beschaedigten "
              f"Textebene.")
    return t + (" Zahlen und Adressen vor der Weitergabe am Original pruefen "
                "und Unsicherheiten benennen, statt sie zu glaetten.]")


# Obergrenze fuer extrahierten Text pro Datei. Ein einzelnes grosses
# Datenmodell-PDF (z.B. "NEXUS KIS Datenmodell – Tabellen", 9 MB) erzeugt sonst
# zig MB Text → hunderttausende Woerter → tausende Chunks → mehrere GB RAM und
# OOM. Darueber hinaus bringt Volltext kaum zusaetzlichen Trefferwert.
MAX_EXTRACT_CHARS = 4_000_000   # ~4 MB Text ≈ max. ~3000 Chunks


def _unlesbar_grund(filepath: Path, max_bytes: int) -> str:
    """Warum lieferte die Extraktion nichts? Antwort in Klartext.

    ``_extract_text_raw`` gibt bei JEDEM Problem schlicht ``None`` zurueck – zu
    gross, unbekanntes Format, kaputte Datei, fehlender Parser sind dort nicht
    unterscheidbar. Fuer die Oberflaeche wird der Grund deshalb hier nachtraeglich
    ermittelt: „4 Dateien fehlgeschlagen" ohne das WARUM ist keine Meldung,
    sondern eine Aufgabe fuer den Leser.

    Rueckgabe ist absichtlich handlungsleitend formuliert (was ist zu tun?),
    nicht als Ausnahmetext.
    """
    try:
        groesse = filepath.stat().st_size
    except OSError as e:
        return f"nicht lesbar ({e.strerror or type(e).__name__})"
    if groesse == 0:
        return "Datei ist leer (0 Byte)"
    if groesse > max_bytes:
        return (f"zu gross: {groesse / 1048576:.1f} MB, erlaubt sind "
                f"{max_bytes / 1048576:.0f} MB (Einstellungen → Skills → Wissen)")
    suffix = filepath.suffix.lower()
    if not suffix:
        return "ohne Dateiendung – Format nicht bestimmbar"
    if suffix not in alle_endungen():
        if suffix == ".onetoc2":
            return ("OneNote-Notizbuch-Index – enthaelt keinen Inhalt, nur die "
                    "Abschnittsliste. Indiziert werden die .one-Dateien daneben.")
        return f"Format {suffix} wird nicht unterstuetzt"
    # Bekanntes Format, passende Groesse -> der Parser selbst ist gescheitert.
    # Haeufigste Ursachen: beschaedigte Datei, passwortgeschuetztes PDF, oder ein
    # rein bildbasiertes PDF ohne OCR-Paket.
    if suffix in EXTENSIONS_PDF:
        return ("PDF nicht lesbar – beschaedigt, passwortgeschuetzt oder ein "
                "reines Scan-PDF ohne Text (dann OCR noetig: tesseract-ocr)")
    if suffix in EXTENSIONS_IMAGE:
        return "Bild ohne erkennbaren Text (OCR benoetigt tesseract-ocr)"
    if suffix in (EXTENSIONS_VIDEO | EXTENSIONS_AUDIO):
        return "Audio/Video ohne Transkript (benoetigt faster-whisper)"
    if suffix in EXTENSIONS_ONENOTE:
        # Der Grund steht im Extraktor: fehlt Java/Tika, ist das die Ursache
        # und der Text nennt den Weg. Sonst war die Datei wirklich leer.
        try:
            from backend.tools.onenote import fehlender_baustein
            hinweis = fehlender_baustein()
        except Exception:  # noqa: BLE001
            hinweis = ""
        return hinweis or ("OneNote-Abschnitt ohne auslesbaren Text (leer, nur "
                           "Handschrift oder nur Bilder ohne erkannten Text)")
    return f"{suffix}-Datei liess sich nicht auslesen (beschaedigt?)"


def _extract_text(filepath: Path, max_bytes: int) -> str | None:
    """Extrahiert Text (Text/PDF/DOCX/XLSX/PPTX/Bild-OCR/Video/Audio) und deckelt
    die Laenge, damit ein einzelnes Riesendokument nicht den Speicher sprengt."""
    text = _extract_text_raw(filepath, max_bytes)
    if text and len(text) > MAX_EXTRACT_CHARS:
        _log.warning(f"Extrahierter Text gekappt ({len(text)} → {MAX_EXTRACT_CHARS} Zeichen): {filepath.name}")
        text = text[:MAX_EXTRACT_CHARS]
    return text


def _extract_text_raw(filepath: Path, max_bytes: int) -> str | None:
    """Rohe Extraktion (ohne Laengen-Deckelung – die macht ``_extract_text``)."""
    try:
        if filepath.stat().st_size > max_bytes:
            return None
    except Exception:
        return None

    suffix = filepath.suffix.lower()

    if suffix in EXTENSIONS_TEXT:
        try:
            return filepath.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return None

    if suffix in EXTENSIONS_PDF:
        return pdf_text_mit_bericht(filepath)[0]

    if suffix in EXTENSIONS_DOCX:
        try:
            import docx
            doc = docx.Document(str(filepath))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paras) or None
        except ImportError:
            return None
        except Exception:
            return None

    return _extract_text_rest(filepath, max_bytes)


def pdf_text_mit_bericht(filepath: Path) -> tuple[str | None, dict]:
    """PDF-Text samt Qualitaetsbericht.

    Eigene Funktion, weil der Aufrufer im Chat den Bericht BRAUCHT (er wird dem
    Modell als Hinweis mitgegeben). Ueber eine ContextVar ginge das nicht:
    ``asyncio.to_thread`` uebergibt eine KOPIE des Kontextes, ein ``set()`` im
    Thread ist im Aufrufer nicht sichtbar.
    """
    bericht: dict = {"geprueft": False, "verdacht": {}, "ocr": False,
                     "seiten_ocr": 0, "seiten_gesamt": 0, "grund": ""}
    try:
        import gc
        import pdfplumber
        texts = []          # nur Seiten MIT Text – fuer die Ausgabe
        seiten = []         # ALLE Seiten, Index 0 = Seite 1 – fuer die Pruefung
        total = 0
        with pdfplumber.open(str(filepath)) as pdf:
            # ueber den INDEX laufen, nicht ueber `for p in pdf.pages`:
            # nur so laesst sich der Platz in der Liste danach freigeben.
            seitenzahl = len(pdf.pages)
            for i in range(seitenzahl):
                p = pdf.pages[i]
                t = p.extract_text()
                # DREI Freigaben, und alle drei werden gebraucht:
                #   flush_cache() – leert den Layout-Cache der Seite
                #   close()       – gibt den pdfminer-Layoutbaum frei
                #   _pages[i]     – die Seiten-INSTANZ selbst; pdfplumber
                #                   haelt jede erzeugte Seite dauerhaft in
                #                   dieser Liste fest.
                # Der dritte Punkt ist der entscheidende und war bis
                # 2026-08-01 nicht da. Nachgemessen an einem 130-MB-PDF
                # (KBV-Bewertungsmassstab, ~9000 Seiten) auf ECHT:
                #   nur flush_cache()          -> Spitze 6000 MB
                #   zusaetzlich close+_pages   -> Spitze  306 MB
                # bei identischem Text (4.000.122 Zeichen) und gleicher
                # Dauer (251 s gegen 243 s). Faktor 20 – der Unterschied
                # zwischen "laeuft" und "OOM-Kill" auf dieser VM.
                try:
                    p.flush_cache()
                    p.close()
                except Exception:
                    pass
                try:
                    pdf._pages[i] = None      # privat, deshalb abgesichert
                except Exception:
                    pass
                # LUECKENLOS mitfuehren: die Qualitaetspruefung ersetzt
                # Seiten einzeln und braucht dafuer Index == Seitennummer-1.
                # Wuerden leere Seiten uebersprungen, landete der OCR-Text
                # von Seite 7 auf Seite 5.
                seiten.append(t or "")
                if t:
                    texts.append(t)
                    total += len(t) + 2
                    if total >= MAX_EXTRACT_CHARS:
                        _log.warning(f"PDF-Extraktion bei {MAX_EXTRACT_CHARS} Zeichen "
                                     f"gestoppt (grosses Dokument): {filepath.name}")
                        break
                # Der Zyklen-Sammler laeuft sonst erst spaet; bei tausenden
                # Seiten haelt das die Spitze unnoetig hoch.
                if i and i % 200 == 0:
                    gc.collect()
        combined = "\n\n".join(texts)
        bericht["seiten_gesamt"] = len(seiten)
        # OCR-Fallback bei gescannten/bildbasierten PDFs (kein/zu wenig Text-Layer)
        if len(combined.strip()) < 80:
            ocr = _ocr_pdf_bytes(filepath.read_bytes())
            if len(ocr.strip()) > len(combined.strip()):
                bericht.update({"geprueft": True, "ocr": True,
                                "seiten_ocr": len(seiten),
                                "grund": "kein Text-Layer – vollstaendig per OCR gelesen"})
                return (ocr or None), bericht
            return (combined or None), bericht
        # Text ist da – aber ist er auch BRAUCHBAR? Die Schwelle oben greift
        # nur bei gescannten Seiten; eine beschaedigte Zeichentabelle liefert
        # reichlich Text und trotzdem Unsinn (siehe Vorfall 2026-08-12).
        try:
            geprueft, bericht = pdf_qualitaet_sichern(filepath.read_bytes(), seiten)
            if bericht.get("ocr"):
                return ("\n\n".join(s for s in geprueft if s) or None), bericht
        except Exception as e:
            # Eine gescheiterte Qualitaetspruefung darf die Extraktion nicht
            # mitreissen – der Text von vorhin ist immer noch besser als nichts.
            _log.warning("PDF-Qualitaetspruefung fehlgeschlagen (%s): %s",
                         filepath.name, e)
            bericht["grund"] = f"Pruefung fehlgeschlagen: {e}"
        return (combined or None), bericht
    except ImportError:
        bericht["grund"] = "pdfplumber fehlt"
        return None, bericht
    except Exception as e:
        bericht["grund"] = f"Extraktion fehlgeschlagen: {e}"
        return None, bericht


def _extract_text_rest(filepath: Path, max_bytes: int) -> str | None:
    """Alle uebrigen Formate (XLSX/PPTX/Bild/Audio/Video)."""
    suffix = filepath.suffix.lower()

    if suffix in EXTENSIONS_XLSX:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            sheets_text = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
                if rows:
                    header = f"[Sheet: {ws.title}]"
                    sheets_text.append(header + "\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(sheets_text) or None
        except ImportError:
            return None
        except Exception:
            return None

    if suffix in EXTENSIONS_PPTX:
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            if any(cells):
                                texts.append("\t".join(cells))
                if texts:
                    slides_text.append(f"[Folie {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides_text) or None
        except ImportError:
            return None
        except Exception:
            return None

    if suffix in EXTENSIONS_ONENOTE:
        # Der Grund wird hier nur protokolliert – der Aufrufer holt ihn bei
        # Bedarf ueber _failure_reason. Ein zweiter Rueckgabewert haette jede
        # Aufrufstelle von _extract_text_raw angefasst.
        from backend.tools.onenote import text_aus_datei
        text, grund = text_aus_datei(filepath)
        if text is None:
            _log.info(f"OneNote nicht gelesen: {filepath.name} – {grund}")
        else:
            _log.debug(f"OneNote gelesen: {filepath.name} – {grund}")
        return text

    if suffix in EXTENSIONS_IMAGE:
        return _ocr_image(filepath)

    if suffix in EXTENSIONS_VIDEO | EXTENSIONS_AUDIO:
        # Video/Audio: max_bytes-Check großzügiger (200MB Default für Medien)
        media_max = max(max_bytes, 200 * 1024 * 1024)
        try:
            if filepath.stat().st_size > media_max:
                _log.warning(f"Mediendatei zu groß: {filepath} ({filepath.stat().st_size / 1024 / 1024:.0f} MB)")
                return None
        except Exception:
            return None
        return _transcribe_media(filepath)

    return None


def _tokenize(text: str) -> list[str]:
    return re.findall(r'\b\w{2,}\b', text.lower())


def _chunk_text(text: str, chunk_size: int = 200, overlap: int = 40) -> list[str]:
    """Zerlegt Text in ueberlappende Wort-Chunks.

    chunk_size MUSS zum Embedding-Modell passen: multilingual-e5-small hat ein
    Limit von 512 Tokens und schneidet laengere Chunks stillschweigend ab. Ein
    800-Wort-Chunk deutscher Fachtexte sind ~2000 Tokens – davon waren 75%
    unsichtbar (gemessen: die Beschreibung von @STR_UCASE lag hinter dem
    Cut-off und war ueber die Vektorsuche nicht auffindbar). 200 Woerter bleiben
    mit Reserve unter dem Limit und schaerfen zugleich das Ranking, weil ein
    Chunk dann ein Thema behandelt statt eines halben Kapitels.
    """
    words = text.split()
    if len(words) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return chunks


def _load_cache() -> dict:
    try:
        if INDEX_CACHE_PATH.exists():
            data = json.loads(INDEX_CACHE_PATH.read_text(encoding="utf-8"))
            if data.get("version") == 1:
                return data
    except Exception:
        pass
    return {"version": 1, "files": {}}


def _is_pending_path(p) -> bool:
    """True fuer den internen Entwurfs-Speicher (data/knowledge/pending/*.json).
    Diese Extraktor-Entwuerfe sind KEIN Wissen und duerfen weder indiziert noch
    in der Dokument-/Gruppenliste auftauchen."""
    return "data/knowledge/pending/" in str(p or "").replace("\\", "/")


def _indexed_rel_paths() -> list:
    """Alle Datei-Pfade, die im INDEX (lokale Wissensdatenbank) stehen.

    Quelle: TF-IDF-Cache (``knowledge_index.json``) + FAISS-Meta
    (``faiss_meta.json``) – BEIDES lokale Dateien. Es wird KEIN Datei-Share
    durchlaufen; die Funktion ist damit immer schnell und unabhaengig davon,
    ob ein Netzlaufwerk erreichbar ist. Genau das ist die richtige Quelle fuer
    die Gruppen-Zaehler (die Gruppen sind logische Tags auf DB-Eintraegen)."""
    paths = set()
    try:
        for p in _load_cache().get("files", {}).keys():
            if p and not _is_pending_path(p):
                paths.add(p)
    except Exception:
        pass
    try:
        _meta = PROJECT_ROOT / "data" / "vector_store" / "faiss_meta.json"
        if _meta.exists() and _meta.stat().st_size > 10:
            for m in json.loads(_meta.read_text(encoding="utf-8")):
                fp = m.get("file_path")
                if fp and not _is_pending_path(fp):
                    paths.add(fp)
    except Exception:
        pass
    return list(paths)


def known_paths_with_disk() -> list:
    """Index-Pfade PLUS aktuell auf der Platte liegende Wissensdateien.

    Gemeinsame Zaehl-/Listen-Basis fuer die Wissensgruppen: Der Index allein
    hinkt der Platte hinterher (z.B. Pending-Extraktor-JSONs), wodurch
    Gruppen-Zaehler kleiner ausfielen als die tatsaechliche Dokumentliste.
    Der Disk-Teil laeuft best-effort – tote Netz-Shares faengt
    _all_files/_safe_exists ab, bei Fehlern bleibt es beim Index."""
    paths = set(_indexed_rel_paths())
    try:
        for f in _all_files(_get_folders()):
            paths.add(str(f))
    except Exception:
        pass
    # Versteckte/interne Dateien ausschliessen – faengt auch evtl. frueher
    # indizierte Alt-Eintraege ab (z.B. das Manifest .groups.json oder die
    # Entwurfs-JSONs unter data/knowledge/pending/).
    return [p for p in paths
            if not os.path.basename(p).startswith(".") and not _is_pending_path(p)]


# Kleiner mtime-Cache fuer den Disk-Scan der Inhalts-Suche
_scan_cache: dict = {}          # path_str -> (mtime, text_lower)
_SCAN_MAX_BYTES = 2_000_000     # groessere Dateien werden beim Disk-Scan uebersprungen
_SCAN_CACHE_BYTES = 262_144     # nur Dateien bis 256 KB im RAM cachen


def content_search_paths(needle: str) -> list:
    """Substring-Suche (case-insensitive) ueber den INHALT der Wissensdateien.

    Quellen (in dieser Reihenfolge):
    1. TF-IDF-Cache-Chunks und FAISS-Meta (bereits extrahierte Texte – deckt
       auch PDF/DOCX/OCR-Inhalte ab, sofern indexiert)
    2. Textformate (.json/.md/.txt/...) zusaetzlich direkt von der Platte –
       deckt neue/noch nicht indexierte Dateien ab, z.B. Pending-Extraktor-
       JSONs. Tote Netz-Shares faengt _all_files/_safe_exists ab.

    Gibt relative Pfade zurueck."""
    needle = (needle or "").strip().lower()
    if len(needle) < 2:
        return []
    hits = set()
    try:
        for path_str, entry in _load_cache().get("files", {}).items():
            for ch in entry.get("chunks") or []:
                if needle in ch.lower():
                    hits.add(path_str)
                    break
    except Exception:
        pass
    try:
        _meta = PROJECT_ROOT / "data" / "vector_store" / "faiss_meta.json"
        if _meta.exists() and _meta.stat().st_size > 10:
            for m in json.loads(_meta.read_text(encoding="utf-8")):
                fp = m.get("file_path")
                if fp and fp not in hits and needle in (m.get("text") or "").lower():
                    hits.add(fp)
    except Exception:
        pass
    # Disk-Scan fuer Textformate (Index kann hinter der Platte herhinken)
    try:
        for f in _all_files(_get_folders()):
            path_str = str(f)
            if path_str in hits or f.suffix.lower() not in EXTENSIONS_TEXT:
                continue
            try:
                st = f.stat()
                if st.st_size > _SCAN_MAX_BYTES:
                    continue
                cached = _scan_cache.get(path_str)
                if cached and cached[0] == st.st_mtime:
                    text = cached[1]
                else:
                    text = f.read_text(encoding="utf-8", errors="ignore").lower()
                    if st.st_size <= _SCAN_CACHE_BYTES:
                        if len(_scan_cache) > 2000:
                            _scan_cache.clear()
                        _scan_cache[path_str] = (st.st_mtime, text)
            except Exception:
                continue
            if needle in text:
                hits.add(path_str)
    except Exception:
        pass
    out = set()
    for p in hits:
        try:
            out.add(str(Path(p).resolve().relative_to(PROJECT_ROOT)))
        except Exception:
            out.add(str(p))
    return sorted(out)


def _save_cache(cache: dict):
    try:
        INDEX_CACHE_PATH.parent.mkdir(parents=True, exist_ok=True)
        INDEX_CACHE_PATH.write_text(
            json.dumps(cache, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    except Exception:
        pass


# ─── Ordner-Verwaltung: Index-Relokation / -Bereinigung ──────────────────────
# Indizierte Dokumente sind nur ueber ihren absoluten Dateipfad mit dem
# Quellordner verknuepft (TF-IDF-Cache-Schluessel + FAISS file_path). Beim
# Umbenennen/Loeschen eines Wissens-Ordners muessen daher beide Indizes und
# die Gruppen-Zuordnungen (relative Pfade in .groups.json) mitgezogen werden.

def _folder_rel(folder: Path) -> str:
    try:
        return str(folder.relative_to(PROJECT_ROOT))
    except ValueError:
        return str(folder)


def relocate_folder_index(old_folder: Path, new_folder: Path) -> dict:
    """Schreibt nach einer Ordner-Umbenennung alle Index-Eintraege um:
    TF-IDF-Cache, FAISS-Metadaten (ohne Neu-Embedding) und Gruppen-Zuordnungen.
    Gibt Zaehler der verschobenen Eintraege zurueck."""
    old_s, new_s = str(old_folder), str(new_folder)

    moved_tfidf = 0
    with _cache_lock:
        cache = _load_cache()
        files = cache.get("files", {})
        renamed = {}
        for p, entry in files.items():
            if p.startswith(old_s + os.sep):
                renamed[new_s + p[len(old_s):]] = entry
                moved_tfidf += 1
            else:
                renamed[p] = entry
        if moved_tfidf:
            cache["files"] = renamed
            _save_cache(cache)

    moved_vec = 0
    vs = _get_vector_store()
    if vs is not None:
        try:
            moved_vec = vs.rename_path_prefix(old_s, new_s)
        except Exception as e:
            _log.warning(f"FAISS-Relokation fehlgeschlagen: {e}")

    moved_groups = 0
    try:
        from backend import knowledge_groups as kg
        moved_groups = kg.relocate_prefix(_folder_rel(old_folder), _folder_rel(new_folder))
    except Exception as e:
        _log.warning(f"Gruppen-Relokation fehlgeschlagen: {e}")

    _log.info(f"Ordner-Index relokalisiert {old_s} -> {new_s}: "
              f"{moved_tfidf} TF-IDF-Dateien, {moved_vec} Vektor-Chunks, {moved_groups} Gruppen-Zuordnungen")
    return {"tfidf_files": moved_tfidf, "vector_chunks": moved_vec,
            "group_assignments": moved_groups}


def purge_folder_index(folder: Path) -> dict:
    """Entfernt alle Index-Eintraege (TF-IDF + FAISS) und Gruppen-Zuordnungen
    unterhalb eines Ordners. Gibt Zaehler der entfernten Eintraege zurueck."""
    folder_s = str(folder)

    removed_tfidf = 0
    with _cache_lock:
        cache = _load_cache()
        files = cache.get("files", {})
        keep = {p: e for p, e in files.items() if not p.startswith(folder_s + os.sep)}
        removed_tfidf = len(files) - len(keep)
        if removed_tfidf:
            cache["files"] = keep
            _save_cache(cache)

    removed_vec = 0
    vs = _get_vector_store()
    if vs is not None:
        try:
            removed_vec = vs.remove_path_prefix(folder_s)
        except Exception as e:
            _log.warning(f"FAISS-Bereinigung fehlgeschlagen: {e}")

    removed_groups = 0
    try:
        from backend import knowledge_groups as kg
        removed_groups = kg.remove_prefix(_folder_rel(folder))
    except Exception as e:
        _log.warning(f"Gruppen-Bereinigung fehlgeschlagen: {e}")

    _log.info(f"Ordner-Index bereinigt {folder_s}: "
              f"{removed_tfidf} TF-IDF-Dateien, {removed_vec} Vektor-Chunks, {removed_groups} Gruppen-Zuordnungen")
    return {"tfidf_files": removed_tfidf, "vector_chunks": removed_vec,
            "group_assignments": removed_groups}


def purge_file_index(file: Path) -> dict:
    """Entfernt eine EINZELNE Datei restlos aus dem Index: TF-IDF-Cache, FAISS
    und ihre Gruppen-Zuordnung. Einzeldatei-Pendant zu ``purge_folder_index`` –
    wird beim Loeschen einer Wissensdatei aufgerufen, damit die Datei nicht als
    Karteileiche in der Zaehl-Basis (``known_paths_with_disk``) bzw. den
    Wissensgruppen zurueckbleibt. Gibt Zaehler der entfernten Eintraege zurueck."""
    file_s = str(file)

    removed_tfidf = 0
    with _cache_lock:
        cache = _load_cache()
        files = cache.get("files", {})
        if file_s in files:
            del files[file_s]
            removed_tfidf = 1
            _save_cache(cache)

    removed_vec = 0
    vs = _get_vector_store()
    if vs is not None:
        try:
            # remove_file() liefert die Anzahl selbst – frueher wurde dafuer
            # ungeschuetzt auf die private Liste vs._meta zugegriffen.
            removed_vec = vs.remove_file(file_s)
        except Exception as e:
            _log.warning(f"FAISS-Bereinigung fehlgeschlagen: {e}")

    removed_group = False
    try:
        from backend import knowledge_groups as kg
        if kg.get_assignment(file_s):
            kg.set_assignment(file_s, [])  # leere Liste = Zuordnung entfernen
            removed_group = True
    except Exception as e:
        _log.warning(f"Gruppen-Bereinigung fehlgeschlagen: {e}")

    _log.info(f"Datei-Index bereinigt {file_s}: "
              f"{removed_tfidf} TF-IDF, {removed_vec} Vektor-Chunks, "
              f"Gruppen-Zuordnung={'ja' if removed_group else 'nein'}")
    return {"tfidf_files": removed_tfidf, "vector_chunks": removed_vec,
            "group_assignment": removed_group}


def relocate_file_index(old_file: Path, new_file: Path) -> dict:
    """Zieht die Index-Eintraege EINER verschobenen Datei auf den neuen Pfad um –
    ohne Neu-Embedding. Einzeldatei-Pendant zu ``relocate_folder_index``.

    Betrifft TF-IDF-Cache-Schluessel, FAISS-Metadaten und die Wissensgruppen-
    Zuordnung. Die Datei selbst muss vom Aufrufer bereits verschoben worden sein
    (``Path.rename()``), damit mtime und Inhalt unveraendert bleiben und der
    naechste inkrementelle Reindex sie nicht erneut einbettet.

    Gibt Zaehler der umgezogenen Eintraege zurueck.
    """
    old_s, new_s = str(old_file), str(new_file)
    if old_s == new_s:
        return {"tfidf_files": 0, "vector_chunks": 0, "group_assignment": False}

    moved_tfidf = 0
    with _cache_lock:
        cache = _load_cache()
        files = cache.get("files", {})
        if old_s in files:
            files[new_s] = files.pop(old_s)
            moved_tfidf = 1
            _save_cache(cache)

    moved_vec = 0
    vs = _get_vector_store()
    if vs is not None:
        try:
            moved_vec = vs.rename_file_path(old_s, new_s)
        except Exception as e:
            _log.warning(f"FAISS-Relokation fehlgeschlagen: {e}")

    # Gruppen-Zuordnung mitnehmen. Modell A: Dateien in einem Ordner erben
    # dessen Gruppen; eine EXPLIZITE Zuordnung (Modell B) haengt dagegen am
    # relativen Dateipfad und muss aktiv umgehaengt werden.
    moved_group = False
    try:
        from backend import knowledge_groups as kg
        old_rel, new_rel = _folder_rel(old_file), _folder_rel(new_file)
        groups = kg.get_assignment(old_rel)
        if groups:
            kg.set_assignment(new_rel, groups)
            kg.set_assignment(old_rel, [])
            moved_group = True
    except Exception as e:
        _log.warning(f"Gruppen-Relokation fehlgeschlagen: {e}")

    _log.info(f"Datei-Index verschoben {old_s} -> {new_s}: "
              f"{moved_tfidf} TF-IDF, {moved_vec} Vektor-Chunks, "
              f"Gruppen-Zuordnung={'ja' if moved_group else 'nein'}")
    return {"tfidf_files": moved_tfidf, "vector_chunks": moved_vec,
            "group_assignment": moved_group}


def _ist_office_hilfsdatei(name: str) -> bool:
    """Office-Sperr- und Wiederherstellungsdateien erkennen.

    Word, Excel und PowerPoint legen neben jedem GEOEFFNETEN Dokument eine
    Datei ``~$name.docx`` an (Besitzerkennung, wenige hundert Byte). Sie traegt
    die Endung eines Office-Dokuments, ist aber keines – python-docx scheitert
    daran zwangslaeufig.

    GEFUNDEN AM 2026-07-31 auf ECHT: von den vier als „fehlgeschlagen"
    gemeldeten Dateien waren DREI solche Sperrdateien. Sie erzeugten bei jedem
    Indexlauf eine Fehlermeldung, die nach einem Problem aussah, aber keines
    war – und die echte vierte Meldung (ein 130-MB-PDF ueber dem Limit) darin
    untergehen liess. Ein Fehlerzaehler, der zu drei Vierteln aus Rauschen
    besteht, wird nicht mehr gelesen.

    Sie verschwinden von selbst, sobald das Dokument geschlossen wird – auf
    einem Netzlaufwerk mit vielen Bearbeitern aber praktisch nie vollstaendig.
    """
    return name.startswith("~$") or name.startswith("~WRL")


def _all_files(folders: list[Path]) -> list[Path]:
    """Gibt alle unterstützten Dateien in den konfigurierten Ordnern zurück.

    EXTENSIONS_IMAGE gehoert dazu: Der Upload nimmt Bilder an und `_extract_text`
    liest sie per OCR. Fehlten sie hier, wurde ein hochgeladenes Bild nie
    indiziert UND beim naechsten Reindex verlor es seine Gruppen-Zuordnung
    (``kg.prune`` arbeitet auf dieser Liste) – es verschwand damit aus
    „Meine Dateien", obwohl es unveraendert auf der Platte lag.
    """
    all_exts = alle_endungen()
    files = []
    for folder in folders:
        # Totes Netzlaufwerk nicht anfassen -> sonst blockiert os.walk minutenlang.
        if not _safe_exists(folder):
            continue
        try:
            for root, dirs, fs in os.walk(folder, onerror=lambda e: None):
                # Versteckte Verzeichnisse nicht betreten (z.B. .git, .cache) und
                # den internen Entwurfs-Speicher (data/knowledge/pending) auslassen –
                # Extraktor-Entwuerfe sind KEIN Wissensdokument.
                dirs[:] = [d for d in dirs if not d.startswith(".")
                           and not _is_pending_path(os.path.join(root, d) + "/")]
                for f in fs:
                    # Versteckte/interne Dateien ueberspringen – z.B. das
                    # Gruppen-Manifest data/knowledge/.groups.json ist KEIN
                    # Wissensdokument und darf weder indiziert noch gelistet werden.
                    if f.startswith("."):
                        continue
                    if _ist_office_hilfsdatei(f):
                        continue
                    if Path(f).suffix.lower() in all_exts:
                        files.append(Path(root) / f)
        except OSError as e:
            _log.warning("Ordner konnte nicht durchsucht werden (übersprungen): %s (%s)", folder, e)
            continue
    _disk_count_cache.update(value=len(files), ts=time.time())
    return files


# ─── Anzahl indizierbarer Dateien auf der Platte ─────────────────────────────
# Die Statistik-Kachel "Dateien" zeigt, was VORHANDEN ist – nicht, was im Index
# steht (das ist die Kachel "Indiziert"). Der Walk ueber mehrere hundert Dateien
# inkl. Netzlaufwerk darf den Stats-Aufruf aber nicht blockieren, deshalb:
# gecacht, Aktualisierung im Hintergrund, erster Aufruf mit hartem Timeout.
_disk_count_cache: dict = {"value": None, "ts": 0.0}
_DISK_COUNT_TTL = 60.0
_disk_count_refreshing = threading.Event()


def _refresh_disk_count() -> None:
    try:
        _all_files(_get_folders())   # aktualisiert _disk_count_cache selbst
    except Exception as e:
        _log.debug(f"Datei-Zaehlung fehlgeschlagen: {e}")
    finally:
        _disk_count_refreshing.clear()


def get_disk_file_count() -> int | None:
    """Anzahl indizierbarer Dateien in den Wissensordnern (None = noch unbekannt)."""
    cached = _disk_count_cache["value"]
    fresh = cached is not None and (time.time() - _disk_count_cache["ts"]) < _DISK_COUNT_TTL
    if fresh:
        return cached
    if cached is not None:
        # Alten Wert sofort ausliefern, im Hintergrund neu zaehlen.
        if not _disk_count_refreshing.is_set():
            _disk_count_refreshing.set()
            threading.Thread(target=_refresh_disk_count, daemon=True).start()
        return cached
    # Erster Aufruf: kurz warten, danach greift der Cache. Laeuft der Walk in den
    # Timeout, fuellt der (weiterlaufende) Daemon-Thread den Cache trotzdem –
    # der naechste Aufruf hat den Wert dann sofort.
    return _bounded_call(lambda: len(_all_files(_get_folders())), timeout=5.0, default=None)


# Kurz gecachte Dateiliste fuer HEISSE Pfade (Suche). Der Ordner-Scan laeuft
# ueber alle Wissensordner inkl. Netzlaufwerken; ihn pro Suche mehrfach
# auszufuehren war der groesste vermeidbare Posten im Suchpfad.
# Der ausdrueckliche Neuaufbau ruft weiterhin _all_files() direkt – er MUSS den
# echten Stand sehen.
_FILES_TTL = 30.0
_files_cache: dict = {"key": None, "value": None, "ts": 0.0}
_files_cache_lock = threading.Lock()


def _all_files_cached(folders: list[Path]) -> list[Path]:
    """``_all_files`` mit kurzer Vorhaltezeit – nur fuer Lesepfade."""
    key = tuple(str(f) for f in folders)
    now = time.time()
    with _files_cache_lock:
        if (_files_cache["key"] == key and _files_cache["value"] is not None
                and now - _files_cache["ts"] < _FILES_TTL):
            return _files_cache["value"]
    value = _all_files(folders)
    with _files_cache_lock:
        _files_cache.update(key=key, value=value, ts=time.time())
    return value


def invalidate_files_cache() -> None:
    """Vorhaltezeit der Dateiliste vorzeitig beenden (nach Upload/Loeschen)."""
    with _files_cache_lock:
        _files_cache.update(key=None, value=None, ts=0.0)


INLINE_LIMIT = 10  # Maximale Dateien die inline (im Suchpfad) indiziert werden

def _rebuild_cache(folders: list[Path], max_bytes: int, force: bool = False) -> dict:
    """Inkrementeller TF-IDF Index-Aufbau (Thread-sicher).

    force=False (Suchpfad): Kein Bulk-Aufbau bei leerem Index, max. INLINE_LIMIT Dateien.
    force=True  (Neu-Indizieren): Alle Dateien werden verarbeitet, kein Limit.
    """
    with _cache_lock:
        cache = _load_cache()
        files = _all_files(folders)
        current_paths = {str(f) for f in files}

        # Gelöschte Dateien entfernen
        for p in list(cache["files"].keys()):
            if p not in current_paths:
                del cache["files"][p]

        # Neue/geänderte Dateien ermitteln
        to_index = []
        for filepath in files:
            path_str = str(filepath)
            try:
                mtime = filepath.stat().st_mtime
            except Exception:
                continue
            cached = cache["files"].get(path_str, {})
            if cached.get("mtime") != mtime:
                to_index.append(filepath)

        if not force:
            # Leerer Index mit vielen Dateien: kein Inline-Bulk-Indexing
            if not cache["files"] and len(to_index) > INLINE_LIMIT:
                _log.debug(f"TF-IDF Index leer ({len(to_index)} Dateien) – bitte Neu-Indizieren ausfuehren")
                return cache
            # Bestehendes Inkrementell: max. INLINE_LIMIT Dateien inline
            if len(to_index) > INLINE_LIMIT:
                _log.info(f"{len(to_index)} geaenderte Dateien – nur {INLINE_LIMIT} inline, Rest via Neu-Indizieren")
                to_index = to_index[:INLINE_LIMIT]

        _set_progress(phase="TF-IDF", done=0, total=len(to_index))

        changed = False
        for i, filepath in enumerate(to_index):
            if _reindex_cancel.is_set():
                _log.info(f"TF-IDF Index: Abbruch nach {i}/{len(to_index)} Dateien")
                break
            path_str = str(filepath)
            _set_progress(done=i + 1, phase=f"TF-IDF: {filepath.name[:40]}")
            try:
                mtime = filepath.stat().st_mtime
                text = _extract_text(filepath, max_bytes)
                if text and text.strip():
                    cache["files"][path_str] = {
                        "mtime": mtime,
                        "chunks": _chunk_text(text),
                        "size": filepath.stat().st_size,
                    }
                else:
                    cache["files"].pop(path_str, None)
                changed = True
            except Exception:
                pass

        if changed:
            _save_cache(cache)

        if not _reindex_cancel.is_set():
            _set_progress(done=len(to_index), total=len(to_index))
        return cache


def _search(query: str, cache: dict, max_results: int) -> list[tuple[float, str, str]]:
    """TF-IDF Suche über alle gecachten Chunks."""
    query_tokens = _tokenize(query)
    if not query_tokens:
        return []

    all_chunks: list[tuple[str, str]] = []
    for path_str, fdata in cache["files"].items():
        for chunk in fdata.get("chunks", []):
            all_chunks.append((path_str, chunk))

    if not all_chunks:
        return []

    doc_count = len(all_chunks)
    doc_freq: Counter = Counter()
    for _, chunk in all_chunks:
        tokens = set(_tokenize(chunk))
        for t in query_tokens:
            if t in tokens:
                doc_freq[t] += 1

    scored: list[tuple[float, str, str]] = []
    for path_str, chunk in all_chunks:
        tokens = _tokenize(chunk)
        if not tokens:
            continue
        tf = Counter(tokens)
        score = sum(
            (tf[qt] / len(tokens)) * (math.log((doc_count + 1) / (doc_freq.get(qt, 0) + 1)) + 1)
            for qt in query_tokens if qt in tf
        )
        if score > 0:
            try:
                rel = str(Path(path_str).relative_to(PROJECT_ROOT))
            except ValueError:
                rel = path_str
            scored.append((score, rel, chunk))

    scored.sort(key=lambda x: x[0], reverse=True)
    return scored[:max_results]


async def rag_search(query: str, max_results: int = 8, groups=None) -> list[tuple[float, str, str]]:
    """Strukturierte RAG-Suche fuer externe Aufrufer (z.B. Support-Assistent).

    Liefert eine Liste von (score, relativer_pfad, chunk). Nutzt denselben
    Vektor-/TF-IDF-Dispatch wie das knowledge_search-Tool, gibt aber Rohdaten
    statt formatiertem Text zurueck.

    ``groups`` (optional): Liste von Gruppen-IDs (Modell B). Ist sie gesetzt,
    werden nur Treffer aus Dateien dieser Gruppen zurueckgegeben (ODER-Filter;
    "ungrouped" moeglich). Die Einschraenkung geht IN die Vektorsuche; nur der
    TF-IDF-Rueckfall filtert noch nachtraeglich (siehe KnowledgeTool.execute).
    """
    query = (query or "").strip()
    if not query:
        return []
    folders = _get_folders()
    max_bytes = _get_max_bytes()
    # "auto": ohne FAISS Rueckfall auf TF-IDF statt gar keiner Suche (siehe
    # gleichlautende Stelle in KnowledgeTool.execute).
    search_mode_cfg = "auto"

    # Ueber-Abfrage nur noch als Rueckfall fuer TF-IDF (Begruendung bei
    # KnowledgeTool.execute: der Nachfilter verliert still Treffer).
    fetch_n = max(max_results * 5, 40) if groups else max_results

    vs = _get_vector_store()
    vector_index_ready = vs is not None and vs.chunk_count() > 0
    need_tfidf_cache = search_mode_cfg == "tfidf" or (
        search_mode_cfg == "auto" and not vector_index_ready)

    # Kein _load_cache() im Normalfall: die Datei traegt alle Chunk-Texte und
    # wurde bisher bei jedem Aufruf komplett geparst.
    cache = await asyncio.to_thread(_rebuild_cache, folders, max_bytes, False) \
        if need_tfidf_cache else None

    results = None
    gruppen_in_suche = False
    if search_mode_cfg in ("auto", "vector"):
        has_vector = await asyncio.to_thread(_rebuild_vector_index, folders, max_bytes)
        if has_vector:
            # Erlaubte Pfade ERST NACH dem Reindex bestimmen – er kann gerade
            # geaenderte Dateien nachgetragen haben.
            allow_paths = None
            if groups and vs is not None:
                try:
                    from backend import knowledge_groups as kg
                    allow_paths = set(kg.filter_paths_by_groups(
                        list(vs.get_indexed_files().keys()), groups))
                except Exception as e:  # noqa: BLE001
                    print(f"[knowledge] Gruppenfilter nicht vorbereitbar, "
                          f"nutze Nachfilter: {e}", flush=True)
                    allow_paths = None
            n = max_results if allow_paths is not None else fetch_n
            results = await asyncio.to_thread(_vector_search, query, n, allow_paths)
            gruppen_in_suche = allow_paths is not None
        elif search_mode_cfg == "auto" and cache:
            results = _search(query, cache, fetch_n)
    elif search_mode_cfg == "tfidf" and cache:
        results = _search(query, cache, fetch_n)
    results = results or []

    # Nachfilter nur, wenn die Suche selbst nicht einschraenken konnte.
    if groups and not gruppen_in_suche:
        try:
            from backend import knowledge_groups as kg
            kept = set(kg.filter_paths_by_groups([r[1] for r in results], groups))
            results = [r for r in results if r[1] in kept][:max_results]
        except Exception:
            results = results[:max_results]
    else:
        results = results[:max_results]
    return results


def _get_static_stats() -> dict:
    """Format-Support + ChromaDB-Client – wird einmalig gecacht (ändert sich nicht)."""
    global _stats_cache
    with _stats_cache_lock:
        if _stats_cache is not None:
            return _stats_cache
        has_pdf = has_docx = has_xlsx = has_pptx = has_video = has_image = False
        try:
            import pdfplumber; has_pdf = True
        except ImportError: pass
        try:
            import docx; has_docx = True
        except ImportError: pass
        try:
            import openpyxl; has_xlsx = True
        except ImportError: pass
        try:
            from pptx import Presentation; has_pptx = True
        except ImportError: pass
        try:
            from faster_whisper import WhisperModel
            if shutil.which("ffmpeg"): has_video = True
        except ImportError: pass
        try:
            import pytesseract  # noqa: F401
            if shutil.which("tesseract"): has_image = True
        except ImportError: pass

        _stats_cache = {
            "pdf_support": has_pdf, "docx_support": has_docx,
            "xlsx_support": has_xlsx, "pptx_support": has_pptx,
            "video_support": has_video, "image_support": has_image,
        }
        return _stats_cache


def _onenote_support() -> bool:
    """Sind Java UND tika-app.jar vorhanden? (Voraussetzung fuer *.one)"""
    try:
        from backend.tools.onenote import finde_java, finde_tika
        return bool(finde_java() and finde_tika())
    except Exception:  # noqa: BLE001
        return False


def get_stats() -> dict:
    """Statistiken für die API – schnell, kein Netzwerk-/Modell-Scan."""
    folders = _get_folders()

    folder_list = []
    for f in folders:
        try:
            rel = str(f.relative_to(PROJECT_ROOT))
        except ValueError:
            rel = str(f)
        folder_list.append({"path": rel, "exists": _safe_exists(f)})

    # Vektor-DB: FAISS verfuegbar? + Index-Inhalt lesen (meta.json)
    vector_db_available = False
    has_vector = False
    vector_chunks = 0
    vector_files = 0
    vector_db_name = ""
    vector_db_version = ""
    vector_model = ""
    faiss_file_paths: set = set()
    faiss_meta_list: list = []
    try:
        import faiss as _faiss
        vector_db_available = True
        vector_db_name = "FAISS"
        vector_db_version = getattr(_faiss, "__version__", "")
        from backend.tools.vector_store import MODEL_NAME as _VS_MODEL
        vector_model = _VS_MODEL
        _meta_path = PROJECT_ROOT / "data" / "vector_store" / "faiss_meta.json"
        if _meta_path.exists() and _meta_path.stat().st_size > 10:
            import json as _json
            with open(_meta_path, "r", encoding="utf-8") as _f:
                faiss_meta_list = _json.load(_f)
            vector_chunks = len(faiss_meta_list)
            has_vector = vector_chunks > 0
            faiss_file_paths = {m["file_path"] for m in faiss_meta_list}
            vector_files = len(faiss_file_paths)
    except Exception:
        pass

    # Datei-/Chunk-Zähler: FAISS-Meta bevorzugen wenn vorhanden, sonst TF-IDF-Cache
    if has_vector:
        total_files = vector_files
        indexed_files = vector_files
        total_chunks = vector_chunks
        # Dateigröße aus Filesystem (FAISS speichert keine Größe). Zeitlich
        # begrenzt, damit ein totes Netzlaufwerk die Stats nicht einfriert.
        def _sum_sizes():
            total = 0
            for p in faiss_file_paths:
                try:
                    total += Path(p).stat().st_size
                except OSError:
                    continue
            return total
        total_size = _bounded_call(_sum_sizes, timeout=3.0, default=0)
    else:
        cache = _load_cache()
        total_files = len(cache["files"])
        indexed_files = len(cache["files"])
        total_chunks = sum(len(d.get("chunks", [])) for d in cache["files"].values())
        total_size = sum(d.get("size", 0) for d in cache["files"].values())

    # "Dateien" = was in den Wissensordnern LIEGT. Frueher stand hier die Anzahl
    # der Dateien im Index – bei einem unvollstaendigen Index sah es dann so aus,
    # als gaebe es nur 10 statt 700+ Dokumente.
    disk_files = get_disk_file_count()
    if disk_files is not None:
        total_files = disk_files

    return {
        "folders": folder_list,
        "total_files": total_files,
        "disk_files": disk_files,
        "indexed_files": indexed_files,
        "total_chunks": total_chunks,
        "total_size_bytes": total_size,
        **_get_static_stats(),
        # BEWUSST NICHT in _get_static_stats(): das dict wird prozessweit
        # gecacht, und die Voraussetzung fuer OneNote sind DATEIEN auf Platte
        # (Java-Binary, tika-app.jar). Ein Administrator, der gerade
        # deploy/tika_setup.sh gefahren hat, saehe die Plakette sonst bis zum
        # Dienstneustart weiter rot und hielte das Skript fuer wirkungslos –
        # dieselbe Falle wie beim Kontext-Schwellwert mit fest verdrahteter 30.
        # Kosten: ein shutil.which und zwei is_file() je Abruf.
        "onenote_support": _onenote_support(),
        # BEWUSST NICHT in _get_static_stats(): das dict wird prozessweit
        # gecacht, und die Voraussetzung fuer OneNote sind DATEIEN auf Platte
        # (Java-Binary, tika-app.jar). Ein Administrator, der gerade
        # deploy/tika_setup.sh gefahren hat, saehe die Plakette sonst bis zum
        # Dienstneustart weiter rot und hielte das Skript fuer wirkungslos –
        # dieselbe Falle wie beim Kontext-Schwellwert mit fest verdrahteter 30.
        # Kosten: ein shutil.which und zwei is_file() je Abruf.
        "vector_db_available": vector_db_available,
        "vector_search": has_vector,
        "vector_files": vector_files,
        "vector_chunks": vector_chunks,
        "vector_db_name": vector_db_name,
        "vector_db_version": vector_db_version,
        "vector_model": vector_model,
        "search_mode": "auto",
        "indexing": get_index_progress()["running"],
        "index_phase": get_index_progress()["phase"],
        "index_failed": get_index_progress().get("failed", 0),
        "last_index_run": get_last_run(),
        # Zustand der Vektorsuche sichtbar machen: ein fehlgeschlagener Aufbau
        # war bisher nur an einer Journal-Zeile erkennbar.
        "vector_store_state": vector_store_status(),
    }


# Ein Lauf, der an einem Fehler scheitert (Embedding-Modell nicht ladbar,
# Netzlaufwerk weg, Speicher voll), hinterlaesst einen LEEREN Index – der
# Neuaufbau beginnt mit vs.clear(). Deshalb automatisch neu ansetzen. Der
# manuelle Abbruch ist davon ausgenommen (siehe _reindex_cancel).
MAX_INDEX_ATTEMPTS = 3      # 1 regulaerer Lauf + 2 automatische Neuversuche
RETRY_DELAY_SEC = 15        # Pause dazwischen (z.B. bis ein Mount zurueck ist)


def force_reindex(resume_count: int = 0, incremental: bool = False,
                  resume_baseline: int = -1) -> dict:
    """Neuaufbau des Wissens-Index:
    - FAISS verfuegbar → nur Vektor-Index (schneller, besser bei 600+ Dateien)
    - FAISS nicht verfuegbar → TF-IDF-Index

    ``incremental=True`` behaelt den bestehenden Index (kein ``vs.clear()``) und
    ergaenzt nur fehlende/geaenderte Dateien – so setzt eine Wiederaufnahme nach
    Absturz dort fort, wo sie war, statt bei 0 zu beginnen.

    Scheitert ein Lauf mit einer Ausnahme, wird er bis zu ``MAX_INDEX_ATTEMPTS``
    mal automatisch wiederholt. ``resume_count`` zaehlt Wiederaufnahmen nach
    einem Prozess-Neustart, ``resume_baseline`` den Dateistand zu deren Beginn
    (fuer die Fortschritts-Pruefung in ``resume_interrupted_reindex``).
    """
    # Re-Entrancy-Schutz: laeuft bereits ein Reindex, NICHT parallel starten
    # (sonst ueberschreiben sich die Fortschritts-Zaehler -> >100%). Stattdessen
    # einen Rerun vormerken, damit neu hinzugefuegte Dateien danach indexiert werden.
    if not _reindex_lock.acquire(blocking=False):
        _reindex_rerun.set()
        _log.info("force_reindex: laeuft bereits – Rerun vorgemerkt")
        return {"skipped": True, "reason": "reindex already running, rerun scheduled"}
    try:
        _reindex_cancel.clear()
        result = _run_with_retries(resume_count, incremental, resume_baseline)
        # Waehrenddessen weitere Anfragen? -> genau einmal nachholen (coalesced).
        # Nach einem Abbruch NICHT nachholen – sonst startet der Lauf, den der
        # Benutzer gerade gestoppt hat, sofort wieder von vorn.
        while _reindex_rerun.is_set() and not _reindex_cancel.is_set():
            _reindex_rerun.clear()
            result = _run_with_retries(resume_count, incremental, resume_baseline)
        return result
    finally:
        _reindex_rerun.clear()
        # Flag zuruecksetzen, sonst wuerde die naechste inline-Indizierung
        # (Suchpfad) den alten Abbruchwunsch erben und sofort abbrechen.
        _reindex_cancel.clear()
        _reindex_lock.release()


def _run_with_retries(resume_count: int = 0, incremental: bool = False,
                      resume_baseline: int = -1) -> dict:
    """Fuehrt den Neuaufbau aus und wiederholt ihn nach einem Fehler automatisch."""
    last_exc: Exception | None = None
    first_started = time.time()
    for attempt in range(1, MAX_INDEX_ATTEMPTS + 1):
        try:
            return _do_force_reindex(attempt=attempt, resume_count=resume_count,
                                     incremental=incremental,
                                     resume_baseline=resume_baseline)
        except Exception as e:
            last_exc = e
            _log.warning(f"Indizierung Versuch {attempt}/{MAX_INDEX_ATTEMPTS} fehlgeschlagen: {e}")
            if _reindex_cancel.is_set() or attempt >= MAX_INDEX_ATTEMPTS:
                break
            # "laeuft" bleibt gesetzt – die Oberflaeche zeigt den Neuversuch an
            # statt den Knopf freizugeben und den Fehler zu verschweigen.
            _set_progress(running=True, phase="Neuversuch", error=str(e),
                          attempt=attempt + 1, max_attempts=MAX_INDEX_ATTEMPTS)
            # Unterbrechbare Pause: ein Abbruch waehrend der Wartezeit greift sofort.
            if _reindex_cancel.wait(RETRY_DELAY_SEC):
                break

    finished = time.time()
    cancelled = _reindex_cancel.is_set()
    # Teilstand erhalten: bei incremental bleibt der Index bestehen, die schon
    # indizierten Dateien sind kein Verlust.
    try:
        vs = _get_vector_store()
        partial_files = vs.file_count() if vs is not None else 0
        partial_chunks = vs.chunk_count() if vs is not None else 0
    except Exception:
        partial_files = partial_chunks = 0
    _set_progress(running=False, phase="Abgebrochen" if cancelled else "Fehler",
                  error=str(last_exc or ""), finished_at=finished, cancelled=cancelled)
    _save_last_run({"started_at": first_started, "finished_at": finished,
                    "status": "cancelled" if cancelled else "error",
                    "error": str(last_exc or "")[:300], "attempts": MAX_INDEX_ATTEMPTS,
                    "resumed": resume_count, "indexed_files": partial_files,
                    "total_chunks": partial_chunks})
    raise last_exc if last_exc else RuntimeError("Indizierung fehlgeschlagen")


def cancel_reindex() -> dict:
    """Bricht einen laufenden Neuaufbau ab (nach der aktuellen Datei).

    Der Index bleibt danach unvollstaendig – ein Neuaufbau leert ihn zuerst.
    """
    if not get_index_progress().get("running"):
        return {"cancelled": False, "reason": "keine Indizierung aktiv"}
    _reindex_cancel.set()
    _reindex_rerun.clear()
    _set_progress(phase="Wird abgebrochen…")
    _log.info("Indizierung: Abbruch angefordert")
    return {"cancelled": True}


# Kurzprotokoll des letzten Laufs – ueberlebt einen Neustart, damit die
# Oberflaeche "Letzter Indexlauf: <Datum/Uhrzeit>" auch nach einem Restart zeigt.
LAST_INDEX_RUN_PATH = PROJECT_ROOT / "data" / "vector_store" / "last_index.json"


def _save_last_run(run: dict) -> None:
    try:
        LAST_INDEX_RUN_PATH.parent.mkdir(parents=True, exist_ok=True)
        LAST_INDEX_RUN_PATH.write_text(json.dumps(run, ensure_ascii=False), encoding="utf-8")
    except Exception as e:
        _log.warning(f"Lauf-Protokoll konnte nicht geschrieben werden: {e}")


def get_last_run() -> dict:
    """Metadaten des letzten Indexlaufs ({} wenn noch nie gelaufen)."""
    try:
        if LAST_INDEX_RUN_PATH.exists():
            return json.loads(LAST_INDEX_RUN_PATH.read_text(encoding="utf-8"))
    except Exception:
        pass
    return {}


def _write_run_checkpoint(done: int, total: int, current_file: str, indexed_files: int) -> None:
    """Schreibt den Zwischenstand des laufenden Reindex auf Platte.

    Stirbt der Prozess danach, weiss die Wiederaufnahme (und die Oberflaeche),
    WIE WEIT es kam und bei WELCHER Datei es zuletzt war."""
    _save_last_run({**_current_run, "finished_at": 0, "status": "running",
                    "done": done, "total": total, "current_file": current_file,
                    "indexed_files": indexed_files})


# Sicherheitsnetz gegen eine echte Endlosschleife: selbst wenn jeder Lauf
# Fortschritt meldet, hoert die automatische Wiederaufnahme nach so vielen
# Anlaeufen auf. Im Normalfall greift vorher die Fortschritts-Pruefung.
MAX_RESUMES = 20


def resume_interrupted_reindex() -> bool:
    """Setzt einen Lauf fort, der durch einen Prozess-Abbruch geendet hat.

    Beim Start aufrufen: steht im Lauf-Protokoll noch ``status: running``, ist
    der Prozess mittendrin gestorben (Neustart, Absturz, OOM – z.B. der
    OOM-Killer). Der Index ist dann unvollstaendig, ohne dass es irgendwo als
    Fehler auftaucht.

    Die Wiederaufnahme laeuft INKREMENTELL (kein ``vs.clear()``): die bereits
    indizierten Dateien bleiben erhalten, es werden nur die fehlenden ergaenzt.
    Fortgesetzt wird nur, solange messbarer Fortschritt entsteht – bringt ein
    Anlauf keine neue Datei in den Index, wird abgebrochen (sonst liefe eine
    Datei, die den Prozess zuverlaessig killt, endlos in dieselbe Wand).

    Gibt True zurueck, wenn eine Wiederaufnahme gestartet wurde.
    """
    run = get_last_run()
    if run.get("status") != "running":
        return False
    if get_index_progress().get("running"):
        return False   # laeuft bereits (z.B. durch Auto-Mount angestossen)

    # Aktuellen Stand aus dem Index lesen (ueberlebt den Absturz auf Platte).
    try:
        vs = _get_vector_store()
        current_files = vs.file_count() if vs is not None else 0
    except Exception:
        current_files = 0

    resumed = int(run.get("resumed") or 0) + 1
    baseline = int(run.get("resume_baseline", -1))   # Stand zu Beginn des letzten Anlaufs
    stalled = baseline >= 0 and current_files <= baseline

    if stalled or resumed > MAX_RESUMES:
        reason = ("kein Fortschritt seit letztem Anlauf – vermutlich scheitert eine "
                  "bestimmte Datei" if stalled else f"{MAX_RESUMES} Anlaeufe erschoepft")
        run.update(status="interrupted", finished_at=run.get("finished_at") or time.time(),
                   indexed_files=current_files, interrupt_reason=reason)
        _save_last_run(run)
        _log.warning(f"Wiederaufnahme der Indizierung gestoppt: {reason} "
                     f"(bei {current_files} Dateien)")
        return False

    _log.warning(
        f"Unterbrochene Indizierung gefunden (Start "
        f"{time.strftime('%d.%m.%Y %H:%M:%S', time.localtime(run.get('started_at') or 0))}, "
        f"zuletzt {current_files} Dateien im Index) – wird inkrementell fortgesetzt "
        f"(Anlauf {resumed})")

    def _run():
        try:
            force_reindex(resume_count=resumed, incremental=True,
                          resume_baseline=current_files)
        except Exception as e:
            _log.error(f"Automatisch fortgesetzte Indizierung fehlgeschlagen: {e}")

    threading.Thread(target=_run, daemon=True, name="reindex-resume").start()
    return True


def _do_force_reindex(attempt: int = 1, resume_count: int = 0,
                      incremental: bool = False, resume_baseline: int = -1) -> dict:
    global _current_run
    started = time.time()
    # Bei einer Wiederaufnahme den urspruenglichen Start beibehalten, damit die
    # "Letzter Indexlauf"-Zeit nicht bei jedem Anlauf springt.
    if incremental:
        prev = get_last_run()
        started = prev.get("started_at") or started
    _current_run = {"started_at": started, "attempt": attempt,
                    "resumed": resume_count, "resume_baseline": resume_baseline,
                    "incremental": incremental}
    _set_progress(running=True, phase="Starte...", done=0, total=0, vector_done=0,
                  vector_total=0, vector_base=0, chunks=0, error="", current_file="",
                  started_at=started, finished_at=0.0, resumed=resume_count,
                  cancelled=False, attempt=attempt, max_attempts=MAX_INDEX_ATTEMPTS,
                  failed=0, failed_list=[])
    # Marker "laeuft" auf die Platte: stirbt der Prozess mittendrin (Neustart,
    # OOM-Killer), erkennt resume_interrupted_reindex() das beim naechsten Start
    # und setzt den Lauf fort. Ein sauberes Ende ueberschreibt den Marker.
    _save_last_run({**_current_run, "finished_at": 0, "status": "running",
                    "indexed_files": resume_baseline if resume_baseline > 0 else 0,
                    "total_chunks": 0})
    # Ausnahmen werden bewusst NICHT hier abgefangen: Endzustand und Protokoll
    # schreibt _run_with_retries – erst wenn alle Versuche verbraucht sind.
    # Sonst zeigte die Oberflaeche zwischen zwei Neuversuchen "Fehler/fertig".
    folders = _get_folders()
    max_bytes = _get_max_bytes()
    vs = _get_vector_store()

    # Die Dateiliste kann sich geaendert haben – der Neuaufbau muss den echten
    # Stand sehen, nicht den kurz gecachten.
    invalidate_files_cache()

    if vs is not None:
        # ── Nur FAISS aufbauen ──────────────────────────────────────────────
        # incremental: bestehenden Index behalten (Wiederaufnahme nach Absturz);
        # der Reindex ueberspringt unveraenderte Dateien automatisch.
        if not incremental:
            # NICHT blind leeren. `vs.clear()` warf frueher ALLES weg – war
            # danach ein Ordner nicht erreichbar, blieb sein Wissen dauerhaft
            # verloren (dieselbe Verwechslung wie im Suchpfad: "nicht
            # erreichbar" ist nicht "geloescht", nur hier mit dem groesseren
            # Hebel, weil der Neuaufbau ALLES anfasst).
            alive = [f for f in folders if _safe_exists(f)]
            if folders and not alive:
                # Gar nichts erreichbar: abbrechen statt leeren. Der Aufrufer
                # sieht einen Fehler, der bestehende Index bleibt unberuehrt.
                raise RuntimeError(
                    "Kein Wissensordner erreichbar – Neuaufbau abgebrochen. "
                    "Der bestehende Index bleibt erhalten. Betroffen: "
                    + ", ".join(str(f) for f in folders))
            if len(alive) == len(folders):
                vs.clear()                      # alles erreichbar -> wie bisher
            else:
                # Teilweise erreichbar: nur die erreichbaren Ordner neu
                # aufbauen, der Rest behaelt seinen Stand.
                prefixes = tuple(str(r).rstrip(os.sep) + os.sep for r in alive)
                stale = [p for p in vs.get_indexed_files() if p.startswith(prefixes)]
                if stale:
                    vs.remove_files(stale)
                _log.warning("Teil-Neuaufbau: %d von %d Ordnern erreichbar, "
                             "der Rest behaelt seinen Indexstand",
                             len(alive), len(folders))
        _rebuild_vector_index(folders, max_bytes, force=True)
        chunk_count = vs.chunk_count()
        file_count  = vs.file_count()
        result = {"indexed_files": file_count, "total_chunks": chunk_count,
                  "vector_info": f"Vektor: {chunk_count} Chunks"}
    else:
        # ── Nur TF-IDF aufbauen (FAISS nicht installiert) ───────────────────
        with _cache_lock:
            INDEX_CACHE_PATH.unlink(missing_ok=True)
        cache = _rebuild_cache(folders, max_bytes, force=True)
        total_chunks = sum(len(d.get("chunks", [])) for d in cache["files"].values())
        result = {"indexed_files": len(cache["files"]), "total_chunks": total_chunks,
                  "vector_info": ""}

    cancelled = _reindex_cancel.is_set()
    _fortschritt = get_index_progress()
    failed = int(_fortschritt.get("failed") or 0)
    failed_list = list(_fortschritt.get("failed_list") or [])

    # Gruppen-Pflege NACH dem Neuaufbau: erst jetzt entspricht der Index dem
    # tatsaechlichen Bestand. Zaehl-Basis ist Index + Platte (nicht nur die
    # Platte) – ein voruebergehend nicht erreichbares Netzlaufwerk darf die
    # Gruppen-Zuordnungen seiner Dateien nicht mitreissen.
    try:
        from backend import knowledge_groups as _kg
        basis = known_paths_with_disk()
        pruned = _kg.prune(basis)
        # Systemgenerierte Dateien der Gruppe "Erlernt" zuordnen. Das lief
        # bisher NUR beim Oeffnen der Gruppenseite – bis dahin galten gelernte
        # Dateien als "ungruppiert", der Gruppenfilter lieferte je nach
        # Vorgeschichte andere Ergebnisse.
        assigned = _kg.auto_assign_system_files(basis)
        if pruned or assigned:
            _log.info(f"Wissensgruppen gepflegt: {pruned} verwaiste Zuordnung(en) "
                      f"entfernt, {assigned} Datei(en) automatisch zugeordnet")
    except Exception as e:
        _log.warning(f"Gruppen-Pflege nach Reindex fehlgeschlagen: {e}")

    finished = time.time()
    _set_progress(running=False, phase="Abgebrochen" if cancelled else "Fertig",
                  finished_at=finished, cancelled=cancelled)
    _save_last_run({"started_at": started, "finished_at": finished,
                    "status": "cancelled" if cancelled else "ok",
                    "attempt": attempt, "resumed": resume_count,
                    "indexed_files": result["indexed_files"],
                    "total_chunks": result["total_chunks"],
                    "failed_files": failed,
                    "failed_list": failed_list})
    result["cancelled"] = cancelled
    result["failed_files"] = failed
    # Die NAMEN mitgeben, nicht nur die Zahl: die Oberflaeche soll den Grund
    # dort zeigen, wo die Meldung steht.
    result["failed_list"] = failed_list
    return result


class KnowledgeTool(BaseTool):
    """Durchsucht die lokale Knowledge Base (RAG)."""

    @property
    def name(self) -> str:
        return "knowledge_search"

    @property
    def description(self) -> str:
        return (
            "IMMER ZUERST AUFRUFEN bei Fragen zu Produkten, Software, Technik oder Kunden! "
            "Durchsucht die lokale Knowledge Base mit Kundendokumentation, Produkthandbüchern, "
            "Installationsanleitungen, technischen Spezifikationen und internen Vorgaben. "
            "Enthält PDFs, DOCX, PPTX, Excel und Textdateien. "
            "VOR jeder Web- oder Google-Suche dieses Tool verwenden – "
            "die Wissensdatenbank hat aktuelle, kundenbezogene Informationen die im Internet nicht zu finden sind."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Suchbegriff(e) zum Durchsuchen der Knowledge Base."
                },
                "max_results": {
                    "type": "integer",
                    "description": "Maximale Anzahl der Ergebnisse (Standard: 8)."
                }
            },
            "required": ["query"]
        }

    async def execute(self, **kwargs) -> str:
        query = kwargs.get("query", "")
        max_results = int(kwargs.get("max_results", 8))

        # Vom Benutzer gewaehlter Wissensgruppen-Filter (Modell B):
        #   None       -> kein Filter (alle Gruppen)
        #   []          -> Benutzer hat ALLE Gruppen abgewaehlt -> kein Wissen
        #   [ids...]    -> nur Treffer aus diesen Gruppen (ODER; "ungrouped" moeglich)
        kb_groups = kwargs.get("_kb_groups")
        if isinstance(kb_groups, list) and len(kb_groups) == 0:
            return ("🔒 Keine Wissensgruppen ausgewählt – es wird kein Wissen aus der "
                    "Knowledge Base verwendet. (Auswahl über den Wissensgruppen-Filter änderbar.)")
        # Ueber-Abfrage NUR noch als Rueckfall. Bis 2026-08-02 galt sie immer:
        # das Fuenffache holen und danach nach Gruppen filtern. Das verliert
        # still Treffer – wer nur eine kleine Gruppe freigegeben hat, dessen
        # bester Treffer kann jenseits der global besten 5·k liegen, und dann
        # kam "keine Treffer" heraus, obwohl passendes Wissen in seiner Gruppe
        # lag. Die Vektorsuche filtert jetzt WAEHREND der Bewertung
        # (FAISS-IDSelector + BM25-Skip). Der TF-IDF-Rueckfall kann das nicht,
        # dort bleibt es bei Ueber-Abfrage + Nachfilter.
        fetch_n = max(max_results * 5, 40) if kb_groups else max_results

        if not query.strip():
            return "❌ Fehler: query-Parameter fehlt. Bitte knowledge_search erneut aufrufen und einen konkreten Suchbegriff aus der Benutzeranfrage als 'query' übergeben (z.B. knowledge_search({'query': 'LDT Import Medistar'}))."

        # Standardordner sicherstellen
        (PROJECT_ROOT / DEFAULT_FOLDER).mkdir(parents=True, exist_ok=True)

        folders = _get_folders()
        max_bytes = _get_max_bytes()

        # Wissenssuche laeuft ueber die Vektor-/Datenbank-Suche. "auto" heisst:
        # ohne FAISS faellt sie auf TF-IDF zurueck, statt gar nichts zu liefern.
        # (Bis 2026-07-30 stand hier fest "vector" – ohne FAISS gab es dann
        # ueberhaupt keine Wissenssuche mehr, weil der Rueckfall-Zweig
        # unerreichbar war.)
        search_mode_cfg = "auto"

        # TF-IDF Cache nur laden wenn benoetigt (nicht bei reinem Vektor-Modus)
        # Spart bei 600+ Dateien das Laden aller Chunks in den RAM
        vs = _get_vector_store()
        vector_index_ready = vs is not None and vs.chunk_count() > 0
        need_tfidf_cache = search_mode_cfg == "tfidf" or (
            search_mode_cfg == "auto" and not vector_index_ready
        )

        # cache=None heisst "nicht geladen". Frueher wurde die komplette
        # knowledge_index.json (enthaelt ALLE Chunk-Texte) bei JEDER Suche
        # geparst, nur um zu pruefen, ob sie leer ist.
        cache = None
        if need_tfidf_cache:
            cache = await asyncio.to_thread(_rebuild_cache, folders, max_bytes, False)

        if not vector_index_ready and not (cache and cache["files"]):
            # Laeuft gerade ein Neuaufbau, ist der Index nur voruebergehend leer –
            # "keine Treffer" waere die falsche Auskunft.
            prog = get_index_progress()
            if prog.get("running"):
                done = prog.get("vector_done") or prog.get("done") or 0
                total = prog.get("vector_total") or prog.get("total") or 0
                return (f"⏳ Der Wissensindex wird gerade neu aufgebaut "
                        f"({done}/{total} Dateien). Deshalb sind aktuell keine "
                        f"Treffer möglich – bitte in einigen Minuten erneut fragen. "
                        f"Es fehlt KEIN Wissen, es ist nur vorübergehend nicht durchsuchbar.")
            files_on_disk = _all_files_cached(folders)
            if files_on_disk:
                return f"⚠️ Knowledge Base hat {len(files_on_disk)} Dateien, aber noch keinen Index. Bitte einmal 'Neu Indizieren' in den Einstellungen ausführen."
            folder_display = ", ".join(
                str(f.relative_to(PROJECT_ROOT)) if str(f).startswith(str(PROJECT_ROOT)) else str(f)
                for f in folders
            )
            return f"📂 Knowledge Base ist leer. Lege Dateien in einen der Ordner ab: {folder_display}"

        results = None
        search_mode = "TF-IDF"
        # Der TF-IDF-Fallback kennt keinen lexikalischen Anker – dort bleibt es
        # bei False (keine Warnung ohne Befund).
        kein_anker = False
        # True, sobald der Gruppenfilter IN der Suche gewirkt hat – dann darf
        # unten nicht noch einmal nachgefiltert werden.
        gruppen_in_suche = False

        def _erlaubte_pfade():
            """Erlaubte Pfade aus den INDIZIERTEN Dateien (~900), nicht aus den
            Treffern. Rueckgabe None = nicht ermittelbar; dann gilt der alte Weg
            (Ueber-Abfrage + Nachfilter). Ein Fehler hier darf den Filter nicht
            stillschweigend aufheben."""
            if not kb_groups or vs is None:
                return None
            try:
                from backend import knowledge_groups as kg
                return set(kg.filter_paths_by_groups(
                    list(vs.get_indexed_files().keys()), kb_groups))
            except Exception as e:  # noqa: BLE001
                print(f"[knowledge] Gruppenfilter nicht vorbereitbar, nutze Nachfilter: {e}",
                      flush=True)
                return None

        allow_paths = None
        if search_mode_cfg in ("auto", "vector"):
            has_vector = await asyncio.to_thread(_rebuild_vector_index, folders, max_bytes)
            if has_vector:
                # Vektor-Index vorhanden → ausschliesslich Vektor verwenden (kein TF-IDF-Fallback)
                # Begruendung: TF-IDF skaliert O(n) mit Dateizahl, Vektor konstant ~35ms
                # ERST HIER die erlaubten Pfade bestimmen: _rebuild_vector_index()
                # kann gerade geaenderte Dateien nachgetragen haben. Frueher
                # ermittelt, fehlten die in der Liste und waeren fuer diesen Lauf
                # unsichtbar gewesen.
                allow_paths = await asyncio.to_thread(_erlaubte_pfade)
                n = max_results if allow_paths is not None else fetch_n
                results = await asyncio.to_thread(_vector_search, query, n, allow_paths)
                gruppen_in_suche = allow_paths is not None
                search_mode = "Hybrid: Vektor+BM25"
                # SOFORT sichern. Das Merkmal haengt am Listen-Objekt, und
                # weiter unten machen GLEICH ZWEI Stellen eine gewoehnliche
                # Liste daraus: der Gruppenfilter (List Comprehension) und
                # `results[:max_results]` (Slicing einer list-Unterklasse gibt
                # list zurueck). Ohne diese Zeile verschwand die Warnung
                # lautlos – am echten Index nachgewiesen.
                kein_anker = getattr(results, "kein_anker", False)
            elif search_mode_cfg == "auto":
                # Kein Vektor-Index → TF-IDF als Fallback. ``cache`` kann None
                # sein, wenn der Index zwischen Pruefung und Suche verschwand.
                if cache is None:
                    cache = await asyncio.to_thread(_rebuild_cache, folders, max_bytes, False)
                results = _search(query, cache, fetch_n)
                search_mode = "TF-IDF"

        elif search_mode_cfg == "tfidf":
            if cache is None:
                cache = await asyncio.to_thread(_rebuild_cache, folders, max_bytes, False)
            results = _search(query, cache, fetch_n)
            search_mode = "TF-IDF"

        # Nachfilter nur noch, wenn die Suche selbst NICHT einschraenken konnte
        # (TF-IDF-Rueckfall oder fehlgeschlagene Vorbereitung oben).
        if kb_groups and results and not gruppen_in_suche:
            try:
                from backend import knowledge_groups as kg
                kept = set(kg.filter_paths_by_groups([r[1] for r in results], kb_groups))
                results = [r for r in results if r[1] in kept]
            except Exception:
                pass
        if results:
            results = results[:max_results]

        if not results:
            # Bestandszahlen aus der Quelle nehmen, die tatsaechlich gesucht hat.
            if vector_index_ready and vs is not None:
                n_files, n_chunks = vs.file_count(), vs.chunk_count()
            elif cache:
                n_files = len(cache["files"])
                n_chunks = sum(len(d.get("chunks", [])) for d in cache["files"].values())
            else:
                n_files = n_chunks = 0
            _grp = " in den gewählten Wissensgruppen" if kb_groups else ""
            return f"🔍 Keine Treffer für '{query}'{_grp} ({n_files} Dateien, {n_chunks} Chunks)."

        output = f"🔍 {len(results)} Treffer für '{query}' ({search_mode}):\n\n"

        # Kein Wort der Anfrage kommt im Bestand vor -> die Treffer beruhen
        # ALLEIN auf Vektor-Aehnlichkeit. Das ist der Fall, in dem die Suche
        # bisher still drei beliebige Chunks lieferte und das Modell darauf
        # eine Antwort baute. Der Hinweis geht an das MODELL, nicht an den
        # Benutzer – deshalb steht er im Werkzeug-Ergebnis und ist als Auftrag
        # formuliert, nicht als Beobachtung.
        #
        # Bewusst KEINE Unterdrueckung der Treffer: nachgemessen am Echt-Index
        # hat das Signal einen Fehlalarm – eine stark vertippte, aber voellig
        # legitime Anfrage ("patiententen anlgen") hat ebenfalls keinen Anker.
        # Wer dort nichts zurueckgibt, bestraft einen Tippfehler mit einer
        # leeren Antwort. Kennzeichnen kostet nichts und wirkt in beiden Faellen.
        if kein_anker:
            output = (
                f"⚠ HINWEIS ZUR QUALITÄT: Kein einziges Wort aus '{query}' kommt im "
                f"Wissensbestand vor. Die folgenden Treffer beruhen ausschliesslich auf "
                f"Vektor-Ähnlichkeit und sind mit hoher Wahrscheinlichkeit NICHT "
                f"einschlägig. Antworte NUR dann inhaltlich, wenn der Text die Frage "
                f"tatsächlich beantwortet – andernfalls sage klar, dass die "
                f"Wissensdatenbank dazu nichts enthält. Erfinde nichts hinzu.\n\n"
            ) + output

        for i, (score, filename, chunk) in enumerate(results, 1):
            output += f"--- [{i}] {filename} (Relevanz: {score:.2f}) ---\n"
            output += chunk.strip()[:CHUNK_OUTPUT_LIMIT] + "\n\n"

        return output


class KnowledgeManageTool(BaseTool):
    """Verwaltet Knowledge-Base-Ordner und den Suchindex."""

    @property
    def name(self) -> str:
        return "knowledge_manage"

    @property
    def description(self) -> str:
        return (
            "Verwaltet die Knowledge Base. "
            "Aktionen: list_folders (Ordner anzeigen), add_folder (Ordner hinzufügen), "
            "remove_folder (Ordner entfernen), reindex (Index neu aufbauen), "
            "list_docs (alle Dokumente auflisten), stats (Statistiken anzeigen)."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["list_folders", "add_folder", "remove_folder",
                             "reindex", "list_docs", "stats"],
                    "description": "Auszuführende Aktion."
                },
                "folder": {
                    "type": "string",
                    "description": "Ordnerpfad für add_folder/remove_folder."
                }
            },
            "required": ["action"]
        }

    async def execute(self, **kwargs) -> str:
        action = kwargs.get("action", "")
        folder_arg = kwargs.get("folder", "").strip()

        if action == "list_folders":
            folders = _get_folders()
            lines = []
            for f in folders:
                try:
                    rel = str(f.relative_to(PROJECT_ROOT))
                except ValueError:
                    rel = str(f)
                lines.append(f"  {'✅' if _safe_exists(f) else '❌'} {rel}")
            return "📁 Knowledge-Ordner:\n" + "\n".join(lines)

        elif action == "add_folder":
            if not folder_arg:
                return "❌ Kein Ordner angegeben."
            states = config.get_skill_states()
            state = states.get("knowledge", {})
            cfg = state.get("config", {})
            folders = [f.strip() for f in cfg.get("folders", DEFAULT_FOLDER).split(",") if f.strip()]
            if folder_arg in folders:
                return f"ℹ️ '{folder_arg}' ist bereits konfiguriert."
            folders.append(folder_arg)
            cfg["folders"] = ",".join(folders)
            state["config"] = cfg
            config.save_skill_state("knowledge", state)
            return f"✅ Ordner '{folder_arg}' hinzugefügt."

        elif action == "remove_folder":
            if not folder_arg:
                return "❌ Kein Ordner angegeben."
            states = config.get_skill_states()
            state = states.get("knowledge", {})
            cfg = state.get("config", {})
            folders = [f.strip() for f in cfg.get("folders", DEFAULT_FOLDER).split(",") if f.strip()]
            if folder_arg not in folders:
                return f"ℹ️ '{folder_arg}' nicht in der Liste."
            folders.remove(folder_arg)
            cfg["folders"] = ",".join(folders) if folders else DEFAULT_FOLDER
            state["config"] = cfg
            config.save_skill_state("knowledge", state)
            return f"✅ Ordner '{folder_arg}' entfernt."

        elif action == "reindex":
            result = await asyncio.to_thread(force_reindex)
            if result.get("skipped"):
                return "ℹ️ Es läuft bereits eine Indizierung – ein Nachlauf ist vorgemerkt."
            _fail = result.get("failed_files") or 0
            _fail_txt = f", ⚠️ {_fail} Datei(en) fehlgeschlagen" if _fail else ""
            return (f"✅ Index neu aufgebaut: {result['indexed_files']} Dateien, "
                    f"{result['total_chunks']} Chunks{result.get('vector_info', '')}{_fail_txt}.")

        elif action == "list_docs":
            folders = _get_folders()
            files = _all_files(folders)
            if not files:
                return "📂 Keine Dokumente gefunden."
            lines = []
            for f in sorted(files):
                size = f.stat().st_size
                size_str = f"{size/1024:.1f} KB" if size >= 1024 else f"{size} B"
                try:
                    rel = str(f.relative_to(PROJECT_ROOT))
                except ValueError:
                    rel = str(f)
                lines.append(f"  📄 {rel} ({size_str})")
            return f"📚 {len(files)} Dokument(e):\n" + "\n".join(lines)

        elif action == "stats":
            stats = get_stats()
            formats = ["Text/Markdown"]
            if stats["pdf_support"]:
                formats.append("PDF")
            else:
                formats.append("PDF ⚠️ (pdfplumber fehlt)")
            if stats["docx_support"]:
                formats.append("DOCX")
            else:
                formats.append("DOCX ⚠️ (python-docx fehlt)")
            if stats["xlsx_support"]:
                formats.append("Excel")
            else:
                formats.append("Excel ⚠️ (openpyxl fehlt)")
            if stats["pptx_support"]:
                formats.append("PowerPoint")
            else:
                formats.append("PowerPoint ⚠️ (python-pptx fehlt)")
            if stats["video_support"]:
                formats.append("Video/Audio")
            else:
                formats.append("Video/Audio ⚠️ (ffmpeg + faster-whisper nötig)")
            # Bild-OCR (Tesseract)
            try:
                import pytesseract as _pt  # noqa: F401
                import shutil as _sh
                _ocr_ok = bool(_sh.which("tesseract"))
            except Exception:
                _ocr_ok = False
            formats.append("Bilder/OCR" if _ocr_ok
                           else "Bilder/OCR ⚠️ (tesseract-ocr + pytesseract nötig)")
            if stats.get("onenote_support"):
                formats.append("OneNote")
            else:
                formats.append("OneNote ⚠️ (Java + Apache Tika nötig, "
                               "deploy/tika_setup.sh)")
            size_mb = stats["total_size_bytes"] / (1024 * 1024)

            # Vektor-Info
            if stats.get("vector_search"):
                vector_line = f"\n  🧠 Vektor-Suche: aktiv ({stats['vector_files']} Dateien, {stats['vector_chunks']} Chunks)"
            else:
                vector_line = "\n  🧠 Vektor-Suche: inaktiv (faiss-cpu/sentence-transformers fehlt)"

            return (
                f"📊 Knowledge Base Statistiken:\n"
                f"  Dateien: {stats['total_files']} ({size_mb:.1f} MB)\n"
                f"  TF-IDF Index: {stats['indexed_files']} Dateien, {stats['total_chunks']} Chunks"
                f"{vector_line}\n"
                f"  Formate: {', '.join(formats)}"
            )

        return f"❌ Unbekannte Aktion: {action}"
