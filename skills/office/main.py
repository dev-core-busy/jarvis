"""Office-Skill – erzeugt und liest Office-Dokumente (Word/Excel/PowerPoint)
und exportiert sie nach PDF.

Ansatz: programmatisch via python-docx / openpyxl / python-pptx (deterministisch,
headless). PDF-Export via LibreOffice (soffice --headless --convert-to pdf).

Dateien landen im Server-Dateisystem unter data/documents/ mit Capability-Name
(<32-Hex>__<Name>.<ext>) und werden via /api/documents/{name} zum Download
ausgeliefert (siehe backend/main.py).
"""

import os
import re
import json
import uuid
import asyncio
import subprocess
from pathlib import Path

from backend.tools.base import BaseTool

# data/documents/ relativ zum Projekt-Root (skills/office/ -> ../../)
DOCS_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "documents"

_UML = str.maketrans({
    "ä": "ae", "ö": "oe", "ü": "ue", "Ä": "Ae", "Ö": "Oe", "Ü": "Ue", "ß": "ss",
})


def _safe_base(name: str, default: str = "dokument") -> str:
    """Macht aus einem (ggf. unsicheren) Namen einen ASCII-sicheren Basisnamen."""
    base = os.path.splitext(os.path.basename(name or ""))[0].translate(_UML)
    base = re.sub(r"[^A-Za-z0-9_\- ]+", "", base).strip().replace(" ", "_")
    return base or default


def _new_path(friendly: str, ext: str):
    """Erzeugt einen neuen Capability-Pfad. Gibt (disk_path, fname, download_name) zurueck."""
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    token = uuid.uuid4().hex  # 32 Hex-Zeichen
    base = _safe_base(friendly)
    fname = f"{token}__{base}.{ext}"
    return DOCS_DIR / fname, fname, f"{base}.{ext}"


def _ok(download_name: str, fname: str, disk_path: Path, extra: str = "") -> str:
    # Eigentuemer beim ERZEUGEN vermerken, nicht erst beim Ausliefern.
    # agent.py::_deliver_docs registriert die Datei ebenfalls, laeuft aber bei
    # Sub-Agenten gar nicht – dort waere die eigene Datei sonst sofort wieder
    # unsichtbar (Eigentuemer-Schranke ist fail-closed) und ein folgendes
    # office_to_pdf/office_read wuerde "Datei nicht gefunden" melden.
    try:
        from backend import documents as _documents, sandbox as _sbx
        _u = _sbx.tool_user()
        if _u:
            _documents.register(fname, _u)
    except Exception:
        pass
    # Markdown-Download-Link, den die Frontends als Download-Chip rendern.
    return (
        f"✅ '{download_name}' wurde erstellt.\n\n"
        f"[📥 {download_name} herunterladen](/api/documents/{fname})"
        + (f"\n\n{extra}" if extra else "")
    )


def _sichtbar(p: Path | None) -> Path | None:
    """Eigentuemer-Schranke fuer data/documents (dieselbe wie am HTTP-Endpunkt).

    Fremde Dateien werden behandelt, als gaebe es sie nicht (None) – die
    Aufrufer melden dann "Datei nicht gefunden". Dass eine Datei existiert, ist
    selbst eine Information; deshalb keine eigene Fehlermeldung.
    """
    if p is None:
        return None
    try:
        from backend import sandbox as _sbx
        rp = p.resolve()
        if rp.parent == _sbx.DOCS_ROOT and not _sbx.may_see_document(rp.name):
            return None
    except Exception:
        pass
    return p


def _resolve_existing(path: str) -> Path | None:
    """Loest einen Eingabepfad zu einer existierenden Datei auf.

    Akzeptiert: reinen Dateinamen in data/documents/, '/api/documents/<name>'
    oder einen beliebigen (absoluten/relativen) Dateisystempfad. Dateien in
    data/documents, die einem ANDEREN Benutzer gehoeren, gelten als nicht
    vorhanden (siehe ``_sichtbar``).
    """
    if not path:
        return None
    path = path.strip()
    if path.startswith("/api/documents/"):
        path = path[len("/api/documents/"):]
    cand = DOCS_DIR / path
    if cand.exists():
        return _sichtbar(cand)
    # Privates /tmp pro Agent-Lauf: was die Shell als /tmp/ergebnis.docx
    # geschrieben hat, liegt auf dem Host im Lauf-Verzeichnis. Ohne diese
    # Uebersetzung meldet office_read/office_to_pdf "Datei nicht gefunden" fuer
    # eine Datei, die der Agent gerade selbst erzeugt hat – und der Widerspruch
    # ist von aussen nicht erklaerbar. Nicht-/tmp-Pfade und blosse Namen bleiben
    # unberuehrt (siehe backend/lauf_tmp.aufloesen).
    try:
        from backend import lauf_tmp as _lt
        path = _lt.aufloesen(path)
    except Exception:  # noqa: BLE001
        pass
    p = Path(path)
    if p.exists():
        return _sichtbar(p)
    # ANZEIGENAME: Auf Platte heisst die Datei '<32-Hex>__<Anzeigename>', der
    # Werkzeug-Erfolgstext nennt aber nur den Anzeigenamen. Das LLM gibt genau
    # den weiter ("office_to_pdf: IT-Projektangebot.docx") und lief bis
    # 2026-07-28 in "Datei nicht gefunden" – die Kette Erstellen→PDF war damit
    # ueber den natuerlichen Weg nicht benutzbar. Daher hier nachschlagen;
    # bei mehreren Treffern der jüngste (= der gerade erzeugte).
    name = Path(path).name
    if name and "/" not in path.strip("/"):
        try:
            treffer = [f for f in DOCS_DIR.glob(f"*__{name}")
                       if f.is_file() and _sichtbar(f) is not None]
            if treffer:
                return max(treffer, key=lambda f: f.stat().st_mtime)
        except Exception:
            pass
    return None


# ─────────────────────────────────────────────────────────────────────────
# Word
# ─────────────────────────────────────────────────────────────────────────
class CreateWordTool(BaseTool):
    @property
    def name(self) -> str:
        return "office_create_word"

    @property
    def description(self) -> str:
        return (
            "Erstellt ein Word-Dokument (.docx). 'content' wird zeilenweise interpretiert: "
            "'# ' = Ueberschrift 1, '## ' = Ueberschrift 2, '### ' = Ueberschrift 3, "
            "'- ' oder '* ' = Aufzaehlung, '1. ' = nummerierte Liste, Leerzeile = neuer Absatz, "
            "sonst normaler Absatz. Gibt eine Download-URL zurueck."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "filename": {"type": "STRING", "description": "Dateiname/Titel des Dokuments (ohne Pfad), z.B. 'Quartalsbericht'."},
                "title": {"type": "STRING", "description": "Optionaler Titel, der als grosse Ueberschrift oben eingefuegt wird."},
                "content": {"type": "STRING", "description": "Inhalt des Dokuments (mit einfacher Markdown-Syntax, siehe Beschreibung)."},
            },
            "required": ["filename", "content"],
        }

    async def execute(self, filename: str = "", title: str = "", content: str = "", **kwargs) -> str:
        if not filename:
            return "Fehler: 'filename' ist Pflicht."
        try:
            from docx import Document
        except Exception as e:
            return f"Fehler: python-docx nicht verfuegbar ({e})."

        doc = Document()
        if title:
            doc.add_heading(title, level=0)

        for raw in (content or "").split("\n"):
            line = raw.rstrip()
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:].strip(), level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:].strip(), level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:].strip(), level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
            elif re.match(r"^\d+\.\s", stripped):
                doc.add_paragraph(re.sub(r"^\d+\.\s", "", stripped), style="List Number")
            else:
                doc.add_paragraph(stripped)

        disk, fname, dl = _new_path(filename, "docx")
        try:
            doc.save(str(disk))
        except Exception as e:
            return f"Fehler beim Speichern: {e}"
        return _ok(dl, fname, disk)


# ─────────────────────────────────────────────────────────────────────────
# Excel
# ─────────────────────────────────────────────────────────────────────────
class CreateExcelTool(BaseTool):
    @property
    def name(self) -> str:
        return "office_create_excel"

    @property
    def description(self) -> str:
        return (
            "Erstellt eine Excel-Tabelle (.xlsx). Entweder 'rows' (2D-Liste von Zellen) "
            "mit optionalen 'headers' und 'sheet_name' fuer EIN Blatt, ODER 'sheets' "
            "(Objekt: Blattname -> {headers:[...], rows:[[...]]}) fuer mehrere Blaetter. "
            "Gibt eine Download-URL zurueck."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "filename": {"type": "STRING", "description": "Dateiname (ohne Pfad), z.B. 'Umsatz'."},
                "sheet_name": {"type": "STRING", "description": "Blattname fuer den Einzelblatt-Modus (Standard 'Tabelle1')."},
                "headers": {"type": "ARRAY", "items": {"type": "STRING"}, "description": "Optionale Kopfzeile (Einzelblatt-Modus)."},
                "rows": {"type": "ARRAY", "items": {"type": "ARRAY", "items": {"type": "STRING"}}, "description": "Datenzeilen als 2D-Liste (Einzelblatt-Modus)."},
                "sheets": {"type": "OBJECT", "description": "Mehrblatt-Modus: { 'Blattname': { 'headers': [...], 'rows': [[...]] } }."},
            },
            "required": ["filename"],
        }

    def _write_sheet(self, ws, headers, rows):
        if headers:
            ws.append(list(headers))
            # Kopfzeile fett
            from openpyxl.styles import Font
            for cell in ws[1]:
                cell.font = Font(bold=True)
        for row in (rows or []):
            ws.append(list(row) if isinstance(row, (list, tuple)) else [row])

    async def execute(self, filename: str = "", sheet_name: str = "", headers=None,
                       rows=None, sheets=None, **kwargs) -> str:
        if not filename:
            return "Fehler: 'filename' ist Pflicht."

        # ── LAUTER Fehlschlag statt leerer Datei (Vorfall ECHT 2026-08-19) ──
        # Das Modell rief dieses Werkzeug mit dem Parameter " sheets" auf – mit
        # FUEHRENDEM LEERZEICHEN. Der landete in **kwargs, wurde wortlos
        # verworfen, `sheets`/`rows`/`headers` blieben None, es wurde eine LEERE
        # Mappe gespeichert – und die Antwort lautete "✅ erstellt". Der Benutzer
        # bekam eine 0-Zeilen-Datei als Ergebnis angeboten, und das Modell hatte
        # keine Chance, den Fehler zu bemerken.
        unbekannt = [k for k in kwargs if not k.startswith("_")]
        if unbekannt:
            return ("Fehler: unbekannte Parameter " + ", ".join(repr(k) for k in unbekannt)
                    + ". Achte auf die genaue Schreibweise (auch auf fuehrende "
                      "Leerzeichen im Parameternamen). Erlaubt sind: filename, "
                      "sheet_name, headers, rows, sheets.")

        # `sheets` als JSON-STRING statt als Objekt ist der zweite Weg in
        # dieselbe leere Datei: `isinstance(sheets, dict)` war False, der Code
        # fiel in den Einzelblatt-Zweig, und dort war ebenfalls nichts zu
        # schreiben. Tolerant parsen statt still verwerfen.
        if isinstance(sheets, str) and sheets.strip():
            try:
                sheets = json.loads(sheets)
            except Exception:
                return ("Fehler: 'sheets' kam als Text an und ist kein gueltiges "
                        "JSON. Uebergib ein Objekt "
                        "{\"Blattname\": {\"headers\": [...], \"rows\": [[...]]}}.")
        if isinstance(rows, str) and rows.strip():
            try:
                rows = json.loads(rows)
            except Exception:
                return ("Fehler: 'rows' kam als Text an und ist kein gueltiges "
                        "JSON. Uebergib eine 2D-Liste [[...], [...]].")

        if sheets is not None and not isinstance(sheets, dict):
            return (f"Fehler: 'sheets' muss ein Objekt sein "
                    f"(Blattname -> {{headers, rows}}), war aber "
                    f"{type(sheets).__name__}.")
        if rows is not None and not isinstance(rows, (list, tuple)):
            return (f"Fehler: 'rows' muss eine 2D-Liste sein, war aber "
                    f"{type(rows).__name__}.")

        # Nichts zu schreiben = Fehler. Eine leere Tabelle ist niemals das,
        # was jemand bestellt hat.
        hat_inhalt = bool(
            (isinstance(sheets, dict) and any(
                (v or {}).get("headers") or (v or {}).get("rows")
                for v in sheets.values() if isinstance(v, dict)))
            or headers or rows
        )
        if not hat_inhalt:
            return ("Fehler: es wurden keine Daten uebergeben – es wurde KEINE "
                    "Datei erzeugt. Gib entweder 'rows' (mit optionalem "
                    "'headers') fuer ein Blatt an oder 'sheets' fuer mehrere.")

        try:
            from openpyxl import Workbook
        except Exception as e:
            return f"Fehler: openpyxl nicht verfuegbar ({e})."

        wb = Workbook()
        if sheets and isinstance(sheets, dict):
            first = True
            for sname, sdef in sheets.items():
                sdef = sdef or {}
                ws = wb.active if first else wb.create_sheet()
                ws.title = str(sname)[:31] or "Tabelle"
                self._write_sheet(ws, sdef.get("headers"), sdef.get("rows"))
                first = False
        else:
            ws = wb.active
            ws.title = (sheet_name or "Tabelle1")[:31]
            self._write_sheet(ws, headers, rows)

        disk, fname, dl = _new_path(filename, "xlsx")
        try:
            wb.save(str(disk))
        except Exception as e:
            return f"Fehler beim Speichern: {e}"
        return _ok(dl, fname, disk)


# ─────────────────────────────────────────────────────────────────────────
# PowerPoint
# ─────────────────────────────────────────────────────────────────────────
# Layout-Ansprache ueber NAMEN, nicht ueber den Index. Grund: `slide_layouts[1]`
# ist nur im Standardtemplate von python-pptx „Titel und Inhalt"; in einer
# echten Firmenvorlage zeigt derselbe Index irgendwohin (oft auf ein
# Bild-Layout). Die Alias-Listen enthalten deutsche UND englische Namen, damit
# sowohl die selbst erzeugte Vorlage als auch eine mitgebrachte .potx passt.
_LAYOUT_ALIAS = {
    "titel":     ["Titelfolie", "Title Slide", "Titeldia"],
    "inhalt":    ["Titel und Inhalt", "Title and Content", "Titel und Inhaltsverzeichnis"],
    "abschnitt": ["Abschnitt", "Section Header", "Abschnittsüberschrift"],
    "zwei":      ["Zwei Inhalte", "Two Content", "Vergleich", "Comparison"],
    "nurtitel":  ["Nur Titel", "Title Only"],
    "leer":      ["Leer", "Blank"],
    "bild":      ["Bild mit Beschriftung", "Picture with Caption"],
}
# Rueckfall-Indizes, falls die Vorlage keinen passenden Namen hat. Bewusst
# konservativ: 0 = erstes (fast immer die Titelfolie), 1 = zweites.
_LAYOUT_FALLBACK = {"titel": 0, "inhalt": 1, "abschnitt": 2, "zwei": 3,
                    "nurtitel": 5, "leer": 6, "bild": 8}


def _layout(prs, art: str):
    """Sucht ein Layout nach Namen (Alias-Liste), sonst per Rueckfall-Index.

    Der Vergleich ist gegen Gross-/Kleinschreibung und Zusaetze tolerant
    ('Titel und Inhalt (intern)' passt auf 'Titel und Inhalt'), weil
    Firmenvorlagen ihre Layouts gern durchnummerieren oder ergaenzen."""
    layouts = list(prs.slide_layouts)
    if not layouts:
        return None
    namen = [(l, (l.name or "").strip().lower()) for l in layouts]
    for wunsch in _LAYOUT_ALIAS.get(art, []):
        w = wunsch.lower()
        for l, n in namen:
            if n == w:
                return l
    for wunsch in _LAYOUT_ALIAS.get(art, []):
        w = wunsch.lower()
        for l, n in namen:
            if w in n:
                return l
    idx = _LAYOUT_FALLBACK.get(art, 1)
    return layouts[idx] if idx < len(layouts) else layouts[-1]


def _titel_setzen(slide, text: str) -> bool:
    """Schreibt den Folientitel in den Titel-Platzhalter (falls vorhanden)."""
    try:
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(text or "")
            return True
    except Exception:  # noqa: BLE001
        pass
    return False


def _text_platzhalter(slide):
    """Alle Platzhalter, die Inhalt aufnehmen koennen – ohne den Titel.

    Ausgewaehlt wird ueber den Platzhalter-TYP, nicht ueber `placeholders[1]`:
    der Index 1 ist im Standardtemplate der Textkoerper, in einer Firmenvorlage
    aber oft die Fusszeile oder eine Bildbox. Ein Datum-, Fusszeilen- oder
    Foliennummern-Platzhalter darf niemals mit Inhalt befuellt werden."""
    from pptx.enum.shapes import PP_PLACEHOLDER as PH
    erlaubt = {PH.BODY, PH.OBJECT, PH.SUBTITLE, PH.VERTICAL_BODY,
               PH.VERTICAL_OBJECT, PH.TABLE, PH.CHART}
    raus = []
    for ph in slide.placeholders:
        try:
            typ = ph.placeholder_format.type
        except Exception:  # noqa: BLE001
            continue
        if typ in (PH.TITLE, PH.CENTER_TITLE, PH.VERTICAL_TITLE):
            continue
        if typ in erlaubt and ph.has_text_frame:
            raus.append(ph)
    return raus


def _fuelle(platzhalter, bullets=None, text: str = "") -> bool:
    """Befuellt einen Platzhalter mit Aufzaehlung oder Freitext.

    Es wird NUR Text gesetzt – keine Schriftgroesse, keine Farbe, kein
    Zeilenabstand. Das ist der ganze Sinn des Vorlagen-Wegs: diese Werte kommen
    aus dem Layout, und wer sie hier ueberschreibt, erzeugt eine Datei, die
    nach dem Wechsel des Designs falsch aussieht.

    Ebenen: ein Eintrag 'a > b' oder ein Tupel ('b', 1) wird als Unterpunkt
    gesetzt (paragraph.level)."""
    tf = platzhalter.text_frame
    eintraege = []
    if bullets and isinstance(bullets, (list, tuple)):
        for b in bullets:
            if isinstance(b, dict):
                eintraege.append((str(b.get("text", "")), int(b.get("level", 0) or 0)))
            elif isinstance(b, (list, tuple)) and len(b) == 2:
                eintraege.append((str(b[0]), int(b[1] or 0)))
            else:
                s = str(b)
                # '> ' am Anfang = eine Ebene tiefer (bequeme Schreibweise fuer
                # das Modell, das sonst verschachtelte Objekte bauen muesste).
                tiefe = 0
                while s.startswith(">"):
                    tiefe += 1
                    s = s[1:].lstrip()
                eintraege.append((s, min(tiefe, 4)))
    elif text:
        eintraege = [(str(text), 0)]
    if not eintraege:
        return False
    tf.text = eintraege[0][0]
    if eintraege[0][1]:
        tf.paragraphs[0].level = eintraege[0][1]
    for inhalt, ebene in eintraege[1:]:
        p = tf.add_paragraph()
        p.text = inhalt
        if ebene:
            p.level = ebene
    return True


def _leere_platzhalter_entfernen(slide) -> None:
    """Entfernt Platzhalter, die leer geblieben sind.

    Ohne das zeigt PowerPoint auf der Folie den Hinweis „Klicken Sie, um Text
    hinzuzufuegen" und der PDF-Export einen leeren Rahmen – beides sieht nach
    unfertiger Arbeit aus. Titel-Platzhalter mit Text bleiben natuerlich."""
    for ph in list(slide.placeholders):
        try:
            if ph.has_text_frame and not ph.text_frame.text.strip():
                ph._element.getparent().remove(ph._element)
        except Exception:  # noqa: BLE001
            continue


# Obergrenze fuer die Folienzahl. KEINE willkuerliche Zahl: ein Foliensatz, den
# ein Mensch noch vortraegt, liegt darunter – und wer mehr braucht, will in
# Wahrheit eine Tabelle. Der Deckel ist die zweite Sicherung hinter
# `_slides_normalisieren`; er wird NIE still angewandt, sondern im Ergebnis
# genannt (sonst haelt der Aufrufer die gekuerzte Datei fuer vollstaendig).
MAX_FOLIEN = 60


def _slides_normalisieren(slides):
    """Bringt das 'slides'-Argument in eine Liste von Folien-Objekten.

    WARUM DAS NOETIG IST – der Vorfall vom 2026-08-10: ein Modell schickte den
    Foliensatz als TEXT (`"[{'title': …}, …]"`, Python-Schreibweise mit
    einfachen Anfuehrungszeichen, also nicht einmal gueltiges JSON). Die
    Schleife lief mit `for sl in slides` ueber die **Zeichen** dieses Strings,
    und weil ein String-Element toleranterweise als `{"title": …}` galt, entstand
    aus JEDEM Zeichen eine Folie: **2586 Folien**, 2 MB, unbrauchbar. Die
    Toleranz sass also auf der falschen Ebene – sie fing das Element ab, nicht
    den Container.

    Ueber einen String zu iterieren ist in Python immer erlaubt und nie
    gemeint. Deshalb wird hier zuerst der Container geprueft:
      * String  -> als JSON und als Python-Literal parsen (``ast.literal_eval``
        ist sicher, es fuehrt keinen Code aus). Scheitert beides, gibt es einen
        **Fehler mit Auszug** – niemals eine Folie je Zeichen.
      * dict    -> eine einzelne Folie.
      * Liste   -> Elemente einzeln normalisieren.

    Rueckgabe: ``(folien, fehler)``; ist ``fehler`` gesetzt, darf nichts
    erzeugt werden."""
    import ast
    import json

    def als_liste(wert, tiefe=0):
        if isinstance(wert, str):
            text = wert.strip()
            if not text:
                return None
            if tiefe > 1:
                return None          # kein endloses Aufdroeseln verschachtelter Texte
            for parser in (json.loads, ast.literal_eval):
                try:
                    return als_liste(parser(text), tiefe + 1)
                except Exception:  # noqa: BLE001
                    continue
            return None
        if isinstance(wert, dict):
            return [wert]
        if isinstance(wert, (list, tuple)):
            return list(wert)
        return None

    if slides is None:
        return [], "Fehler: 'slides' fehlt."

    roh = als_liste(slides)
    if roh is None:
        auszug = str(slides)[:120].replace("\n", " ")
        return [], (
            "Fehler: 'slides' konnte nicht gelesen werden. Erwartet wird eine LISTE "
            "von Folien-Objekten, z.B. [{\"title\": \"…\", \"bullets\": [\"…\"]}] – "
            "kein Text, der eine Liste beschreibt. "
            f"Empfangen wurde {type(slides).__name__}: {auszug}…"
        )

    folien = []
    verworfen = 0
    for element in roh:
        if isinstance(element, str):
            text = element.strip()
            # Ein Element kann seinerseits ein Objekt als Text sein.
            innen = als_liste(text)
            if innen and all(isinstance(e, dict) for e in innen):
                folien.extend(innen)
                continue
            if not text:
                continue
            # Kurz und einzeilig ist eine Ueberschrift, alles andere Fliesstext.
            folien.append({"title": text} if len(text) <= 80 and "\n" not in text
                          else {"content": text})
        elif isinstance(element, dict):
            folien.append(element)
        else:
            verworfen += 1

    if not folien:
        return [], "Fehler: 'slides' enthaelt keine verwertbare Folie."
    return folien, None


_SHELL_SUBST = re.compile(r"\$\(\s*date[^)]*\)")


def _text_bereinigen(text: str) -> str:
    """Ersetzt eine unaufgeloeste Shell-Substitution durch das heutige Datum.

    Im selben Vorfall stand auf der Titelfolie woertlich
    ``Stand: $(date +%d.%m.%Y)`` – das Modell hatte damit gerechnet, dass eine
    Shell den Ausdruck aufloest. Hier laeuft keine Shell. Ein sichtbarer
    Platzhalter auf der ersten Folie ist der peinlichste Fehler in einer
    Praesentation, und das Datum ist die einzige Absicht, die dahinterstecken
    kann – deshalb wird genau dieser Fall (und nur er) eingesetzt."""
    if not text or "$(" not in text:
        return text
    from datetime import date
    return _SHELL_SUBST.sub(date.today().strftime("%d.%m.%Y"), text)


# Diagrammtypen, deutsch UND englisch. Bewusst eine kleine Auswahl: es sind
# genau die vier, die auf einer Folie funktionieren. Ein Streudiagramm mit 500
# Punkten gehoert nicht in eine Praesentation.
_CHART_TYPEN = {
    "saeulen": "COLUMN_CLUSTERED", "säulen": "COLUMN_CLUSTERED",
    "spalten": "COLUMN_CLUSTERED", "column": "COLUMN_CLUSTERED",
    "bar": "COLUMN_CLUSTERED",          # gaengige Verwechslung: 'bar' meint meist Saeulen
    "balken": "BAR_CLUSTERED", "barh": "BAR_CLUSTERED",
    "linie": "LINE_MARKERS", "line": "LINE_MARKERS", "verlauf": "LINE_MARKERS",
    "kreis": "PIE", "pie": "PIE", "torte": "PIE", "doughnut": "DOUGHNUT",
    "ring": "DOUGHNUT",
}


def _zahl(wert):
    """Wandelt einen Wert in eine Zahl – auch '1.216' und '76 %'.

    Das Modell liefert Zahlen regelmaessig als formatierten Text. `float()`
    macht aus '1.216' die Zahl 1.216 statt 1216 – ein stiller Faktor 1000
    (dieselbe Falle wie in backend/tools/chart.py::parse_number)."""
    if isinstance(wert, (int, float)):
        return float(wert)
    s = str(wert or "").strip().replace("%", "").replace("€", "").replace(" ", " ")
    s = s.replace(" ", "")
    if not s:
        return None
    neg = s.startswith("(") and s.endswith(")")
    s = s.strip("()")
    # Bei gemischten Trennzeichen ist das RECHTESTE das Dezimaltrennzeichen.
    if "," in s and "." in s:
        s = (s.replace(".", "").replace(",", ".") if s.rfind(",") > s.rfind(".")
             else s.replace(",", ""))
    elif "," in s:
        s = s.replace(",", ".")
    elif s.count(".") == 1 and len(s.split(".")[1]) == 3:
        s = s.replace(".", "")          # '1.216' = Tausenderpunkt
    try:
        z = float(s)
    except ValueError:
        return None
    return -z if neg else z


def _diagramm_einfuegen(slide, prs, spec, neben_text: bool = False) -> bool:
    """Legt ein NATIVES PowerPoint-Diagramm auf die Folie.

    Nativ und nicht als Bild: das Diagramm bleibt in PowerPoint anklickbar und
    bearbeitbar, sein Datenblatt liegt in der Datei, und die Reihenfarben kommen
    aus dem THEME – also automatisch aus dem Hausdesign. Ein eingebettetes PNG
    waere beim Zoomen unscharf und truege seine Farben fest eingebrannt.

    ``spec``: {'typ': 'saeulen|balken|linie|kreis', 'kategorien': [...],
    'werte': [...]} – oder mehrere Reihen ueber 'reihen':
    [{'name': …, 'werte': [...]}, …]. Fehlt etwas Wesentliches, passiert
    NICHTS und die Folie bleibt wie sie ist: eine halbe Grafik ist schlechter
    als keine."""
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from skills.office import vorlage as _v
    except Exception:  # noqa: BLE001
        return False

    if not isinstance(spec, dict):
        return False
    kategorien = spec.get("kategorien") or spec.get("categories") or spec.get("labels")
    if not isinstance(kategorien, (list, tuple)) or not kategorien:
        return False
    kategorien = [str(k) for k in kategorien]

    reihen = []
    roh_reihen = spec.get("reihen") or spec.get("series")
    if isinstance(roh_reihen, (list, tuple)) and roh_reihen:
        for r in roh_reihen:
            if isinstance(r, dict):
                werte = r.get("werte") or r.get("values") or r.get("data")
                reihen.append((str(r.get("name") or r.get("titel") or "Reihe"), werte))
    else:
        werte = spec.get("werte") or spec.get("values") or spec.get("data")
        reihen.append((str(spec.get("name") or spec.get("titel") or spec.get("title") or "Wert"), werte))

    daten = CategoryChartData()
    daten.categories = kategorien
    gueltig = 0
    for name, werte in reihen:
        if not isinstance(werte, (list, tuple)):
            continue
        zahlen = [_zahl(w) for w in werte]
        # Auf die Kategorienzahl bringen: fehlende Werte als Luecke (None),
        # ueberzaehlige abschneiden. Sonst lehnt python-pptx die Reihe ab und
        # es entstuende gar kein Diagramm.
        zahlen = (zahlen + [None] * len(kategorien))[:len(kategorien)]
        if all(z is None for z in zahlen):
            continue
        daten.add_series(name, zahlen)
        gueltig += 1
    if not gueltig:
        return False

    typ_name = _CHART_TYPEN.get(str(spec.get("typ") or spec.get("type") or "").strip().lower(),
                                "COLUMN_CLUSTERED")
    try:
        typ = getattr(XL_CHART_TYPE, typ_name)
    except Exception:  # noqa: BLE001
        typ = XL_CHART_TYPE.COLUMN_CLUSTERED

    # Platz: der Inhaltsbereich des Satzspiegels. Steht daneben Text, nimmt das
    # Diagramm die rechte Haelfte – sonst die volle Breite.
    links, oben = _v.RAND_LINKS, _v.INHALT_Y
    breite, hoehe = _v.INHALT_B, _v.INHALT_H
    if neben_text:
        spalte = (_v.INHALT_B - _v.SPALTEN_LUFT) // 2
        links = _v.RAND_LINKS + spalte + _v.SPALTEN_LUFT
        breite = spalte
    try:
        rahmen = slide.shapes.add_chart(typ, links, oben, breite, hoehe, daten)
    except Exception:  # noqa: BLE001
        return False

    try:
        diagramm = rahmen.chart
        # Legende nur, wenn es mehr als eine Reihe gibt – bei einer einzigen
        # wiederholt sie nur den Folientitel und kostet Flaeche.
        diagramm.has_legend = gueltig > 1 or typ_name in ("PIE", "DOUGHNUT")
        if diagramm.has_legend:
            diagramm.legend.position = XL_LEGEND_POSITION.BOTTOM
            diagramm.legend.include_in_layout = False
        titel = spec.get("titel") or spec.get("title")
        diagramm.has_title = bool(titel)
        if titel:
            diagramm.chart_title.text_frame.text = str(titel)
        # Werte anschreiben, solange es lesbar bleibt. Bei einem Kreisdiagramm
        # ist der Prozentwert die eigentliche Aussage.
        punkte = len(kategorien) * max(gueltig, 1)
        if typ_name in ("PIE", "DOUGHNUT") and len(kategorien) <= 8:
            plot = diagramm.plots[0]
            plot.has_data_labels = True
            plot.data_labels.show_percentage = True
            plot.data_labels.show_value = False
        elif punkte <= 24:
            plot = diagramm.plots[0]
            plot.has_data_labels = True
            # OOXML-Zahlenformate sind IMMER US-notiert: Komma = Tausender,
            # Punkt = Dezimaltrenner. PowerPoint lokalisiert die Anzeige selbst.
            # Das deutsche Muster '#.##0' bedeutet dort drei Nachkommastellen –
            # aus 923 wurde im PDF-Test '923,000'.
            plot.data_labels.number_format = "#,##0"
            plot.data_labels.number_format_is_linked = False
    except Exception:  # noqa: BLE001
        pass       # ein Diagramm ohne Feinschliff ist besser als keines
    return True


def _logo_auf_titelfolie(slide, prs) -> None:
    """Setzt das Branding-Logo oben rechts auf die Titelfolie.

    Position und Hoehe stammen aus der Firmenvorlage (dort sitzt das Logo im
    Master, also auf jeder Folie). Hier steht es nur auf der TITELFOLIE, weil
    ``MasterShapes`` in python-pptx keine Bilder aufnehmen kann – die
    Medien-Beziehung haengt an der Folie, nicht am Master.

    Die rechte Kante liegt auf dem Satzspiegel, damit Logo und Textblock
    dieselbe Fluchtlinie haben; die Breite ergibt sich aus dem
    Seitenverhaeltnis der Datei (ein festes Mass wuerde ein quadratisches Logo
    verzerren)."""
    try:
        from skills.office.vorlage import _logo_pfad, LOGO_Y, LOGO_H, LOGO_RECHTS
        logo = _logo_pfad()
        if not logo:
            return
        bild = slide.shapes.add_picture(str(logo), 0, LOGO_Y, height=LOGO_H)
        bild.left = max(0, LOGO_RECHTS - bild.width)
        bild.name = "Logo"
    except Exception:  # noqa: BLE001
        pass       # ein fehlendes/kaputtes Logo darf die Praesentation nicht kosten


class CreatePowerPointTool(BaseTool):
    @property
    def name(self) -> str:
        return "office_create_powerpoint"

    @property
    def description(self) -> str:
        return (
            "Erstellt eine PowerPoint-Praesentation (.pptx) im HAUSDESIGN – 16:9, echte "
            "Masterfolien, Layouts und Platzhalter, Farben und Schrift aus der Vorlage. "
            "'slides' ist eine Liste von Folien-Objekten: "
            "{ 'title': 'Folientitel', 'bullets': ['Punkt 1','Punkt 2'] } oder "
            "{ 'title': ..., 'content': 'Freitext' }. Unterpunkte mit '> ' voranstellen "
            "('> Detail'). Optional je Folie 'layout' ('inhalt' Standard, 'abschnitt' fuer "
            "einen Kapiteltrenner, 'zwei' fuer zwei Spalten, 'nurtitel', 'leer') und "
            "'notes' fuer Sprechernotizen. Eine Titelfolie am Anfang entsteht ueber die "
            "Parameter 'title'/'subtitle' (nicht als Folien-Objekt); wer sie doch als Folie "
            "schickt, nimmt { 'layout': 'titel', 'title': ..., 'subtitle': ... }. "
            "ZAHLEN GEHOEREN IN EIN DIAGRAMM, nicht in eine Aufzaehlung: je Folie "
            "{ 'chart': { 'typ': 'saeulen|balken|linie|kreis', 'kategorien': [...], "
            "'werte': [...] } } – mehrere Datenreihen ueber 'reihen': [{'name':…, "
            "'werte':[...]}, …]. Das Diagramm ist in PowerPoint bearbeitbar und nimmt "
            "seine Farben aus dem Hausdesign. Stehen daneben 'bullets', ruecken beide "
            "nebeneinander. "
            "'slides' MUSS eine echte Liste sein – KEIN Text, der eine Liste enthaelt. "
            "KEINE Farb-, Schrift- oder Groessenangaben mitschicken – die kommen aus der "
            "Vorlage; eigene Werte brechen das Design beim Bearbeiten. "
            "Gibt eine Download-URL zurueck."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "filename": {"type": "STRING", "description": "Dateiname (ohne Pfad)."},
                "title": {"type": "STRING", "description": "Optionaler Titel fuer eine Titelfolie am Anfang."},
                "subtitle": {"type": "STRING", "description": "Optionaler Untertitel der Titelfolie."},
                "slides": {"type": "ARRAY", "items": {"type": "OBJECT"},
                           "description": "Liste der Inhaltsfolien als echtes Array von "
                                          "Objekten (siehe Beschreibung) – nicht als Text. "
                                          "Je Folie optional 'chart' fuer ein Diagramm."},
                "template": {"type": "STRING", "description": "Optionaler Vorlagen-Dateiname (leer = Hausvorlage). Verfuegbare zeigt office_template_info."},
            },
            "required": ["filename", "slides"],
        }

    async def execute(self, filename: str = "", title: str = "", subtitle: str = "",
                       slides=None, template: str = "", **kwargs) -> str:
        if not filename:
            return "Fehler: 'filename' ist Pflicht."
        try:
            from pptx import Presentation
        except Exception as e:
            return f"Fehler: python-pptx nicht verfuegbar ({e})."

        # VOR allem anderen: der Foliensatz muss eine Liste sein. Kommt er als
        # Text, wuerde die Schleife unten ueber die Zeichen laufen.
        slides, fehler = _slides_normalisieren(slides)
        if fehler:
            return fehler
        zuviel = 0
        if len(slides) > MAX_FOLIEN:
            zuviel = len(slides) - MAX_FOLIEN
            slides = slides[:MAX_FOLIEN]

        # Vorlage: 16:9, Branding-Farben, echte Masterfolien. Faellt das
        # Erzeugen aus (Rechte, fehlendes lxml), wird OHNE Vorlage
        # weitergearbeitet – eine Praesentation im Standarddesign ist besser
        # als eine Fehlermeldung.
        vorlage_hinweis = ""
        try:
            from skills.office import vorlage as _vorlage
            prs = Presentation(str(_vorlage.loese_vorlage(template)))
        except Exception as e:  # noqa: BLE001
            prs = Presentation()
            vorlage_hinweis = f" (Hinweis: Hausvorlage nicht nutzbar – {e})"

        if title:
            slide = prs.slides.add_slide(_layout(prs, "titel"))
            _titel_setzen(slide, _text_bereinigen(title))
            felder = _text_platzhalter(slide)
            if subtitle and felder:
                _fuelle(felder[0], text=_text_bereinigen(subtitle))
            _leere_platzhalter_entfernen(slide)
            _logo_auf_titelfolie(slide, prs)

        for sl in slides:
            art = str(sl.get("layout") or "").strip().lower()
            if art not in _LAYOUT_ALIAS:
                art = "inhalt"
            slide = prs.slides.add_slide(_layout(prs, art))
            _titel_setzen(slide, sl.get("title", ""))

            felder = _text_platzhalter(slide)
            bullets = sl.get("bullets")
            # 'subtitle' als Alias: die Werkzeugbeschreibung nennt es im selben
            # Absatz wie die Folien-Felder (dort ist der TOP-LEVEL-Parameter
            # gemeint), und ein Modell schickt es deshalb regelmaessig je Folie
            # mit – bei 'layout: titel' ist das sogar die naheliegende
            # Schreibweise. Ohne den Alias fiele der Text wortlos weg und die
            # Titelfolie bliebe ohne Kicker (beim Abnahmelauf auf ECHT genau so
            # passiert).
            inhalt = sl.get("content") or sl.get("text") or sl.get("subtitle") or ""
            diagramm = sl.get("chart") or sl.get("diagramm")
            hat_text = bool(bullets) or bool(inhalt)

            if felder:
                if art == "zwei" and len(felder) > 1 and isinstance(bullets, (list, tuple)) and len(bullets) > 1:
                    # Zwei-Spalten-Layout: die Aufzaehlung in der Mitte teilen,
                    # sonst bliebe die rechte Spalte leer (und wuerde entfernt).
                    mitte = (len(bullets) + 1) // 2
                    _fuelle(felder[0], bullets=list(bullets[:mitte]))
                    _fuelle(felder[1], bullets=list(bullets[mitte:]))
                else:
                    _fuelle(felder[0], bullets=bullets, text=inhalt)
            _leere_platzhalter_entfernen(slide)

            if diagramm:
                _diagramm_einfuegen(slide, prs, diagramm, neben_text=hat_text)

            notizen = sl.get("notes") or sl.get("notizen")
            if notizen:
                try:
                    slide.notes_slide.notes_text_frame.text = str(notizen)
                except Exception:  # noqa: BLE001
                    pass

        disk, fname, dl = _new_path(filename, "pptx")
        try:
            prs.save(str(disk))
        except Exception as e:
            return f"Fehler beim Speichern: {e}"
        # Ein Deckel, der nicht genannt wird, ist ein stiller Datenverlust: der
        # Aufrufer haelt die gekuerzte Datei sonst fuer vollstaendig.
        if zuviel:
            vorlage_hinweis += (f" (Hinweis: {zuviel} weitere Folien wurden NICHT "
                                f"uebernommen – hoechstens {MAX_FOLIEN} je Datei. "
                                f"Fuer mehr Material eine Tabelle erzeugen.)")
        return _ok(dl, fname, disk, extra=vorlage_hinweis)


class TemplateInfoTool(BaseTool):
    """Zeigt, was eine Vorlage anbietet.

    Ohne dieses Werkzeug muesste das Modell die Layout- und Platzhalternamen
    einer mitgebrachten Firmenvorlage raten. Es liest nur – es erzeugt keine
    Datei und veraendert keine Vorlage."""

    @property
    def name(self) -> str:
        return "office_template_info"

    @property
    def description(self) -> str:
        return (
            "Zeigt die verfuegbaren PowerPoint-Vorlagen und die Layouts der gewaehlten "
            "Vorlage (Name, Platzhalter, Foliengroesse). Nuetzlich, wenn eine eigene "
            "Firmenvorlage hinterlegt wurde und unklar ist, welche Layouts es gibt."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "template": {"type": "STRING", "description": "Vorlagen-Dateiname (leer = Hausvorlage)."},
            },
            "required": [],
        }

    async def execute(self, template: str = "", **kwargs) -> str:
        try:
            from pptx import Presentation
            from skills.office import vorlage as _vorlage
        except Exception as e:  # noqa: BLE001
            return f"Fehler: python-pptx/Vorlagenmodul nicht verfuegbar ({e})."
        try:
            pfad = _vorlage.loese_vorlage(template)
            prs = Presentation(str(pfad))
        except Exception as e:  # noqa: BLE001
            return f"Fehler: Vorlage nicht lesbar ({e})."

        breite_cm = round(prs.slide_width / 360000, 1)
        hoehe_cm = round(prs.slide_height / 360000, 1)
        seiten = "16:9" if abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.02 else (
            "4:3" if abs(prs.slide_width / prs.slide_height - 4 / 3) < 0.02 else "andere")
        zeilen = [
            f"Vorlage: {pfad.name} ({breite_cm}×{hoehe_cm} cm, {seiten})",
            f"Vorhandene Vorlagen: {', '.join(_vorlage.verfuegbare()) or '(keine)'}",
            "",
            "Layouts (Name → befuellbare Platzhalter):",
        ]
        for i, l in enumerate(prs.slide_layouts):
            namen = []
            for ph in l.placeholders:
                try:
                    namen.append(f"{ph.placeholder_format.type}".split(".")[-1].split(" ")[0])
                except Exception:  # noqa: BLE001
                    continue
            zeilen.append(f"  [{i}] {l.name} → {', '.join(namen) or '(keine)'}")
        zeilen += [
            "",
            "Kurznamen fuer 'layout' in office_create_powerpoint: "
            + ", ".join(sorted(_LAYOUT_ALIAS.keys())),
        ]
        return "\n".join(zeilen)


# ─────────────────────────────────────────────────────────────────────────
# Lesen
# ─────────────────────────────────────────────────────────────────────────
# Ab dieser Textmenge ist eine Tabelle ueber office_read nicht mehr sinnvoll
# lesbar: der Agent kappt Werkzeug-Ergebnisse, und was ankaeme, waere ein
# Bruchteil, den das Modell fuer das Ganze haelt. Statt eines Fragments gibt es
# dann den STRUKTUR-Ueberblick und den Verweis auf die Tabellen-Werkzeuge.
_XLSX_TEXT_MAX = 18000


def _xlsx_zu_gross(p: Path, wb, blocks: list) -> str:
    """Antwort fuer eine Tabelle, die als Text nicht sinnvoll uebergeben werden kann.

    DAS IST KEIN FEHLER, SONDERN DIE RICHTIGE ANTWORT. Im Vorfall vom
    2026-08-19 lieferte office_read fuer eine Mappe mit 362.195 Zellen ein
    0,4-%-Fragment; das Modell baute daraus eine "zusammengefuehrte" Tabelle
    mit zwei Zeilen. Ein ehrlicher Verweis auf das passende Werkzeug ist jeder
    Teilantwort ueberlegen.
    """
    zeilen = [f"Diese Tabelle ist zu gross, um sie als Text zu uebergeben "
              f"({sum(len(b) for b in blocks)}+ Zeichen). "
              f"Es wurde NICHTS davon gelesen – hier nur der Aufbau:", ""]
    try:
        for ws in wb.worksheets:
            zeilen.append(f"  Blatt '{ws.title}': {ws.max_row or 0} Zeilen "
                          f"x {ws.max_column or 0} Spalten")
    except Exception:  # noqa: BLE001
        pass
    try:
        wb.close()
    except Exception:  # noqa: BLE001
        pass
    zeilen += [
        "",
        "BENUTZE STATTDESSEN:",
        "  xlsx_inspect     – Aufbau, Kopfzeilen, Datentypen, Beispielzeilen",
        "  xlsx_read_range  – einen begrenzten Ausschnitt wirklich lesen",
        "  xlsx_merge       – zwei Tabellen zusammenfuehren (Daten laufen nicht",
        "                     durch dich; Formeln und Layout bleiben erhalten)",
        "  xlsx_edit        – einzelne Zellen schreiben",
        "",
        "Tippe die Daten NICHT ab und baue die Tabelle NICHT mit "
        "office_create_excel nach – bei dieser Groesse wird das Ergebnis "
        "zwangslaeufig unvollstaendig.",
    ]
    return "\n".join(zeilen)


class ReadDocumentTool(BaseTool):
    # Der Ergebnis-Deckel des Agenten (5.000) wuerde den ohnehin schon
    # gekuerzten Text ein zweites Mal kappen. 20.000 entspricht dem eigenen
    # Deckel unten, damit genau EINE Kuerzung stattfindet – und die wird
    # beziffert.
    ergebnis_max = 21000

    @property
    def name(self) -> str:
        return "office_read"

    @property
    def description(self) -> str:
        return (
            "Liest den Textinhalt eines Office-Dokuments (.docx, .xlsx, .pptx) und gibt ihn "
            "als Text zurueck. 'path' kann ein Dateiname aus data/documents/, eine "
            "/api/documents/-URL oder ein beliebiger Server-Pfad sein. "
            "FUER TABELLENDATEN ist xlsx_inspect der richtige Einstieg – dieses "
            "Werkzeug macht aus einer Mappe Fliesstext und kann sie nicht bearbeiten."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "path": {"type": "STRING", "description": "Pfad/Name des zu lesenden Dokuments."},
            },
            "required": ["path"],
        }

    async def execute(self, path: str = "", **kwargs) -> str:
        p = _resolve_existing(path)
        if not p:
            return f"Fehler: Datei nicht gefunden: {path}"
        ext = p.suffix.lower()
        try:
            if ext == ".docx":
                from docx import Document
                doc = Document(str(p))
                parts = [par.text for par in doc.paragraphs if par.text.strip()]
                for tbl in doc.tables:
                    for row in tbl.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
                text = "\n".join(parts)
            elif ext == ".xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(str(p), read_only=True, data_only=True)
                blocks = []
                zu_gross = None
                for ws in wb.worksheets:
                    blocks.append(f"# Blatt: {ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        blocks.append(" | ".join("" if c is None else str(c) for c in row))
                    # Sobald absehbar ist, dass ohnehin fast alles wegfaellt,
                    # wird gar nicht erst weitergelesen – die Antwort ist dann
                    # der Verweis auf xlsx_inspect, kein Fragment.
                    if sum(len(b) for b in blocks) > _XLSX_TEXT_MAX * 4:
                        zu_gross = True
                        break
                if zu_gross or sum(len(b) for b in blocks) > _XLSX_TEXT_MAX:
                    return _xlsx_zu_gross(p, wb, blocks)
                text = "\n".join(blocks)
            elif ext == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(p))
                blocks = []
                for i, slide in enumerate(prs.slides, 1):
                    blocks.append(f"# Folie {i}")
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            blocks.append(shape.text_frame.text)
                text = "\n".join(blocks)
            else:
                return f"Fehler: Nicht unterstuetzte Endung '{ext}' (docx/xlsx/pptx)."
        except Exception as e:
            return f"Fehler beim Lesen: {e}"

        # KUERZUNG WIRD BEZIFFERT – UND STEHT VORNE.
        # Vorher stand hier ein blosses "… [gekuerzt]" am ENDE. Zwei Deckel
        # hintereinander machten daraus im Vorfall vom 2026-08-19 eine
        # doppelte Luege: office_read kuerzte still von 1.265.130 auf 20.000,
        # danach kappte der Agent auf 5.000 und meldete "5.000 von 20.014" –
        # das Modell hielt sich fuer gut informiert und sah 0,4 % der Datei.
        # Der Hinweis muss deshalb (a) die ECHTE Gesamtgroesse nennen und
        # (b) AM ANFANG stehen, damit ihn der Deckel des Agenten nicht
        # ebenfalls abschneidet.
        if len(text) > 20000:
            text = (f"[GEKUERZT: {20000} von {len(text)} Zeichen "
                    f"({20000 * 100 // max(1, len(text))} %). Der Rest FEHLT. "
                    f"Ziehe daraus keine Schluesse ueber die Gesamtheit – frage "
                    f"gezielt nach dem fehlenden Teil.]\n\n" + text[:20000])
        return text or "(leeres Dokument)"


# ─────────────────────────────────────────────────────────────────────────
# PDF-Export via LibreOffice
# ─────────────────────────────────────────────────────────────────────────
# Wo LibreOffice liegen kann. `soffice` steht bei Debian-Paketen in /usr/bin,
# eigenstaendige Installationen (LibreOffice-Download, Flatpak-Export) legen es
# nur unter /opt bzw. /usr/lib ab.
_SOFFICE_KANDIDATEN = (
    "soffice", "libreoffice",
    "/usr/bin/soffice", "/usr/lib/libreoffice/program/soffice",
    "/opt/libreoffice/program/soffice", "/snap/bin/libreoffice",
)


def _find_soffice() -> str | None:
    """Sucht die LibreOffice-Binaerdatei; None = nicht installiert."""
    import shutil as _sh
    for k in _SOFFICE_KANDIDATEN:
        if k.startswith("/"):
            if Path(k).is_file():
                return k
        else:
            gefunden = _sh.which(k)
            if gefunden:
                return gefunden
    return None


class ExportPdfTool(BaseTool):
    @property
    def name(self) -> str:
        return "office_to_pdf"

    @property
    def description(self) -> str:
        return (
            "Exportiert ein Office-Dokument (.docx/.xlsx/.pptx) nach PDF (via LibreOffice). "
            "'path' kann ein Dateiname aus data/documents/, eine /api/documents/-URL oder ein "
            "Server-Pfad sein. Gibt eine Download-URL fuer das PDF zurueck."
        )

    def parameters_schema(self) -> dict:
        return {
            "type": "OBJECT",
            "properties": {
                "path": {"type": "STRING", "description": "Pfad/Name des zu konvertierenden Dokuments."},
            },
            "required": ["path"],
        }

    async def execute(self, path: str = "", **kwargs) -> str:
        src = _resolve_existing(path)
        if not src:
            return f"Fehler: Datei nicht gefunden: {path}"
        if src.suffix.lower() not in (".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"):
            return f"Fehler: Format '{src.suffix}' wird fuer PDF-Export nicht unterstuetzt."

        # Fehlt LibreOffice, gab es bis 2026-07-28 nur das rohe
        # "[Errno 2] No such file or directory: 'soffice'" – daraus konnte weder
        # das LLM noch der Nutzer ableiten, WAS zu tun ist. Jetzt eine Meldung,
        # die den Grund nennt und den Weg zeigt (auf ECHT war LibreOffice nie
        # installiert, auf DEV schon – daher "geht bei mir, dort nicht").
        binary = _find_soffice()
        if not binary:
            return (
                "Fehler: PDF-Export nicht moeglich – auf diesem Server ist LibreOffice "
                "nicht installiert (weder 'soffice' noch 'libreoffice' gefunden). "
                "Ein Administrator behebt das unter Einstellungen → Skills: der Office-Skill "
                "zeigt dann die Plakette 'Abhängigkeit fehlt' und daneben den Knopf ⤓ "
                "'Fehlende Abhängigkeiten nachinstallieren'. Alternativ von Hand: "
                "'apt install libreoffice-writer libreoffice-calc libreoffice-impress'. "
                "Laeuft der Root-Broker, muss die apt-Installation ggf. einmal unter "
                "Einstellungen → Sicherheit → Root-Freigaben freigegeben werden. "
                "Das Office-Dokument selbst wurde erzeugt und kann heruntergeladen werden."
            )

        DOCS_DIR.mkdir(parents=True, exist_ok=True)
        token = uuid.uuid4().hex
        # Eigenes UserInstallation-Profil, um Konflikte mit dem Desktop-LibreOffice zu vermeiden
        profile = f"/tmp/lo_jarvis_{token}"
        cmd = [
            binary, "--headless", "--norestore", "--convert-to", "pdf",
            "--outdir", str(DOCS_DIR),
            f"-env:UserInstallation=file://{profile}",
            str(src),
        ]
        try:
            proc = await asyncio.to_thread(
                subprocess.run, cmd,
                capture_output=True, text=True, timeout=120,
            )
        except subprocess.TimeoutExpired:
            return "Fehler: PDF-Export hat das Zeitlimit (120s) ueberschritten."
        except Exception as e:
            return f"Fehler beim PDF-Export: {e}"

        # soffice legt <stem>.pdf in outdir ab
        produced = DOCS_DIR / (src.stem + ".pdf")
        if not produced.exists():
            return f"Fehler: PDF wurde nicht erzeugt. soffice: {proc.stderr or proc.stdout}".strip()

        # In Capability-Schema umbenennen (Download-Name aus Original-Basis ableiten)
        base = src.stem.split("__", 1)[-1] if "__" in src.stem else src.stem
        disk, fname, dl = _new_path(base, "pdf")
        try:
            produced.rename(disk)
        except Exception:
            # Fallback: Inhalt kopieren
            disk.write_bytes(produced.read_bytes())
            produced.unlink(missing_ok=True)
        return _ok(dl, fname, disk)


def get_tools():
    tools = [
        CreateWordTool(),
        CreateExcelTool(),
        CreatePowerPointTool(),
        TemplateInfoTool(),
        ReadDocumentTool(),
        ExportPdfTool(),
    ]
    # Tabellen-Werkzeuge (ansehen/bearbeiten statt neu aufbauen). Bewusst in
    # einem eigenen Modul und hier ANGEHAENGT statt eingebaut: faellt der Import
    # aus (alte openpyxl-Version, Teil-Deploy), bleibt der Office-Skill mit
    # seinen bisherigen Werkzeugen benutzbar, statt komplett auszufallen.
    try:
        from skills.office.tabellen import get_tabellen_tools
        tools.extend(get_tabellen_tools())
    except Exception as e:  # noqa: BLE001
        print(f"[Office] Tabellen-Werkzeuge nicht geladen: {e}", flush=True)
    # Formular-PDFs (eine Seite = ein Datensatz). Aus demselben Grund
    # angehaengt und einzeln abgesichert wie die Tabellen-Werkzeuge: faellt
    # der Import aus (pdfplumber/pytesseract fehlt), bleibt der Office-Skill
    # mit allen uebrigen Werkzeugen benutzbar.
    try:
        from skills.office.pdf_formular import get_pdf_formular_tools
        tools.extend(get_pdf_formular_tools())
    except Exception as e:  # noqa: BLE001
        print(f"[Office] PDF-Formular-Werkzeug nicht geladen: {e}", flush=True)
    return tools
