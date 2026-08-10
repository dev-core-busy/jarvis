"""PowerPoint-Vorlage im Hausdesign (NEXUS-Profil + Branding-Farbe).

WARUM ES DAS GIBT: ``office_create_powerpoint`` rief bisher ``Presentation()``
OHNE Argument auf. Damit gilt das eingebaute Standarddesign von python-pptx –
und das ist der Grund, warum erzeugte Decks „nicht wie eine Firmenpräsentation"
aussehen:

  * **4:3** (9144000 × 6858000 EMU). Seit rund zehn Jahren ist 16:9 der
    Standard; 4:3 erkennt jeder Betrachter sofort als veraltet.
  * **Calibri** und das Office-Standard-Blau – kein Bezug zum Branding.
  * Alle Farben liegen als *Theme*-Werte im Master. Wer sie pro Folie
    überschreibt (Textfarbe je Absatz setzen), bekommt eine Datei, die beim
    Bearbeiten in PowerPoint auseinanderfällt: der Designer zeigt dann andere
    Farben als die Folien. Deshalb wird hier das THEME geändert, nicht die
    einzelne Folie.

ERGEBNIS: eine .pptx mit echten Masterfolien, Layouts und Platzhaltern, die
PowerPoint wie eine von Hand angelegte Vorlage behandelt. Der Agent befüllt nur
Platzhalter – Schriftgrößen, Aufzählungsebenen und Abstände kommen aus der
Vorlage.

WARUM DAS THEME-XML DIREKT IM ZIP GEPATCHT WIRD: python-pptx hat keine API für
``a:clrScheme``/``a:fontScheme``. Der Weg über die internen Part-Objekte wäre an
die Version gebunden; eine .pptx ist ein ZIP, und ``ppt/theme/theme1.xml`` ist
eine gewöhnliche Datei darin. Das ist stabil und prüfbar (der Test liest die
Farben genau so wieder aus).

WOHER DIE WERTE STAMMEN: Alle Farben, Schriften, Schriftgrößen und Positionen
unten sind aus der Firmenvorlage ``NEXUS_PowerPoint-Template_LAB_2025.potx``
ausgelesen (Theme ``nexus``, Farbschema ``NEXUS``, Schriftschema
``NEXUS-Font``) – nicht geschätzt. Übernommen wurde das GESTALTUNGSSYSTEM
(Farbrollen, Typo-Stufen, Satzspiegel, Kapitelfolie), nicht das Bildmaterial:
die Originalvorlage bringt 17 Grafiken (Hexagon-Welt, Logos, Vollbilder) und
873 KB mit. Eine generierte Vorlage bleibt bei ~40 KB, ist prüfbar und bleibt
white-label-fähig – ein anderer Branding-Akzent schlägt auf ``accent1`` durch.

SCHRIFTWAHL – bewusst ANDERS als bei matplotlib: dort zählt, was auf dem SERVER
installiert ist (`backend/plotstyles/jarvis.mplstyle` nennt DejaVu/Liberation).
Eine .pptx wird dagegen auf einem FREMDEN Rechner geöffnet, und dort steht die
CI-Schrift. Deshalb tragen Theme und Textstile die Originalschriften
``HelveticaNeue LT 75 Bold`` / ``55 Roman``.

  ACHTUNG PDF-EXPORT: Auf dem Server ist HelveticaNeue LT NICHT installiert
  (``fc-match`` liefert Noto Sans). ``office_to_pdf`` setzt also eine
  Ersatzschrift – Zeilenumbrüche und Textlängen können dort vom Original
  abweichen. Für Empfänger mit lizenzierter Schrift (Windows/macOS) stimmt das
  Deck; wer maßhaltige PDFs vom Server braucht, hinterlegt eine metrisch
  kompatible Schrift (Nimbus Sans / TeX Gyre Heros) oder stellt
  ``SCHRIFT_TITEL``/``SCHRIFT_TEXT`` auf ``Arial``.
"""

import re
import zipfile
from pathlib import Path

# Vorlagen liegen neben den erzeugten Dokumenten, aber in einem EIGENEN Ordner:
# data/documents wird nach Frist aufgeraeumt (documents.cleanup_old) und ist
# eigentuemergebunden – eine Vorlage darf weder verfallen noch einem Benutzer
# gehoeren.
VORLAGEN_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "vorlagen"
STANDARD_NAME = "standard.pptx"

# 16:9 in EMU (13.333 × 7.5 Zoll)
BREITE_16_9 = 12192000
HOEHE_16_9 = 6858000

# ─────────────────────────────────────────────────────────────────────────────
# DESIGNPROFIL – ausgewertet aus der Firmenvorlage (Theme "nexus")
# ─────────────────────────────────────────────────────────────────────────────
# Die Hausfarbe. Sie belegt accent1 UND die Hyperlinkfarbe – so haelt es die
# Originalvorlage. Ein konfigurierter Branding-Akzent ersetzt sie (siehe
# branding_farben): das ist die einzige Stelle, an der eine Fremdmarke das
# Profil uebersteuert.
STANDARD_AKZENT = "B80F2E"

# accent2..accent6 in der Reihenfolge der Originalvorlage. DAS IST EINE
# BEWUSSTE ABWEICHUNG von frontend/js/charts.js und jarvis.mplstyle: jene Reihe
# (Blau/Gruen/Orange/Rot/Cyan) ist fuer Bildschirm-Diagramme gemacht. Ein
# Diagramm, das IN PowerPoint eingefuegt wird, zieht seine Farben aus genau
# diesen Slots – mit der Web-Reihe saehe es auf der Folie wie ein Fremdkoerper
# aus. Innerhalb einer Praesentation gewinnt also das Hausdesign.
PALETTE = ["4F6792", "1F2336", "E8ECF0", "9C9D9F", "BA4C61"]
FOL_HLINK = "BA4C61"          # besuchter Link (accent6), wie im Original

# Grundfarben. dk1/lt1 sind im Original 'sysClr windowText/window', also
# Schwarz auf Weiss – bewusst uebernommen statt "weicher" gemacht: der Text
# eines Jarvis-Decks soll neben einem von Hand gebauten Deck nicht auffallen.
HELL = "FFFFFF"
DUNKEL = "000000"
GRAU = "9C9D9F"               # accent5, Sekundaertext/Linien
FLAECHE = "E8ECF0"            # accent4, helle Fuellflaeche
GRAU_BG = "E0E8EB"            # Folienhintergrund der grauen Layout-Variante

SCHRIFT_TITEL = "HelveticaNeue LT 75 Bold"   # majorFont
SCHRIFT_TEXT = "HelveticaNeue LT 55 Roman"   # minorFont

# ── Satzspiegel (EMU) ────────────────────────────────────────────────────────
# Aus den Layouts "Standard"/"Title" der Originalvorlage. Der linke Rand von
# 700679 EMU (1,84 cm) ist das praegende Mass: er gilt fuer Titel, Subline und
# Inhalt gleichermassen, dadurch stehen alle Textkanten auf einer Linie.
RAND_LINKS = 700679
# Die Firmenvorlage nennt 10822650 – bei ihrer Folienbreite (12190413 EMU)
# bleibt rechts dadurch 0,09 cm weniger Rand als links. Das ist eine Rundung
# aus dem 4:3-Ursprung, kein Gestaltungswille; hier wird die Breite aus dem
# Rand GERECHNET, damit der Satzspiegel exakt mittig steht.
INHALT_B = BREITE_16_9 - 2 * RAND_LINKS
TITEL_Y, TITEL_H = 673096, 504916
SUB_Y, SUB_H = 1215038, 395288        # Subline/Kicker unter dem Titel
INHALT_Y, INHALT_H = 1827356, 3833215
SPALTEN_LUFT = 342900                 # 0,95 cm zwischen zwei Inhaltsspalten

# Titelfolie: KICKER OBEN, grosser Titel darunter – nicht umgekehrt. Genau das
# unterscheidet die Vorlage vom Office-Standard (dort steht der Untertitel
# unter dem Titel).
T_KICKER_Y, T_KICKER_H = 1957837, 733247
T_TITEL_Y, T_TITEL_H = 2997585, 1325563

# Kapitel-/Abschnittsfolie: farbiger Kasten unten links, Titel weiss darin.
KAP_X, KAP_Y, KAP_B, KAP_H = 698500, 4198652, 5553064, 1865327

# Akzentstrich unter dem Titel der Titelfolie (Ersatz fuer das Vollbild der
# Originalvorlage): ein Sechstel der Satzbreite, 0,12 cm stark.
TITELSTRICH_B = INHALT_B // 6
TITELSTRICH_H = 45720

# Foliennummer rechts unten (Master).
NUM_X, NUM_Y, NUM_B, NUM_H = 11153375, 6485059, 423019, 236416

# Logo oben rechts: Hoehe und Grundlinie aus der Firmenvorlage. Die BREITE
# ergibt sich aus dem Seitenverhaeltnis der jeweiligen Logodatei – ein festes
# Mass wuerde ein quadratisches Logo verzerren.
LOGO_Y = 333517
LOGO_H = 265625
LOGO_RECHTS = RAND_LINKS + INHALT_B      # rechte Kante des Satzspiegels

# ── Typografie (Hundertstel-Punkt) ───────────────────────────────────────────
# titleStyle/bodyStyle der Originalvorlage. Die Office-Vorgabe (Titel 44 pt,
# Text 28 pt) ist fuer 4:3 gemacht und wirkt auf 16:9 grob.
TYPO_TITEL = 3200                       # Folientitel, fett
TYPO_BODY = [1800, 1600, 1400, 1200, 1200]   # Aufzaehlungsebenen 1..5
TYPO_TITELFOLIE = 4800                  # Titel der Titelfolie
TYPO_KICKER = 2800                      # Zeile ueber dem Titel der Titelfolie
TYPO_SUB = 2400                         # Unterzeile / Spaltenueberschrift
TYPO_KAPITEL = 3600                     # Titel im Kapitelkasten

# Einrueckschritt je Aufzaehlungsebene und Absatzabstand – beide aus dem
# Original (271463 EMU = 0,71 cm; 6 pt nach jedem Absatz).
BULLET_EINZUG = 271463
ABSTAND_NACH = 600                      # spcAft in Hundertstel-Punkt
ABSTAND_VOR = 700                       # spcBef, nur vor einem Hauptpunkt
KAP_INNEN = 274320                      # Innenabstand im Kapitelkasten (0,76 cm)

# Deutsche Layout-Namen. Der Agent spricht Layouts ueber den NAMEN an (siehe
# main.py::_layout); englische Namen aus dem Standardtemplate waeren fuer eine
# deutschsprachige Oberflaeche eine unnoetige Huerde. Der englische Name bleibt
# als Alias in main.py erhalten, damit eine FREMDE Firmenvorlage weiter passt.
LAYOUT_NAMEN = {
    "Title Slide": "Titelfolie",
    "Title and Content": "Titel und Inhalt",
    "Section Header": "Abschnitt",
    "Two Content": "Zwei Inhalte",
    "Comparison": "Vergleich",
    "Title Only": "Nur Titel",
    "Blank": "Leer",
    "Content with Caption": "Inhalt mit Beschriftung",
    "Picture with Caption": "Bild mit Beschriftung",
}


def _hex(wert, fallback: str) -> str:
    """Normiert eine Farbangabe auf 6 Hex-Ziffern ohne '#'.

    Branding-Werte kommen aus einem Eingabefeld: '#9B59B6', '9b59b6' und '#abc'
    sind alle moeglich. Alles andere (rgb(), Farbnamen, Leerstring) faellt auf
    den Vorgabewert zurueck – eine ungueltige Farbe im Theme macht die Datei
    fuer PowerPoint unlesbar."""
    s = str(wert or "").strip().lstrip("#")
    if re.fullmatch(r"[0-9a-fA-F]{6}", s):
        return s.upper()
    if re.fullmatch(r"[0-9a-fA-F]{3}", s):
        return "".join(c * 2 for c in s).upper()
    return fallback


def branding_farben() -> dict:
    """Liest die Branding-Werte und leitet daraus die Theme-Farben ab.

    Genommen wird der Akzent aus dem HELL-Modus (``colors_light``), sonst der
    aus dem Dunkel-Modus. Hintergrund und Textfarbe kommen NICHT aus dem
    Branding: die Chat-Oberflaeche ist dunkel, eine Praesentation muss hell
    sein. Uebernommen wird also genau das, was uebertragbar ist – die
    Markenfarbe."""
    cfg = {}
    try:
        from backend.config import config
        states = config.get_skill_states()
        st = states.get("branding") or {}
        if st.get("enabled"):
            cfg = st.get("config") or {}
    except Exception:  # noqa: BLE001
        cfg = {}
    hell = (cfg.get("colors_light") or {}) if isinstance(cfg.get("colors_light"), dict) else {}
    dunkel = (cfg.get("colors") or {}) if isinstance(cfg.get("colors"), dict) else {}
    akzent = _hex(hell.get("accent") or dunkel.get("accent"), STANDARD_AKZENT)
    return {
        "akzent": akzent,
        "firma": str(cfg.get("company_name") or "").strip(),
        "hell": HELL,
        "dunkel": DUNKEL,
        "grau": GRAU,
        "flaeche": FLAECHE,
    }


def _theme_xml(alt: bytes, farben: dict,
               schrift_titel: str = SCHRIFT_TITEL,
               schrift_text: str = SCHRIFT_TEXT) -> bytes:
    """Ersetzt Farb- und Schriftschema im Theme-XML.

    Ueberschrift und Grundtext bekommen UNTERSCHIEDLICHE Schriften
    (majorFont/minorFont) – so haelt es die Firmenvorlage, und nur so wirkt der
    Fettschnitt der Ueberschrift auch dann, wenn ein Empfaenger die
    Schriftfamilie nur teilweise installiert hat."""
    from lxml import etree

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns = {"a": A}
    root = etree.fromstring(alt)

    def srgb(el, wert):
        """Setzt den Farbwert eines Schema-Slots (dk1/lt1/accent1 …).
        Der Slot kann <a:srgbClr> ODER <a:sysClr> enthalten – letzteres bei
        dk1/lt1 im Standardtemplate ('windowText'/'window'). sysClr wird
        ersetzt, sonst blieben Text- und Hintergrundfarbe unveraendert."""
        for kind in list(el):
            el.remove(kind)
        neu = etree.SubElement(el, f"{{{A}}}srgbClr")
        neu.set("val", wert)

    clr = root.find(".//a:themeElements/a:clrScheme", ns)
    if clr is not None:
        werte = {
            "dk1": farben["dunkel"], "lt1": farben["hell"],
            "dk2": farben["dunkel"], "lt2": farben["hell"],
            "accent1": farben["akzent"],
            "accent2": PALETTE[0], "accent3": PALETTE[1], "accent4": PALETTE[2],
            "accent5": PALETTE[3], "accent6": PALETTE[4],
            "hlink": farben["akzent"], "folHlink": FOL_HLINK,
        }
        for slot, wert in werte.items():
            el = clr.find(f"a:{slot}", ns)
            if el is not None:
                srgb(el, wert)

    fonts = root.find(".//a:themeElements/a:fontScheme", ns)
    if fonts is not None:
        for teil, schrift in (("a:majorFont", schrift_titel),
                              ("a:minorFont", schrift_text)):
            f = fonts.find(teil, ns)
            if f is None:
                continue
            latin = f.find("a:latin", ns)
            if latin is not None:
                latin.set("typeface", schrift)
                # panose/pitchFamily des Standardtemplates gehoeren zu Calibri
                # und wuerden PowerPoint eine falsche Ersatzschrift waehlen
                # lassen.
                for attr in ("panose", "pitchFamily", "charset"):
                    if attr in latin.attrib:
                        del latin.attrib[attr]
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _patch_theme(pptx_pfad: Path, farben: dict,
                 schrift_titel: str = SCHRIFT_TITEL,
                 schrift_text: str = SCHRIFT_TEXT) -> None:
    """Schreibt das neue Theme in die .pptx (ZIP wird neu gepackt).

    Es wird in eine TEMPORAERE Datei geschrieben und erst danach umbenannt –
    bricht der Vorgang ab, bleibt die alte (funktionierende) Vorlage stehen
    statt eines halben ZIPs."""
    tmp = pptx_pfad.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_pfad) as quelle:
        namen = quelle.namelist()
        ziel_themes = [n for n in namen if re.fullmatch(r"ppt/theme/theme\d+\.xml", n)]
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as ziel:
            for eintrag in quelle.infolist():
                daten = quelle.read(eintrag.filename)
                if eintrag.filename in ziel_themes:
                    try:
                        daten = _theme_xml(daten, farben, schrift_titel, schrift_text)
                    except Exception:  # noqa: BLE001
                        pass       # lieber unveraendertes Theme als kaputte Datei
                ziel.writestr(eintrag, daten)
    tmp.replace(pptx_pfad)


def _auf_breitbild_skalieren(prs, faktor: float) -> None:
    """Zieht Platzhalter in Master und Layouts auf die neue Folienbreite.

    DAS IST KEIN FEINSCHLIFF, SONDERN PFLICHT: ``prs.slide_width`` zu setzen
    aendert nur die Folienabmessung – die Platzhalter behalten ihre absoluten
    Positionen aus dem 4:3-Ausgangstemplate. Ohne diese Skalierung endet jede
    Aufzaehlung bei etwa zwei Dritteln der Folienbreite und rechts bleibt ein
    breiter, leerer Streifen. Genau so sah die erste Fassung im PDF-Test aus.

    Skaliert wird nur WAAGERECHT (links + Breite): die Folienhoehe bleibt bei
    6858000 EMU unveraendert, senkrecht passt also alles schon.

    ANGEFASST WERDEN NUR FORMEN MIT EIGENER GEOMETRIE (``spPr/a:xfrm``). Ein
    Layout-Platzhalter ohne eigenes ``xfrm`` ERBT Position und Groesse vom
    gleichnamigen Master-Platzhalter. Wer ihn trotzdem setzt, bekommt zwei
    Fehler auf einmal – gemessen an der ersten Fassung:
      * die Breite wird ZWEIMAL skaliert (Master 8229600 → 10972525, danach
        liest das Layout diesen geerbten Wert und macht 14630400 daraus – mehr
        als die Folie breit ist), und
      * ``top`` fällt auf 0, weil python-pptx beim Anlegen des neuen ``xfrm``
        nur die gesetzte Achse kennt. Im PDF standen die Titel dadurch am
        oberen Folienrand und wurden angeschnitten.
    Geerbte Platzhalter folgen dem Master von allein – man muss sie also gar
    nicht anfassen."""
    def eigene_geometrie(shape) -> bool:
        try:
            return shape._element.spPr.xfrm is not None
        except Exception:  # noqa: BLE001
            return False

    def anpassen(formen):
        for shape in formen:
            if not eigene_geometrie(shape):
                continue
            try:
                if shape.left is None or shape.width is None:
                    continue
                shape.left = int(shape.left * faktor)
                shape.width = int(shape.width * faktor)
            except Exception:  # noqa: BLE001
                continue

    for master in prs.slide_masters:
        anpassen(master.shapes)
        for layout in master.slide_layouts:
            anpassen(layout.shapes)


def _ph_stil(ph, sz=None, farbe=None, fett=None, algn=None, anker=None) -> None:
    """Legt Schriftgroesse/-farbe/Ausrichtung EINES Layout-Platzhalters fest.

    Warum ueber ``lstStyle`` im Layout und nicht ueber ``paragraphs[0].font``:
    ein Layout-Platzhalter hat in der Regel gar keinen Absatz, und ein Absatz
    auf der FOLIE erbt seine Eigenschaften aus der Listenformatierung des
    Layouts – nicht aus einem dort zufaellig vorhandenen Absatz. Genau daran
    ist eine frühere Fassung gescheitert (die Titel blieben zentriert).

    Das gilt nur fuer Werte, die zur STELLE gehoeren (Titelfolie ist groesser
    als eine Inhaltsfolie). Alles Durchgaengige steht in den Master-Textstilen
    (siehe ``_typografie``)."""
    from lxml import etree

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        txBody = ph.text_frame._txBody
        if anker is not None:
            bodyPr = txBody.find(f"{{{A}}}bodyPr")
            if bodyPr is not None:
                bodyPr.set("anchor", anker)
        lst = txBody.find(f"{{{A}}}lstStyle")
        if lst is None:
            lst = etree.SubElement(txBody, f"{{{A}}}lstStyle")
            # lstStyle muss VOR den Absaetzen stehen (Schema-Reihenfolge:
            # bodyPr, lstStyle, p…) – sonst lehnt PowerPoint die Datei ab.
            txBody.remove(lst)
            txBody.insert(1, lst)
        lvl1 = lst.find(f"{{{A}}}lvl1pPr")
        if lvl1 is None:
            lvl1 = etree.SubElement(lst, f"{{{A}}}lvl1pPr")
        if algn:
            # ALLE vorhandenen Ebenen, nicht nur die erste: der Untertitel der
            # Office-Titelfolie bringt einen eigenen lstStyle mit, in dem lvl1
            # BIS lvl9 auf 'ctr' stehen. Nur lvl1 zu setzen liesse acht
            # zentrierte Ebenen zurueck – sie schlagen zu, sobald jemand im
            # Kicker eine Unterebene benutzt, und der Testlauf hat genau das
            # gefunden.
            lvl1.set("algn", algn)
            for el in lst:
                if re.fullmatch(r"lvl\dpPr", etree.QName(el).localname or ""):
                    el.set("algn", algn)
        if sz is None and farbe is None and fett is None:
            return
        d = lvl1.find(f"{{{A}}}defRPr")
        if d is None:
            d = etree.SubElement(lvl1, f"{{{A}}}defRPr")
        if sz is not None:
            d.set("sz", str(int(sz)))
        if fett is not None:
            d.set("b", "1" if fett else "0")
        if farbe is not None:
            for k in list(d.findall(f"{{{A}}}solidFill")):
                d.remove(k)
            fill = etree.Element(f"{{{A}}}solidFill")
            etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", farbe)
            d.insert(0, fill)   # solidFill steht vor latin/ea/cs
    except Exception:  # noqa: BLE001
        pass


def _typografie(prs, farben: dict) -> None:
    """Schreibt Titel- und Textstile der Firmenvorlage in den Master.

    HIER LIEGT DER GRÖSSTE TEIL DES DESIGNS. Die Office-Vorgabe (Titel 44 pt
    zentriert, Text 28/24/20 pt mit runden Aufzaehlungspunkten) ist fuer 4:3
    gemacht; auf 16:9 wirkt sie grob und fuellt eine Folie mit vier Stichworten.
    Die Firmenvorlage setzt Titel 32 pt fett linksbuendig, Grundtext 18 pt und
    verzichtet auf der ersten Ebene ganz auf ein Aufzaehlungszeichen.

    ZWEI BEWUSSTE ABWEICHUNGEN vom Original, beide zugunsten der Lesbarkeit
    einer maschinell erzeugten Folie:
      * Das Original setzt Ebene 2 auf 14 pt OHNE Zeichen und Einrueckung und
        gibt erst ab Ebene 3 das ``+`` – eine Unterebene waere damit von der
        Hauptebene nicht zu unterscheiden. Der Agent nutzt Ebenen aber
        regelmaessig (``> Unterpunkt``). Deshalb: absteigende Groessenreihe und
        das ``+`` schon ab Ebene 2, mit demselben Einrueckschritt (271463 EMU).
      * Das Original laesst Ebene 3–5 bei 18 pt, also GROESSER als Ebene 2.

    Die Schrift wird hier ausdruecklich gesetzt (nicht ueber ``+mn-lt``
    geerbt), weil auch das Original in den Textstilen den Roman-Schnitt nennt
    und den Titel ueber ``b="1"`` fett schaltet. Wuerde der Titel die
    Bold-Variante des Themes erben UND ``b="1"`` tragen, rechnete PowerPoint
    einen zweiten, kuenstlichen Fettschnitt darauf."""
    from lxml import etree

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    text = farben.get("dunkel", DUNKEL)

    def schriftzeilen() -> str:
        return "".join(f'<a:{t} typeface="{SCHRIFT_TEXT}"/>' for t in ("latin", "ea", "cs"))

    def absatz(tag: str, groesse: int, marL: int, fett: bool, bullet: bool,
               vorher: int = 0) -> str:
        # Kind-Reihenfolge nach CT_TextParagraphProperties: lnSpc, spcBef,
        # spcAft, bu*, defRPr. Ein vertauschtes Kind macht die Datei fuer
        # PowerPoint ungueltig – deshalb wird der Absatz komplett neu gebaut
        # statt vorhandene Elemente zu ergaenzen.
        bu = (f'<a:buFont typeface="{SCHRIFT_TEXT}"/><a:buChar char="+"/>'
              if bullet else "<a:buNone/>")
        einzug = f' indent="{-BULLET_EINZUG}"' if bullet else ' indent="0"'
        return (
            f'<a:{tag} marL="{marL}"{einzug} algn="l" defTabSz="914400" rtl="0" '
            f'eaLnBrk="1" latinLnBrk="0" hangingPunct="1">'
            f'<a:lnSpc><a:spcPct val="100000"/></a:lnSpc>'
            f'<a:spcBef><a:spcPts val="{vorher}"/></a:spcBef>'
            f'<a:spcAft><a:spcPts val="{ABSTAND_NACH}"/></a:spcAft>'
            f'{bu}'
            f'<a:defRPr sz="{groesse}" b="{"1" if fett else "0"}" i="0" kern="1200">'
            f'<a:solidFill><a:srgbClr val="{text}"/></a:solidFill>'
            f'{schriftzeilen()}'
            f'</a:defRPr>'
            f'</a:{tag}>'
        )

    # Ebene 1 bekommt zusaetzlich Abstand DAVOR. Das Original setzt ihn auf 0 –
    # dort steht auf einer Folie meist ein durchgehender Text. Eine erzeugte
    # Aufzaehlung wechselt dagegen staendig zwischen Haupt- und Unterpunkten,
    # und ohne diesen Abstand beginnt der naechste Hauptpunkt unmittelbar unter
    # dem letzten Unterpunkt – die Gliederung ist dann nicht mehr ablesbar
    # (im PDF-Test genau so gesehen).
    koerper = "".join(
        absatz(f"lvl{i + 1}pPr", groesse,
               0 if i == 0 else BULLET_EINZUG * i,
               False, i > 0, ABSTAND_VOR if i == 0 else 0)
        for i, groesse in enumerate(TYPO_BODY)
    )
    xml = (
        f'<p:txStyles xmlns:p="{P}" xmlns:a="{A}">'
        f'<p:titleStyle>{absatz("lvl1pPr", TYPO_TITEL, 0, True, False)}</p:titleStyle>'
        f'<p:bodyStyle>{koerper}</p:bodyStyle>'
        f'<p:otherStyle><a:lvl1pPr><a:defRPr sz="{TYPO_BODY[0]}"/></a:lvl1pPr></p:otherStyle>'
        f'</p:txStyles>'
    )

    for master in prs.slide_masters:
        try:
            alt = master.element.find(f"{{{P}}}txStyles")
            if alt is not None:
                master.element.remove(alt)
            # txStyles ist das LETZTE Kind von <p:sldMaster> (nach cSld,
            # clrMap, sldLayoutIdLst, transition, timing, hf) – anhaengen ist
            # damit schema-konform.
            master.element.append(etree.fromstring(xml))
        except Exception:  # noqa: BLE001
            continue


def _koerper_und_titel(layout):
    """Zerlegt ein Layout in (Titel, Textkoerper, hat_bild).

    Der Bild-Hinweis entscheidet, ob das Raster ueberhaupt greifen darf: in
    'Bild mit Beschriftung' laege der Text sonst auf dem Bildplatzhalter."""
    from pptx.enum.shapes import PP_PLACEHOLDER as PH

    titel = None
    koerper = []
    hat_bild = False
    for ph in layout.placeholders:
        try:
            typ = ph.placeholder_format.type
        except Exception:  # noqa: BLE001
            continue
        if typ in (PH.TITLE, PH.CENTER_TITLE, PH.VERTICAL_TITLE):
            if titel is None:
                titel = ph
        elif typ in (PH.PICTURE, PH.CHART, PH.TABLE, PH.MEDIA_CLIP):
            hat_bild = True
        elif typ in (PH.BODY, PH.OBJECT, PH.SUBTITLE):
            koerper.append(ph)
    koerper.sort(key=lambda p: ((p.top or 0), (p.left or 0)))
    return titel, koerper, hat_bild


def _setze(shape, links: int, oben: int, breite: int, hoehe: int) -> None:
    """Setzt Position und Groesse – IMMER alle vier Werte.

    Fehlt eine Achse, legt python-pptx ein ``a:xfrm`` mit nur der gesetzten
    Groesse an und die andere faellt auf 0 zurueck; im PDF klebten die Titel
    dadurch am oberen Folienrand (siehe ``_auf_breitbild_skalieren``)."""
    try:
        shape.left, shape.top, shape.width, shape.height = links, oben, breite, hoehe
    except Exception:  # noqa: BLE001
        pass


def _raster(prs, farben: dict) -> None:
    """Legt den Satzspiegel der Firmenvorlage ueber die Layouts.

    Ohne diesen Schritt sitzen die Platzhalter weiter an den Office-Positionen
    (Titel mittig auf halber Hoehe, Rand 0,64 cm) – das Deck traegt dann zwar
    die richtigen Farben, ist aber sofort als Fremdkoerper zu erkennen. Der
    praegende Wert ist der linke Rand: Titel, Unterzeile und Inhalt beginnen
    auf DERSELBEN Kante.

    Angefasst werden nur die namentlich bekannten Layouts. Alles andere
    (Bildlayouts, senkrechter Text) bleibt bei der reinen Breitbild-Skalierung
    – dort wuerde ein pauschales Raster Text auf Bildplatzhalter legen."""
    rechts = BREITE_16_9 - RAND_LINKS - INHALT_B      # symmetrischer Rand
    spalte = (INHALT_B - SPALTEN_LUFT) // 2

    for layout in prs.slide_layouts:
        name = (layout.name or "").strip().lower()
        titel, koerper, hat_bild = _koerper_und_titel(layout)

        if name in ("titelfolie", "title slide"):
            # KICKER OBEN, Titel gross darunter – das ist der auffaelligste
            # Unterschied zum Office-Standard (dort steht der Untertitel unter
            # dem Titel und beides ist zentriert).
            if titel is not None:
                _setze(titel, RAND_LINKS, T_TITEL_Y, INHALT_B, T_TITEL_H)
                _ph_stil(titel, sz=TYPO_TITELFOLIE, fett=True, algn="l", anker="t")
            if koerper:
                _setze(koerper[0], RAND_LINKS, T_KICKER_Y, INHALT_B, T_KICKER_H)
                _ph_stil(koerper[0], sz=TYPO_KICKER, farbe=farben.get("akzent", STANDARD_AKZENT),
                         fett=False, algn="l", anker="b")
            for rest in koerper[1:]:
                _setze(rest, RAND_LINKS, HOEHE_16_9 - 900000, INHALT_B, SUB_H)
            # Akzentstrich unter dem Titel. Die Firmenvorlage fuellt die untere
            # Folienhaelfte mit einem Vollbild; ohne Bildmaterial bliebe hier
            # eine leere Flaeche und die Titelfolie saehe unfertig aus (im
            # PDF-Test genau so). Der Strich nimmt die Kante des
            # Kapitelkastens als Motiv auf und bleibt zurueckhaltend genug,
            # dass ein spaeter eingefuegtes Bild ihn nicht stoert.
            _rechteck(layout, RAND_LINKS, T_TITEL_Y + T_TITEL_H + 220000,
                      TITELSTRICH_B, TITELSTRICH_H,
                      farben.get("akzent", STANDARD_AKZENT), "Akzentstrich", 902)
            continue

        if name in ("abschnitt", "section header"):
            _kapitelkasten(layout, titel, koerper, farben)
            continue

        if titel is not None:
            _setze(titel, RAND_LINKS, TITEL_Y, INHALT_B, TITEL_H)

        if hat_bild or not koerper:
            continue
        if name in ("titel und inhalt", "title and content"):
            _setze(koerper[0], RAND_LINKS, INHALT_Y, INHALT_B, INHALT_H)
        elif name in ("zwei inhalte", "two content") and len(koerper) >= 2:
            _setze(koerper[0], RAND_LINKS, INHALT_Y, spalte, INHALT_H)
            _setze(koerper[1], BREITE_16_9 - rechts - spalte, INHALT_Y, spalte, INHALT_H)
        elif name in ("vergleich", "comparison") and len(koerper) >= 4:
            # Sortiert nach (oben, links): zwei Spaltenueberschriften, darunter
            # die beiden Inhalte.
            kopf_h = SUB_H
            luft = 120000
            for i, ph in enumerate(koerper[:4]):
                x = RAND_LINKS if i % 2 == 0 else BREITE_16_9 - rechts - spalte
                if i < 2:
                    _setze(ph, x, INHALT_Y, spalte, kopf_h)
                    _ph_stil(ph, sz=TYPO_SUB, fett=True,
                             farbe=farben.get("akzent", STANDARD_AKZENT))
                else:
                    _setze(ph, x, INHALT_Y + kopf_h + luft, spalte,
                           INHALT_H - kopf_h - luft)


def _rechteck(layout, x: int, y: int, breite: int, hoehe: int,
              farbe: str, name: str, kennung: int) -> bool:
    """Haengt ein farbiges Rechteck in den spTree eines LAYOUTS.

    ``LayoutShapes`` hat – wie ``MasterShapes`` – in python-pptx (1.0.2) kein
    ``add_shape``; die Methode gibt es nur auf Folien. Der oft genutzte Umweg
    „Form auf einer Wegwerf-Folie erzeugen und das Element verschieben" braucht
    danach das Loeschen dieser Folie (dafuer hat python-pptx keine API). Fuer
    ein Rechteck ohne Verknuepfung ist direktes XML der kuerzere Weg.

    EINGEFUEGT WIRD AN INDEX 2 – hinter ``nvGrpSpPr`` und ``grpSpPr``, also VOR
    allen Platzhaltern. Angehaengt laege die Flaeche ueber dem Text und deckte
    ihn zu."""
    from lxml import etree

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    xml = f"""<p:sp xmlns:p="{P}" xmlns:a="{A}">
  <p:nvSpPr>
    <p:cNvPr id="{kennung}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr userDrawn="1"/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{breite}" cy="{hoehe}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{farbe}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    try:
        layout.shapes._spTree.insert(2, etree.fromstring(xml))
        return True
    except Exception:  # noqa: BLE001
        return False


def _fusszeile(prs) -> None:
    """Setzt Datum, Fusszeile und Foliennummer auf eine gemeinsame Grundlinie.

    Die Firmenvorlage kennt unten nur EINEN Eintrag: die Foliennummer rechts
    aussen (11153375 EMU). Das Office-Template verteilt stattdessen drei
    Kaesten ueber die halbe Folienbreite – bei 16:9 enden sie mitten im Satz.

    Praktisch sichtbar wird das selten (PowerPoint blendet Fusszeilen erst nach
    dem Einschalten ein, und ``_leere_platzhalter_entfernen`` in main.py raeumt
    sie auf der Folie ohnehin weg). Es kostet aber nichts, sie richtig zu
    setzen – wer die Fusszeile einschaltet, bekommt sonst ein zerfallenes
    Fussfeld und sucht den Fehler in seinem Deck."""
    from pptx.enum.shapes import PP_PLACEHOLDER as PH

    breit = INHALT_B - NUM_B - 200000
    bereiche = [(prs.slide_masters[0].placeholders,)]
    bereiche += [(l.placeholders,) for l in prs.slide_layouts]
    for (phs,) in bereiche:
        for ph in phs:
            try:
                typ = ph.placeholder_format.type
            except Exception:  # noqa: BLE001
                continue
            if typ == PH.SLIDE_NUMBER:
                _setze(ph, NUM_X, NUM_Y, NUM_B, NUM_H)
                _ph_stil(ph, algn="r")
            elif typ == PH.FOOTER:
                _setze(ph, RAND_LINKS, NUM_Y, breit, NUM_H)
                _ph_stil(ph, algn="l")
            elif typ == PH.DATE:
                # Auf dieselbe Grundlinie, aber rechts neben der Fusszeile –
                # NICHT darunter: zwei Zeilen unter dem Satzspiegel wirken wie
                # ein Umbruchfehler.
                _setze(ph, RAND_LINKS + breit - 1800000, NUM_Y, 1800000, NUM_H)
                _ph_stil(ph, algn="r")


def _kapitelkasten(layout, titel, koerper, farben: dict) -> None:
    """Abschnittsfolie: farbiger Kasten unten links, Titel weiss darin.

    DAS IST DAS WIEDERERKENNUNGSMERKMAL der Firmenvorlage – dort liegt der
    Kasten (5553064 × 1865327 EMU) auf einem Vollbild. Ohne Bildmaterial steht
    er hier auf weissem Grund; Position, Groesse und Farbe sind uebernommen.

    TITEL UND ZUSATZTEXT LIEGEN BEIDE IM KASTEN. Die erste Fassung setzte den
    Zusatztext darunter auf die weisse Flaeche – im PDF-Test lief er in die
    Kastenkante (eine Zeile bei 24 pt passt nicht in die 1,1 cm, die unter dem
    Kasten bis zum Folienrand bleiben) und stand halb auf Rot, halb auf Weiss.
    Innerhalb des Kastens ist ausserdem der Bezug eindeutig: die Unterzeile
    gehoert zur Kapitelueberschrift, nicht zur folgenden Folie."""
    akzent = farben.get("akzent", STANDARD_AKZENT)
    if not _rechteck(layout, KAP_X, KAP_Y, KAP_B, KAP_H, akzent, "Kapitelkasten", 901):
        return

    # Der Kasten wird geteilt: oben der Titel (unten buendig), darunter die
    # Unterzeile (oben buendig). So sitzen beide an der gemeinsamen Mittellinie
    # und der Abstand bleibt gleich, egal wie lang der Titel ist.
    oben_h = int(KAP_H * 0.58)
    if titel is not None:
        _setze(titel, KAP_X, KAP_Y, KAP_B, oben_h)
        _ph_stil(titel, sz=TYPO_KAPITEL, farbe=HELL, fett=True, algn="l", anker="b")
        # ``unten`` ist NICHT 0: sonst klebt die Unterzeile am Titel, weil der
        # Titel unten und die Unterzeile oben buendig sitzt und zwischen beiden
        # nur die Trennlinie der Boxen liegt.
        _innenabstand(titel, links=KAP_INNEN, rechts=KAP_INNEN,
                      oben=KAP_INNEN, unten=70000)
    if koerper:
        _setze(koerper[0], KAP_X, KAP_Y + oben_h, KAP_B, KAP_H - oben_h)
        _ph_stil(koerper[0], sz=TYPO_BODY[0], farbe=HELL, algn="l", anker="t")
        _innenabstand(koerper[0], links=KAP_INNEN, rechts=KAP_INNEN,
                      oben=0, unten=KAP_INNEN)
    # Weitere Platzhalter (kommen im Office-Layout nicht vor) wandern aus dem
    # Kasten heraus, statt ihn zu ueberdecken.
    for rest in koerper[1:]:
        _setze(rest, RAND_LINKS, TITEL_Y, INHALT_B, TITEL_H)


def _innenabstand(ph, links=None, rechts=None, oben=None, unten=None) -> None:
    """Setzt die Textraender EINES Platzhalters (``bodyPr``-Attribute).

    Damit haelt der Text im Kapitelkasten Abstand zur Farbkante, ohne dass die
    Box kleiner sein muss als der Kasten – sonst waeren beim Verschieben des
    Kastens zwei Werte zu pflegen."""
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        bodyPr = ph.text_frame._txBody.find(f"{{{A}}}bodyPr")
        if bodyPr is None:
            return
        for name, wert in (("lIns", links), ("rIns", rechts),
                           ("tIns", oben), ("bIns", unten)):
            if wert is not None:
                bodyPr.set(name, str(int(wert)))
    except Exception:  # noqa: BLE001
        pass


def _logo_pfad():
    """Pfad zum Branding-Logo (helle Variante bevorzugt) oder None.

    Fuer eine helle Folie ist die HELL-Variante gemeint: das ist die fuer
    hellen Untergrund gemachte Fassung. Fehlt sie, wird die Dunkel-Variante
    genommen – die ist auf Weiss oft unsichtbar, deshalb nur als Rueckfall und
    nur, wenn ueberhaupt eine existiert."""
    try:
        from backend.main import _branding_logo_path
    except Exception:  # noqa: BLE001
        return None
    for variante in ("light", "dark"):
        try:
            p = _branding_logo_path(variante, "compact")
        except Exception:  # noqa: BLE001
            p = None
        if p and Path(p).exists():
            return Path(p)
    return None


def erzeuge(ziel: Path = None) -> Path:
    """Erzeugt die neutrale Vorlage und gibt ihren Pfad zurueck."""
    from pptx import Presentation
    from pptx.util import Emu

    ziel = Path(ziel) if ziel else (VORLAGEN_DIR / STANDARD_NAME)
    ziel.parent.mkdir(parents=True, exist_ok=True)
    farben = branding_farben()

    prs = Presentation()
    # 16:9 – der eigentliche Grund, warum die alten Decks „alt" aussahen.
    alt_breite = prs.slide_width
    prs.slide_width = Emu(BREITE_16_9)
    prs.slide_height = Emu(HOEHE_16_9)
    # …und die Platzhalter mitziehen, sonst bleibt rechts ein Drittel leer.
    if alt_breite and alt_breite != BREITE_16_9:
        _auf_breitbild_skalieren(prs, BREITE_16_9 / float(alt_breite))

    # Layouts deutsch benennen (nur die bekannten; unbekannte bleiben).
    for layout in prs.slide_layouts:
        neu = LAYOUT_NAMEN.get(layout.name)
        if neu:
            layout.name = neu

    # Reihenfolge: erst die Textstile (gelten fuer alles), dann das Raster –
    # das Raster setzt an einzelnen Stellen abweichende Groessen (Titelfolie,
    # Kapitelkasten) und muss deshalb zuletzt kommen. Beides NACH dem
    # Skalieren, denn beide arbeiten bereits in 16:9-Koordinaten und duerften
    # nicht noch einmal gestreckt werden.
    _typografie(prs, farben)
    _raster(prs, farben)
    _fusszeile(prs)
    # Das LOGO kommt NICHT in den Master, sondern beim Erzeugen auf die
    # Titelfolie (main.py). Zwei Gruende: MasterShapes kann keine Bilder
    # aufnehmen (die Medien-Beziehung haengt an der Folie, nicht am Master),
    # und ein Logo auf JEDER Folie ist im Corporate-Design die Ausnahme –
    # ueblich ist gross auf dem Titel, dezent oder gar nicht danach.
    prs.save(str(ziel))
    # Theme (Farben + Schrift) erst NACH dem Speichern – python-pptx schreibt
    # die Datei komplett neu und wuerde einen vorher gepatchten Theme-Part
    # wieder mit dem Original ueberschreiben.
    _patch_theme(ziel, farben)
    return ziel


def sicherstellen(pfad: Path = None) -> Path:
    """Gibt die Standardvorlage zurueck und erzeugt sie beim ersten Mal.

    Bewusst KEINE Neuerzeugung bei jedem Aufruf: die Datei darf von Hand
    ausgetauscht werden (echte Firmenvorlage). Wer die Branding-Farben neu
    einziehen will, loescht sie oder ruft ``erzeuge()`` auf."""
    ziel = Path(pfad) if pfad else (VORLAGEN_DIR / STANDARD_NAME)
    if ziel.exists() and ziel.stat().st_size > 0:
        return ziel
    return erzeuge(ziel)


def verfuegbare() -> list:
    """Alle hinterlegten Vorlagen (Dateinamen), Standardvorlage zuerst."""
    try:
        dateien = sorted(p.name for p in VORLAGEN_DIR.glob("*.pptx")
                         if not p.name.startswith(".") and not p.name.endswith(".tmp.pptx"))
    except Exception:  # noqa: BLE001
        return []
    if STANDARD_NAME in dateien:
        dateien.remove(STANDARD_NAME)
        dateien.insert(0, STANDARD_NAME)
    return dateien


def loese_vorlage(name: str = "") -> Path:
    """Loest einen Vorlagennamen auf; leer = Standardvorlage.

    Pfadanteile werden VERWORFEN (nur der Dateiname zaehlt): der Name kommt aus
    einem LLM-Aufruf, und '../../data/settings.json' darf hier nicht landen.
    Eine unbekannte Vorlage faellt auf die Standardvorlage zurueck, statt den
    Lauf scheitern zu lassen – ein Deck im Hausdesign ist besser als kein Deck.
    """
    if not name:
        return sicherstellen()
    sauber = Path(str(name)).name
    if not sauber.lower().endswith(".pptx"):
        sauber += ".pptx"
    kandidat = VORLAGEN_DIR / sauber
    if kandidat.exists():
        return kandidat
    return sicherstellen()
