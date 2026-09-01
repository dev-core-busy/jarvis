#!/usr/bin/env python3
"""Live-Probe auf DEV: baut das Modell ein Schaubild-Deck jetzt auf der Hausvorlage?

DER FALL (ECHT, 2026-09-01, Lauf 17882583414110010): "Baue auf einer Folie ein
echtes Schaubild mit Kaesten und Verbindungspfeilen" – das Modell griff
regelkonform zu python-pptx (Punkt 16 verlangt das fuer Formen) und oeffnete
`Presentation()` OHNE Argument. Ergebnis: accent1 4F81BD, Calibri, kein
Branding.

Gemessen wird hier NICHT der Wortlaut der Antwort, sondern das THEME der
erzeugten Datei – das ist die Eigenschaft, um die es geht. Ohne
Positivkontrolle waere ein "keine Datei gefunden" von "Erfolg" nicht zu
unterscheiden, deshalb bricht die Probe dann mit Exit 2 ab.

Aufruf auf DEV:  runuser -u jarvis -- ./venv/bin/python tests/live_pptx_hausvorlage_dev.py
"""

import asyncio
import re
import sys
import time
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

AUFGABE = ("Erstelle eine PowerPoint-Präsentation mit 4 Folien über die Vorteile von "
           "Prozessautomatisierung. Baue auf einer Folie ein echtes Schaubild mit Kästen "
           "und Verbindungspfeilen (kein Aufzählungstext) für den Ablauf: "
           "Auslöser → Prüfung → Verarbeitung → Benachrichtigung.")

_ok = _fail = 0


def pruefe(bedingung, text, detail=""):
    global _ok, _fail
    if bedingung:
        _ok += 1
        print(f"  ✓ {text}")
    else:
        _fail += 1
        print(f"  ✗ {text}" + (f" – {detail}" if detail else ""))


def accent1(p: Path) -> str:
    x = zipfile.ZipFile(str(p)).read("ppt/theme/theme1.xml").decode("utf-8", "replace")
    t = re.search(r"<a:accent1>.*?</a:accent1>", x, re.S).group(0)
    return re.search(r'val="([0-9A-Fa-f]{6})"', t).group(1).upper()


def hauptschrift(p: Path) -> str:
    x = zipfile.ZipFile(str(p)).read("ppt/theme/theme1.xml").decode("utf-8", "replace")
    f = re.search(r"<a:fontScheme.*?</a:fontScheme>", x, re.S).group(0)
    return re.findall(r'typeface="([^"]+)"', f)[0]


def _formen_zaehlen(p: Path):
    """(Formen, davon mit Theme-Akzent). Eine Form mit festem RGB-Wert folgt
    dem Branding NICHT – genau daran erkennt man ein handgebautes Schaubild."""
    try:
        from pptx import Presentation
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:  # noqa: BLE001
        return 0, 0
    ganz = akzent = 0
    prs = Presentation(str(p))
    for folie in prs.slides:
        for sh in folie.shapes:
            # AUF DEN SHAPE_TYPE FILTERN, nicht auf auto_shape_type: ein
            # Picture (hier das Logo der Titelfolie) wirft dort NICHT und
            # rutschte als achte "Form ohne Theme-Farbe" durch – ein Fehler
            # der Messung, nicht des Decks.
            if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            ganz += 1
            try:
                if sh.fill.fore_color.theme_color in (MSO_THEME_COLOR.ACCENT_1,
                                                      MSO_THEME_COLOR.ACCENT_2,
                                                      MSO_THEME_COLOR.ACCENT_5):
                    akzent += 1
            except Exception:  # noqa: BLE001
                pass
    return ganz, akzent


async def main():
    from backend.agent import JarvisAgent
    from skills.office import vorlage as VO

    soll = VO.branding_farben()["akzent"]
    print(f"Branding-Akzent dieser Installation: {soll}")

    start = time.time()
    vorher = {p: p.stat().st_mtime for p in Path("/tmp").rglob("*.pptx")}
    docs = ROOT / "data" / "documents"
    vorher |= {p: p.stat().st_mtime for p in docs.glob("*.pptx")}

    agent = JarvisAgent()
    antwort = await agent.run_task_headless(AUFGABE)
    dauer = time.time() - start
    print(f"\nLauf beendet nach {dauer:.1f}s, Antwort {len(antwort or '')} Zeichen\n")

    neu = [p for p in list(Path("/tmp").rglob("*.pptx")) + list(docs.glob("*.pptx"))
           if p not in vorher]
    if not neu:
        print("ABBRUCH: keine .pptx entstanden – ohne Datei ist nichts gemessen.")
        print((antwort or "")[:1500])
        sys.exit(2)

    for p in sorted(neu, key=lambda x: x.stat().st_mtime):
        a, s = accent1(p), hauptschrift(p)
        medien = len([n for n in zipfile.ZipFile(str(p)).namelist()
                      if n.startswith("ppt/media")])
        # WOHER kommt die Datei? Das Werkzeug legt sie mit Capability-Namen in
        # data/documents, ein selbstgebautes Skript nach /tmp. Der Ort ist damit
        # der Beleg fuer den GEWAEHLTEN WEG – und genau der war der Vorfall.
        weg = "office_create_powerpoint" if p.parent.name == "documents" else "python-pptx (Shell)"
        formen, akzent_formen = _formen_zaehlen(p)
        print(f"  DATEI {p.name}: accent1={a}, Schrift={s}, Medien={medien}, "
              f"Formen={formen} (davon Theme-Akzent {akzent_formen}), Weg={weg}")
        pruefe(a == soll, f"{p.name}: Theme traegt den Hausakzent {soll}", f"gemessen {a}")
        pruefe(a != "4F81BD", f"{p.name}: NICHT das python-pptx-Standardtheme")
        pruefe("Calibri" not in s, f"{p.name}: nicht Calibri", f"gemessen {s}")
        if formen:
            pruefe(akzent_formen == formen,
                   f"{p.name}: alle {formen} Formen nehmen die Theme-Farbe",
                   f"nur {akzent_formen}")

    print("\nAntwort-Auszug:", (antwort or "")[:400].replace("\n", " | "))
    print(f"\n{'=' * 62}\nErgebnis: {_ok}/{_ok + _fail} Pruefungen bestanden")
    sys.exit(1 if _fail else 0)


asyncio.run(main())
