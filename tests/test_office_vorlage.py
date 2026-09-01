#!/usr/bin/env python3
"""Tests fuer die PowerPoint-Hausvorlage (skills/office/vorlage.py + main.py).

DER AUSGANGSPUNKT: ``office_create_powerpoint`` rief ``Presentation()`` ohne
Argument auf – also das eingebaute Standarddesign von python-pptx: 4:3,
Calibri, Office-Blau, kein Bezug zum Branding. Jetzt wird eine Vorlage mit
echten Masterfolien benutzt.

DIE WICHTIGSTEN PRUEFUNGEN sind die beiden Regressionsfaelle, die erst der
PDF-Blick gezeigt hat:
 1. Platzhalter duerfen NICHT breiter als die Folie werden (die erste Fassung
    skalierte geerbte Layout-Platzhalter ein zweites Mal: 8229600 → 10972525 →
    14630400 bei 12192000 Folienbreite), und
 2. ``top`` darf nicht auf 0 fallen (beim Anlegen eines neuen ``xfrm`` kennt
    python-pptx nur die gesetzte Achse – im PDF klebten die Titel dadurch am
    oberen Rand).

Teil 1 laeuft ueberall (Quelltext), Teil 2 nur mit python-pptx (DEV im venv).
"""

import asyncio
import re
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

_ok = 0
_fail = 0


def pruefe(bedingung, text, detail=""):
    global _ok, _fail
    if bedingung:
        _ok += 1
        print(f"  ✓ {text}")
    else:
        _fail += 1
        print(f"  ✗ {text}" + (f" – {detail}" if detail else ""))


# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 1. Quelltext ===")

V = (ROOT / "skills" / "office" / "vorlage.py").read_text()
M = (ROOT / "skills" / "office" / "main.py").read_text()

pruefe("BREITE_16_9 = 12192000" in V, "Vorlage ist 16:9 (nicht das 4:3-Standardtemplate)")
pruefe("eigene_geometrie" in V,
       "Skalierung fasst nur Formen mit EIGENEM xfrm an (sonst doppelt + top=0)")
pruefe('SCHRIFT_TITEL = "HelveticaNeue LT 75 Bold"' in V
       and 'SCHRIFT_TEXT = "HelveticaNeue LT 55 Roman"' in V,
       "Schriften der Firmenvorlage (majorFont/minorFont getrennt)")
pruefe("PDF-EXPORT" in V,
       "die fehlende Serverschrift ist als Fallstrick dokumentiert")
pruefe("colors_light" in V and "accent" in V, "Akzentfarbe kommt aus dem Branding")
pruefe("_typografie" in V and "_raster" in V,
       "Designprofil wirkt ueber Master-Textstile UND Satzspiegel")
pruefe("VORLAGEN_DIR" in V and "vorlagen" in V and '"documents"' not in V,
       "Vorlagen liegen NICHT in data/documents (dort gilt Frist + Eigentuemer)")
pruefe("titleStyle" in V,
       "Titel-Ausrichtung wird im Master-titleStyle gesetzt (Layout-Absatz wirkt nicht)")
pruefe("txBody.insert(1, lst)" in V,
       "lstStyle wird an der schema-richtigen Stelle eingefuegt (bodyPr, lstStyle, p)")

pruefe("Presentation(str(_vorlage.loese_vorlage(template)))" in M,
       "das Werkzeug nutzt die Vorlage")
pruefe("_LAYOUT_ALIAS" in M and "Title and Content" in M,
       "Layouts werden ueber NAMEN gesucht, deutsch UND englisch")
pruefe("_LAYOUT_FALLBACK" in M, "Rueckfall-Index, wenn eine Fremdvorlage andere Namen nutzt")
pruefe("placeholder_format.type" in M,
       "Platzhalter werden ueber den TYP gewaehlt, nicht ueber placeholders[1]")
pruefe("PH.DATE" not in M and "erlaubt = {PH.BODY" in M,
       "Datum/Fusszeile/Foliennummer werden nie mit Inhalt befuellt")
pruefe("_leere_platzhalter_entfernen" in M,
       "leere Platzhalter werden entfernt (sonst 'Klicken Sie, um Text hinzuzufuegen')")
pruefe("TemplateInfoTool()" in M, "office_template_info ist registriert")
i_desc = M.find("KEINE Farb-, Schrift- oder Groessenangaben")
pruefe(i_desc > 0, "die Werkzeug-Beschreibung verbietet eigene Stilangaben ausdruecklich")
# Der Kern des Vorlagen-Wegs: es wird NUR Text gesetzt.
pruefe(".font.size" not in M.split("class CreatePowerPointTool")[1].split("class TemplateInfoTool")[0],
       "das Werkzeug setzt keine Schriftgroessen (die kommen aus der Vorlage)")

# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 2. Erzeugte Vorlage ===")

try:
    from pptx import Presentation
    from skills.office import vorlage as VO
    HAT_PPTX = True
except Exception as e:  # noqa: BLE001
    HAT_PPTX = False
    print(f"  … uebersprungen: {type(e).__name__} ({e}) – auf DEV im venv ausfuehren")

if HAT_PPTX:
    # _hex: Farbnormierung
    faelle = [("#9B59B6", "9B59B6"), ("9b59b6", "9B59B6"), ("#abc", "AABBCC"),
              ("", "FALLBACK"), (None, "FALLBACK"), ("rgb(1,2,3)", "FALLBACK"),
              ("#12345", "FALLBACK"), ("  #10B981  ", "10B981")]
    for roh, erwartet in faelle:
        got = VO._hex(roh, "FALLBACK")
        pruefe(got == erwartet, f"_hex({roh!r}) -> {got}", f"erwartet {erwartet}")

    import tempfile
    tmp = Path(tempfile.mkdtemp(prefix="jvorl_"))
    ziel = tmp / "probe.pptx"
    VO.erzeuge(ziel)
    pruefe(ziel.exists() and ziel.stat().st_size > 10000, "Vorlage wird erzeugt")

    prs = Presentation(str(ziel))
    pruefe(prs.slide_width == 12192000 and prs.slide_height == 6858000,
           f"Format 16:9 ({prs.slide_width}×{prs.slide_height})")
    namen = [l.name for l in prs.slide_layouts]
    pruefe("Titelfolie" in namen and "Titel und Inhalt" in namen,
           f"Layouts deutsch benannt ({namen[:3]})")
    ab_layout = [l for l in prs.slide_layouts if l.name == "Abschnitt"][0]
    pruefe("Kapitelkasten" in [s.name for s in ab_layout.shapes],
           "Abschnittsfolie hat den farbigen Kasten der Firmenvorlage")
    # Er muss VOR den Platzhaltern stehen, sonst deckt er den Titel zu – aber
    # HINTER einem etwaigen Hintergrundbild (das wird spaeter eingehaengt).
    ab_namen = [s.name for s in ab_layout.shapes]
    i_kasten = ab_namen.index("Kapitelkasten")
    i_titel = next((i for i, n in enumerate(ab_namen) if n.startswith("Title")), 99)
    pruefe(i_kasten < i_titel,
           f"Kapitelkasten liegt hinter dem Text ({ab_namen[:3]})")
    tl_layout = [l for l in prs.slide_layouts if l.name == "Titelfolie"][0]
    pruefe("Akzentstrich" in [s.name for s in tl_layout.shapes],
           "Titelfolie hat den Akzentstrich (Ersatz fuer das Vollbild)")

    # ── REGRESSION 1: nichts ist breiter als die Folie ──────────────────────
    zu_breit = []
    top_null = []
    for l in prs.slide_layouts:
        for sh in l.shapes:
            try:
                if sh.left is None or sh.width is None:
                    continue
                if sh.left + sh.width > prs.slide_width + 1000:
                    zu_breit.append(f"{l.name}/{sh.name}={sh.left + sh.width}")
                if sh.top == 0 and sh.height and sh.height < prs.slide_height:
                    top_null.append(f"{l.name}/{sh.name}")
            except Exception:  # noqa: BLE001
                continue
    pruefe(not zu_breit, "kein Platzhalter ragt aus der Folie (keine Doppel-Skalierung)",
           ", ".join(zu_breit[:3]))
    pruefe(not top_null, "kein Platzhalter auf top=0 gerutscht", ", ".join(top_null[:3]))

    # Volle Breite wird genutzt: rechter Rand = linker Rand (symmetrisch)
    inhalt = [l for l in prs.slide_layouts if l.name == "Titel und Inhalt"][0]
    titel = list(inhalt.placeholders)[0]
    links = titel.left
    rechts = prs.slide_width - (titel.left + titel.width)
    pruefe(abs(links - rechts) < 20000,
           f"Ränder symmetrisch (links {links}, rechts {rechts})")
    pruefe(titel.width > 10000000, f"Titel nutzt die Breitbild-Breite ({titel.width})")

    # ── Theme im ZIP ────────────────────────────────────────────────────────
    x = zipfile.ZipFile(str(ziel)).read("ppt/theme/theme1.xml").decode("utf-8")
    def slot(name):
        m = re.search(r"<a:" + name + r">.*?val=\"([0-9A-Fa-f]{6})\"", x, re.S)
        return (m.group(1).upper() if m else "")
    erwartet_akzent = VO.branding_farben()["akzent"]
    pruefe(slot("accent1") == erwartet_akzent,
           f"accent1 ist die Branding-/Standardfarbe ({slot('accent1')})")
    pruefe(slot("lt1") == "FFFFFF", f"heller Hintergrund ({slot('lt1')})")
    pruefe(slot("dk1") == VO.DUNKEL, f"dunkle Textfarbe ({slot('dk1')})")
    pruefe(slot("accent2") == VO.PALETTE[0],
           "Folgefarben in der Reihenfolge der Firmenvorlage (accent2..6)")
    pruefe(slot("hlink") == erwartet_akzent and slot("folHlink") == VO.FOL_HLINK,
           "Hyperlinkfarben wie in der Firmenvorlage")
    pruefe("windowText" not in x and "sysClr" not in x.split("</a:clrScheme>")[0],
           "sysClr (windowText/window) wurde ersetzt – sonst blieben Text/Hintergrund alt")
    pruefe(f'<a:latin typeface="{VO.SCHRIFT_TITEL}"' in x,
           "majorFont = Ueberschriftenschnitt der Firmenvorlage")
    pruefe(f'<a:latin typeface="{VO.SCHRIFT_TEXT}"' in x,
           "minorFont = Textschnitt der Firmenvorlage")
    pruefe("Calibri" not in x, "keine Calibri-Reste im Theme")

    # ── Titel-Ausrichtung ───────────────────────────────────────────────────
    mx = zipfile.ZipFile(str(ziel)).read("ppt/slideMasters/slideMaster1.xml").decode("utf-8")
    tstyle = mx.split("<p:titleStyle>")[1].split("</p:titleStyle>")[0] if "<p:titleStyle>" in mx else ""
    pruefe('algn="l"' in tstyle, "Master-Titelstil ist linksbuendig")
    pruefe(f'sz="{VO.TYPO_TITEL}"' in tstyle and 'b="1"' in tstyle,
           f"Titel {VO.TYPO_TITEL // 100} pt fett (nicht die Office-Vorgabe 44 pt)")
    bstyle = mx.split("<p:bodyStyle>")[1].split("</p:bodyStyle>")[0] if "<p:bodyStyle>" in mx else ""
    pruefe(f'sz="{VO.TYPO_BODY[0]}"' in bstyle and f'sz="{VO.TYPO_BODY[1]}"' in bstyle,
           "Textstufen der Firmenvorlage im Master")
    pruefe("<a:buNone/>" in bstyle.split("</a:lvl1pPr>")[0],
           "Ebene 1 ohne Aufzaehlungszeichen (wie die Firmenvorlage)")
    pruefe('char="+"' in bstyle, "Ebenen ab 2 mit dem '+' der Firmenvorlage")
    # Die Titelfolie ist im Hausdesign LINKSBUENDIG – anders als im
    # Office-Standard. Eine zentrierte Titelfolie waere hier ein Rueckfall.
    tl = [l for l in prs.slide_layouts if l.name == "Titelfolie"][0]
    pruefe('algn="ctr"' not in tl.element.xml, "Titelfolie ist linksbuendig")
    tl_titel, tl_koerper, _ = VO._koerper_und_titel(tl)
    pruefe(tl_koerper and tl_titel.top > tl_koerper[0].top,
           "Titelfolie: Kicker steht UEBER dem Titel")

    # ── Abschnitt: Titel oben ───────────────────────────────────────────────
    ab = [l for l in prs.slide_layouts if l.name == "Abschnitt"][0]
    phs = {}
    for ph in ab.placeholders:
        phs[str(ph.placeholder_format.type)] = ph.top
    t_top = next((v for k, v in phs.items() if "TITLE" in k), None)
    b_top = next((v for k, v in phs.items() if "BODY" in k or "OBJECT" in k), None)
    pruefe(t_top is not None and b_top is not None and t_top < b_top,
           f"Abschnittsfolie: Titel steht ueber dem Zusatztext ({t_top} < {b_top})")

    # ── Satzspiegel: alle Textkanten auf einer Linie ────────────────────────
    kanten = {}
    for lname in ("Titel und Inhalt", "Nur Titel", "Zwei Inhalte", "Vergleich"):
        lay = [l for l in prs.slide_layouts if l.name == lname][0]
        t, k, _ = VO._koerper_und_titel(lay)
        if t is not None:
            kanten.setdefault(lname, []).append(t.left)
        if k:
            kanten[lname].append(k[0].left)
    schief = {n: v for n, v in kanten.items() if len(set(v)) > 1}
    pruefe(not schief, "Titel und Inhalt beginnen auf derselben linken Kante", str(schief))
    pruefe(all(v[0] == VO.RAND_LINKS for v in kanten.values()),
           f"linke Kante ist der Satzspiegel der Firmenvorlage ({VO.RAND_LINKS})")

    zwei = [l for l in prs.slide_layouts if l.name == "Zwei Inhalte"][0]
    _, sp, _ = VO._koerper_und_titel(zwei)
    pruefe(len(sp) >= 2 and sp[0].width == sp[1].width, "zwei gleich breite Spalten")
    pruefe(len(sp) >= 2 and sp[1].left + sp[1].width == VO.RAND_LINKS + VO.INHALT_B,
           "rechte Spalte endet am Satzspiegel")

    # Der Kapitelkasten darf den Titel nicht ueberragen und der Zusatztext
    # nicht aus ihm herauslaufen – genau das war im ersten PDF-Test der Fall.
    ab2 = [l for l in prs.slide_layouts if l.name == "Abschnitt"][0]
    kasten = [s for s in ab2.shapes if s.name == "Kapitelkasten"][0]
    ab_t, ab_k, _ = VO._koerper_und_titel(ab2)
    drin = []
    for ph in ([ab_t] if ab_t is not None else []) + ab_k[:1]:
        drin.append(ph.top >= kasten.top
                    and ph.top + ph.height <= kasten.top + kasten.height + 1000
                    and ph.left >= kasten.left
                    and ph.left + ph.width <= kasten.left + kasten.width + 1000)
    pruefe(drin and all(drin), "Titel und Unterzeile liegen IM Kapitelkasten")

    # Foliennummer rechts unten (Firmenvorlage), nicht auf Office-Position
    nummer = None
    for ph in prs.slide_masters[0].placeholders:
        if "SLIDE_NUMBER" in str(ph.placeholder_format.type):
            nummer = ph
    pruefe(nummer is not None and nummer.left == VO.NUM_X and nummer.top == VO.NUM_Y,
           "Foliennummer sitzt rechts unten")

    # ── sicherstellen / loese_vorlage ───────────────────────────────────────
    vor = ziel.stat().st_mtime_ns
    gleich = VO.sicherstellen(ziel)
    pruefe(gleich == ziel and ziel.stat().st_mtime_ns == vor,
           "sicherstellen() erzeugt eine vorhandene Vorlage NICHT neu (Handaustausch bleibt)")
    echt = VO.VORLAGEN_DIR / VO.STANDARD_NAME
    pruefe(VO.loese_vorlage("../../data/settings.json").name.endswith(".pptx"),
           "Pfadanteile im Vorlagennamen werden verworfen")
    pruefe(VO.loese_vorlage("gibtsnicht") == echt or VO.loese_vorlage("gibtsnicht").exists(),
           "unbekannte Vorlage faellt auf die Standardvorlage zurueck")

    # ═════════════════════════════════════════════════════════════════════════
    print("\n=== 3. Layout-Auflösung und Befuellung ===")

    from skills.office.main import (_layout, _fuelle, _text_platzhalter,
                                    _leere_platzhalter_entfernen, _titel_setzen,
                                    CreatePowerPointTool, TemplateInfoTool)

    pruefe(_layout(prs, "inhalt").name == "Titel und Inhalt", "Kurzname 'inhalt' trifft")
    pruefe(_layout(prs, "abschnitt").name == "Abschnitt", "Kurzname 'abschnitt' trifft")
    pruefe(_layout(prs, "zwei").name == "Zwei Inhalte", "Kurzname 'zwei' trifft")
    pruefe(_layout(prs, "leer").name == "Leer", "Kurzname 'leer' trifft")
    pruefe(_layout(prs, "quatsch") is not None, "unbekannte Art liefert trotzdem ein Layout")

    # Englische Fremdvorlage: das Standardtemplate von python-pptx
    fremd = Presentation()
    pruefe(_layout(fremd, "inhalt").name == "Title and Content",
           "englische Layoutnamen werden ebenfalls getroffen")
    # Teiltreffer (Firmenvorlagen haengen gern etwas an)
    fremd.slide_layouts[1].name = "Titel und Inhalt (intern)"
    pruefe(_layout(fremd, "inhalt").name == "Titel und Inhalt (intern)",
           "Teiltreffer im Layoutnamen wird erkannt")

    # Befuellung inkl. Ebenen
    p2 = Presentation(str(ziel))
    s = p2.slides.add_slide(_layout(p2, "inhalt"))
    _titel_setzen(s, "Titel")
    felder = _text_platzhalter(s)
    pruefe(len(felder) == 1, f"genau EIN befuellbares Feld im Inhalts-Layout ({len(felder)})")
    _fuelle(felder[0], bullets=["Eins", "> Unter", ">> Tiefer", {"text": "Vier", "level": 1},
                                ("Fuenf", 2)])
    stufen = [p.level for p in felder[0].text_frame.paragraphs]
    texte = [p.text for p in felder[0].text_frame.paragraphs]
    pruefe(stufen == [0, 1, 2, 1, 2], f"Aufzaehlungsebenen: {stufen}")
    pruefe(texte == ["Eins", "Unter", "Tiefer", "Vier", "Fuenf"], f"Texte: {texte}")
    pruefe(all("font" not in p._pPr.xml if p._pPr is not None else True
               for p in felder[0].text_frame.paragraphs),
           "keine Schriftangaben in den Absaetzen (Vorlage bestimmt das Aussehen)")

    # Leere Platzhalter verschwinden
    s2 = p2.slides.add_slide(_layout(p2, "inhalt"))
    _titel_setzen(s2, "Nur Titel gesetzt")
    vorher = len(list(s2.placeholders))
    _leere_platzhalter_entfernen(s2)
    nachher = len(list(s2.placeholders))
    pruefe(nachher < vorher, f"leere Platzhalter entfernt ({vorher} -> {nachher})")
    pruefe(s2.shapes.title is not None and s2.shapes.title.text == "Nur Titel gesetzt",
           "der gefuellte Titel bleibt erhalten")

    # ── Ende zu Ende ────────────────────────────────────────────────────────
    print("\n=== 4. Ende zu Ende ===")
    r = asyncio.run(CreatePowerPointTool().execute(
        filename="test_vorlage_e2e",
        title="Titel", subtitle="Unter",
        slides=[
            {"title": "Erste", "bullets": ["A", "> a1"], "notes": "Notiz"},
            {"layout": "zwei", "title": "Zwei", "bullets": ["L1", "L2", "R1", "R2"]},
            {"layout": "abschnitt", "title": "Kapitel"},
        ]))
    pruefe("/api/documents/" in r, f"Download-Link geliefert ({r[:60]})")
    pruefe("Hinweis:" not in r, f"Vorlage war nutzbar (kein Rueckfall) – {r[-90:]}")
    m = re.search(r"/api/documents/([0-9a-f]{32}__[^)\s]+)", r)
    erzeugt = Path(ROOT / "data" / "documents" / m.group(1)) if m else None
    if erzeugt and erzeugt.exists():
        e = Presentation(str(erzeugt))
        pruefe(e.slide_width == 12192000, "erzeugte Datei ist 16:9")
        pruefe(len(e.slides) == 4, f"vier Folien ({len(e.slides)})")
        pruefe(e.slides[0].slide_layout.name == "Titelfolie", "erste Folie nutzt die Titelfolie")
        pruefe(e.slides[2].slide_layout.name == "Zwei Inhalte", "Zwei-Spalten-Layout genutzt")
        pruefe(e.slides[3].slide_layout.name == "Abschnitt", "Abschnitts-Layout genutzt")
        # Zwei-Spalten: beide Seiten befuellt
        spalten = _text_platzhalter(e.slides[2])
        gefuellt = [p for p in spalten if p.text_frame.text.strip()]
        pruefe(len(gefuellt) == 2, f"beide Spalten haben Inhalt ({len(gefuellt)})")
        pruefe(e.slides[1].has_notes_slide
               and "Notiz" in e.slides[1].notes_slide.notes_text_frame.text,
               "Sprechernotiz uebernommen")
        # Keine leeren Rahmen mehr
        leere = [sh.name for s3 in e.slides for sh in s3.placeholders
                 if sh.has_text_frame and not sh.text_frame.text.strip()]
        pruefe(not leere, f"keine leeren Platzhalter in der Datei ({leere[:3]})")
        try:
            erzeugt.unlink()
        except Exception:  # noqa: BLE001
            pass
    else:
        pruefe(False, "erzeugte Datei gefunden", str(erzeugt))

    info = asyncio.run(TemplateInfoTool().execute())
    pruefe("16:9" in info and "Titel und Inhalt" in info,
           "office_template_info nennt Format und Layouts")
    pruefe("abschnitt" in info and "nurtitel" in info, "…und die Kurznamen")

    import shutil
    shutil.rmtree(tmp, ignore_errors=True)

# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 5. Upload ueber den Branding-Reiter ===")

MAIN = (ROOT / "backend" / "main.py").read_text()
BR_JS = (ROOT / "frontend" / "js" / "branding.js").read_text()
SET_HTML = (ROOT / "frontend" / "settings.html").read_text()
I18N = (ROOT / "frontend" / "js" / "i18n.js").read_text()

for route, methode in [("/api/branding/pptx-templates", "get"),
                       ("/api/branding/pptx-template", "post"),
                       ("/api/branding/pptx-template", "delete"),
                       ("/api/branding/pptx-template/regenerate", "post")]:
    marke = f'@app.{methode}("{route}")'
    i = MAIN.find(marke)
    pruefe(i > 0, f"{methode.upper()} {route} existiert")
    if i > 0:
        # Alle vier sind Admin-Sache: sie legen Dateien ab, die JEDE spaetere
        # Praesentation bestimmen (Masterfolien = ausfuehrbares Aussehen).
        pruefe("require_local_auth" in MAIN[i:i + 700],
               f"{methode.upper()} {route} verlangt Admin (require_local_auth)")

i_up = MAIN.find('@app.post("/api/branding/pptx-template")')
fenster = MAIN[i_up:i_up + 2600]
pruefe("_pptx_tpl_pruefen" in fenster, "Upload prueft die Datei vor dem Ablegen")
pruefe("ppt/presentation.xml" in MAIN,
       "Pruefung schaut in den ZIP-Inhalt (nicht nur auf die Endung)")
pruefe("_PPTX_TPL_MAX_BYTES" in MAIN and "25 * 1024 * 1024" in MAIN, "Groessengrenze gesetzt")
pruefe('re.sub(r"[^A-Za-z0-9_\\-. ]+"' in fenster,
       "Dateiname wird entschaerft (kommt aus einem Datei-Dialog)")
pruefe("Path(file.filename" in fenster and ".name" in fenster,
       "nur der Basisname wird verwendet (keine Pfadanteile)")
pruefe('ziel.with_suffix(".upload.tmp")' in fenster and "tmp.replace(ziel)" in fenster,
       "erst danebenschreiben, dann umbenennen (keine halbe Vorlage bei Abbruch)")
pruefe("potx" in MAIN and 'f"{rein[:60]}.pptx"' in fenster,
       ".potx wird als .pptx abgelegt (der Skill sucht *.pptx)")
i_del = MAIN.find('@app.delete("/api/branding/pptx-template")')
pruefe("Path(str(name or \"\")).name" in MAIN[i_del:i_del + 900],
       "Loeschen akzeptiert keinen Pfad")
pruefe("status_code=404" in MAIN[i_del:i_del + 1600], "unbekannte Vorlage -> 404")
i_reg = MAIN.find('@app.post("/api/branding/pptx-template/regenerate")')
pruefe("v.erzeuge()" in MAIN[i_reg:i_reg + 800],
       "Neuerzeugen ruft erzeuge() (nicht sicherstellen – das wuerde nichts tun)")
pruefe("VORLAGEN_DIR" in MAIN and "data/branding" not in MAIN[i_up:i_up + 2000],
       "Vorlagen landen in data/vorlagen, nicht bei den Logos")

pruefe("uploadPptxTemplate" in BR_JS and "loadPptxTemplates" in BR_JS,
       "branding.js kann hochladen und listen")
pruefe("regeneratePptxTemplate" in BR_JS and "deletePptxTemplate" in BR_JS,
       "…neu erzeugen und entfernen")
pruefe("name.textContent = t.name" in BR_JS,
       "Vorlagenname wird per textContent gesetzt (Fremdinhalt, kein innerHTML)")
pruefe("this.loadPptxTemplates();" in BR_JS, "die Liste wird beim Oeffnen geladen")
# Nach der DEFINITION suchen, nicht nach dem ersten Vorkommen (das ist
# der Aufruf in renderPptxTemplates).
i_delfn = BR_JS.find("deletePptxTemplate: function")
pruefe(i_delfn > 0 and "window.confirm" in BR_JS[i_delfn:i_delfn + 700],
       "Entfernen fragt nach")
pruefe("br-pptx-file" in SET_HTML and "br-pptx-default" in SET_HTML
       and "br-pptx-regen" in SET_HTML and "br-pptx-list" in SET_HTML,
       "Markup im Branding-Reiter vorhanden")
pruefe(".pptx,.potx" in SET_HTML, "Datei-Dialog filtert auf Vorlagen")
for key in ("branding.pptx_heading", "branding.pptx_as_default", "branding.pptx_regen",
            "branding.pptx_none", "branding.pptx_badge_default"):
    pruefe(I18N.count(f"'{key}'") == 2, f"i18n {key} in DE und EN")

# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 6. Vorfall 2026-08-10: 2586 Folien, keine Grafik, kein Hintergrund ===")
# Ein Modell schickte 'slides' als TEXT. Die Schleife lief ueber die ZEICHEN,
# und weil ein String-Element als {"title": …} galt, wurde jedes Zeichen eine
# Folie: 2586 Stueck, 2 MB. Dazu fehlten Diagramme (das Werkzeug konnte keine)
# und jeder Hintergrund.

if HAT_PPTX:
    from skills.office.main import (_slides_normalisieren, _text_bereinigen,
                                    _zahl, MAX_FOLIEN)

    def _erzeugte_datei(antwort: str) -> Path:
        """Pfad der Datei hinter dem Download-Link."""
        m = re.search(r"/api/documents/([0-9a-f]{32}__[^)\s]+)", antwort)
        return ROOT / "data" / "documents" / m.group(1) if m else None

    # ── a) Der gemeldete Aufruf, woertlich aus dem Audit-Log von ECHT ────────
    ECHT = ("[{'title': 'Gesamtübersicht', 'layout': 'inhalt', 'bullets': "
            "['1.216 offene Tickets der letzten 14 Tage', 'Blocker (6), High (45)']}, "
            "{'title': 'Prioritäten', 'bullets': ['High 45']}]")
    folien, fehler = _slides_normalisieren(ECHT)
    pruefe(fehler is None and len(folien) == 2,
           f"Python-Schreibweise als Text -> 2 Folien (nicht 2586) [{len(folien)}]")
    pruefe(all(isinstance(f, dict) for f in folien), "…und es sind echte Folien-Objekte")
    pruefe(folien[0].get("title") == "Gesamtübersicht", "…Inhalt kommt korrekt an")

    f2, e2 = _slides_normalisieren('[{"title": "A"}, {"title": "B"}]')
    pruefe(e2 is None and len(f2) == 2, "gueltiges JSON als Text wird ebenfalls gelesen")

    # Kein Parser der Welt macht daraus eine Liste -> FEHLER, nie Zeichen-Folien
    for muell in ("das ist einfach nur text", "{kaputt", "[{'a': "):
        fm, em = _slides_normalisieren(muell)
        pruefe(em is not None and not fm,
               f"unlesbares 'slides' -> Fehlermeldung statt Zeichen-Folien ({muell[:12]}…)")
    fe, ee = _slides_normalisieren("")
    pruefe(ee is not None, "leerer Text -> Fehler")
    fn, en = _slides_normalisieren(None)
    pruefe(en is not None, "fehlendes 'slides' -> Fehler")

    # Was weiterhin funktionieren MUSS
    fd, _ = _slides_normalisieren({"title": "Einzelfolie"})
    pruefe(len(fd) == 1, "ein einzelnes dict ist eine Folie")
    fl, _ = _slides_normalisieren([{"title": "X"}, "Kurzer Titel"])
    pruefe(len(fl) == 2 and fl[1].get("title") == "Kurzer Titel",
           "String-Element bleibt eine Folie mit Titel")
    lang = "Zeile\nZeile2 mit viel mehr Text " * 5
    fg, _ = _slides_normalisieren([lang])
    pruefe(fg[0].get("content") == lang.strip(),
           "mehrzeiliger String wird Fliesstext, nicht Titel")

    # ── b) Deckel gegen ein Riesendeck ──────────────────────────────────────
    viele = [{"title": f"F{i}"} for i in range(MAX_FOLIEN + 25)]
    r_viel = asyncio.run(CreatePowerPointTool().execute(
        filename="test_deckel", slides=viele))
    p_viel = _erzeugte_datei(r_viel)
    pruefe(len(Presentation(str(p_viel)).slides) == MAX_FOLIEN,
           f"Folienzahl ist auf {MAX_FOLIEN} gedeckelt")
    pruefe("25" in r_viel and ("nicht" in r_viel.lower() or "weitere" in r_viel.lower()),
           "…und die weggelassenen Folien werden GENANNT (kein stiller Schnitt)")
    p_viel.unlink()

    # ── c) Unaufgeloeste Shell-Substitution ─────────────────────────────────
    from datetime import date
    heute = date.today().strftime("%d.%m.%Y")
    pruefe(_text_bereinigen("Stand: $(date +%d.%m.%Y) | 1.216 Tickets")
           == f"Stand: {heute} | 1.216 Tickets",
           "$(date …) wird durch das heutige Datum ersetzt")
    pruefe(_text_bereinigen("Kosten: 5 $(netto)") == "Kosten: 5 $(netto)",
           "andere $()-Ausdruecke bleiben unangetastet")
    pruefe(_text_bereinigen("") == "" and _text_bereinigen("normal") == "normal",
           "Text ohne Substitution bleibt gleich")

    # ── d) Zahlen aus Modelltext ────────────────────────────────────────────
    for roh, erwartet in [("1.216", 1216.0), ("923", 923.0), ("76 %", 76.0),
                          ("1.234,50", 1234.5), ("1,5", 1.5), ("(1.216)", -1216.0),
                          (45, 45.0), ("", None), ("keine Zahl", None)]:
        pruefe(_zahl(roh) == erwartet, f"_zahl({roh!r}) -> {_zahl(roh)}", f"erwartet {erwartet}")

    # ── e) Diagramme ────────────────────────────────────────────────────────
    r_ch = asyncio.run(CreatePowerPointTool().execute(
        filename="test_chart", slides=[
            {"title": "Saeulen", "chart": {"typ": "saeulen",
             "kategorien": ["Blocker", "High", "Middle"], "werte": [6, 45, 233]}},
            {"title": "Kreis", "chart": {"typ": "kreis",
             "kategorien": ["A", "B"], "werte": ["923", "233"]}},
            {"title": "Neben Text", "bullets": ["Punkt"], "chart": {"typ": "balken",
             "kategorien": ["X"], "werte": [3]}},
            {"title": "Kaputt", "chart": {"typ": "saeulen", "kategorien": []}},
        ]))
    p_ch = _erzeugte_datei(r_ch)
    prs_ch = Presentation(str(p_ch))
    charts = [[sh for sh in s.shapes if sh.has_chart] for s in prs_ch.slides]
    pruefe(len(charts[0]) == 1 and len(charts[1]) == 1 and len(charts[2]) == 1,
           "Diagramm wird eingefuegt (Saeulen, Kreis, Balken)")
    pruefe(len(charts[3]) == 0, "Diagramm ohne Kategorien wird weggelassen (keine halbe Grafik)")
    pruefe(charts[0][0].chart.plots[0].categories[0] == "Blocker", "Kategorien stimmen")
    werte = list(charts[1][0].chart.plots[0].series[0].values)
    pruefe(werte == [923.0, 233.0], f"Zahlen als Text werden gelesen ({werte})")
    pruefe(charts[0][0].left == VO.RAND_LINKS and charts[0][0].top == VO.INHALT_Y,
           "Diagramm sitzt im Satzspiegel")
    pruefe(charts[2][0].left > VO.RAND_LINKS + 1000000,
           "…und rueckt zur Seite, wenn daneben Text steht")
    pruefe(charts[1][0].chart.has_legend, "Kreisdiagramm hat eine Legende")
    pruefe(not charts[0][0].chart.has_legend,
           "eine einzelne Reihe bekommt KEINE Legende (sie wiederholt nur den Titel)")
    # OOXML-Zahlenformate sind US-notiert; '#.##0' ergab im PDF '923,000'.
    xml_ch = charts[0][0].chart._chartSpace.xml
    pruefe("#,##0" in xml_ch and "#.##0" not in xml_ch,
           "Zahlenformat ist US-notiert (#,##0) – sonst steht 923,000 auf der Folie")
    p_ch.unlink()

    # ── f) Hintergrundmaterial aus einer Firmenvorlage ──────────────────────
    # Eigene Mini-Vorlage, damit der Test ohne die echte .potx laeuft.
    import struct, zlib as _zlib

    def _png(breite, hoehe):
        """Gueltiges PNG in der gewuenschten Groesse (mit Rauschen, damit es
        ueber BILD_MIN_BYTES kommt – ein einfarbiges komprimiert zu stark)."""
        zeilen = bytearray()
        for y in range(hoehe):
            zeilen.append(0)
            zeilen.extend(bytes(((x * 7 + y * 13) % 251 for x in range(breite * 3))))
        def block(typ, daten):
            return (struct.pack(">I", len(daten)) + typ + daten
                    + struct.pack(">I", _zlib.crc32(typ + daten) & 0xFFFFFFFF))
        return (b"\x89PNG\r\n\x1a\n"
                + block(b"IHDR", struct.pack(">IIBBBBB", breite, hoehe, 8, 2, 0, 0, 0))
                + block(b"IDAT", _zlib.compress(bytes(zeilen), 6))
                + block(b"IEND", b""))

    # Realistische Masse: unter BILD_MIN_BREITE gilt ein Bild als Symbol.
    voll_png, band_png = _png(1600, 900), _png(1600, 236)
    klein_png = _png(384, 384)
    pruefe(VO._bildmasse(voll_png) == (1600, 900), "PNG-Masse werden aus dem Kopf gelesen")
    pruefe(VO._bildmasse(b"kein bild") is None, "Nicht-Bild liefert keine Masse")

    import tempfile as _tf
    tmp6 = Path(_tf.mkdtemp(prefix="jvorl6_"))   # eigenes Verzeichnis: das aus
    fake = tmp6 / "firma_test.pptx"              # Abschnitt 2 ist schon abgeraeumt
    fp = Presentation()
    fs = fp.slides.add_slide(fp.slide_layouts[6])
    from io import BytesIO
    fs.shapes.add_picture(BytesIO(voll_png), 0, 0, width=914400)
    fs.shapes.add_picture(BytesIO(band_png), 0, 914400, width=914400)
    fs.shapes.add_picture(BytesIO(klein_png), 0, 1828800, width=914400)  # Symbol
    fp.save(str(fake))

    bilder = VO.design_bilder(fake)
    pruefe(VO._bildmasse(bilder.get("vollbild", b"")) == (1600, 900),
           "Vollbild wird ueber das Seitenverhaeltnis 16:9 erkannt")
    pruefe(VO._bildmasse(bilder.get("band", b"")) == (1600, 236),
           "Zierband wird ueber das breite Verhaeltnis erkannt")
    pruefe(VO._bildmasse(bilder.get("vollbild", b"")) != (384, 384),
           "ein Symbol wird NICHT als Hintergrund genommen")
    pruefe(VO.design_bilder(tmp6 / "gibtsnicht.pptx") == {},
           "fehlende Firmenvorlage -> keine Bilder, kein Fehler")

    ziel_bg = tmp6 / "mit_bg.pptx"
    VO.erzeuge(ziel_bg)
    prs_bg = Presentation(str(ziel_bg))
    VO._hintergruende(prs_bg, bilder)
    def namen(lname):
        return [s.name for s in
                [l for l in prs_bg.slide_layouts if l.name == lname][0].shapes]
    pruefe("Hintergrund" in namen("Titelfolie"), "Titelfolie bekommt das Vollbild")
    pruefe("Hintergrund" in namen("Abschnitt"), "Abschnittsfolie bekommt das Vollbild")
    pruefe("Zierband" in namen("Titel und Inhalt"), "Inhaltsfolie bekommt das Zierband")
    pruefe("Hintergrund" not in namen("Titel und Inhalt"),
           "…aber KEIN Vollbild hinter der Aufzaehlung (unlesbar)")
    pruefe(namen("Titelfolie").index("Hintergrund") == 0,
           "Hintergrund liegt ganz hinten (deckt den Text nicht zu)")
    lay_i = [l for l in prs_bg.slide_layouts if l.name == "Titel und Inhalt"][0]
    band_shape = [s for s in lay_i.shapes if s.name == "Zierband"][0]
    inhalt_unten = VO.INHALT_Y + VO.INHALT_H
    pruefe(band_shape.top >= inhalt_unten - 1000,
           f"Zierband beginnt UNTER dem Inhaltsbereich ({band_shape.top} >= {inhalt_unten})")

    import shutil as _sh
    _sh.rmtree(tmp6, ignore_errors=True)

# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 7. Vorfall 2026-09-01: python-pptx-Weg ohne Hausvorlage ===")

# WAS PASSIERT IST (Lauf 17882583414110010 auf ECHT, 12:25): die Aufgabe
# verlangte "ein echtes Schaubild mit Kaesten und Verbindungspfeilen". Das
# Modell hat sich REGELKONFORM verhalten – Punkt 16 des System-Prompts sagt
# fuer Formen/Connectors ausdruecklich "MUSST du python-pptx via
# shell_execute verwenden" – und dabei `Presentation()` OHNE Argument
# geoeffnet. Gemessen an der ausgelieferten Datei: accent1 4F81BD (statt
# B80F2E), Calibri, 0 Medien, 11 englische Default-Layouts.
#
# Die Luecke war UNSERE, an zwei Stellen: der Prompt sagte kein Wort davon,
# dass auch der python-pptx-Weg auf der Hausvorlage aufsetzt – und
# office_template_info nannte den Dateinamen, aber NIE den Pfad. Das Modell
# hatte es im Lauf aufgerufen und wollte branden; es konnte gar nicht wissen,
# was es an Presentation(...) uebergeben soll.
#
# Der Waechter prueft die REGEL, nicht den Wortlaut: (a) wo der Prompt
# python-pptx fuer Praesentationen erlaubt, muss die Hausvorlage stehen,
# (b) `Presentation()` leer darf im Prompt nur als VERBOT vorkommen,
# (c) office_template_info muss den Pfad ausgeben – und Teil B misst, dass
# die Anleitung, die es gibt, wirklich zu einem gebrandeten Deck fuehrt.

import io
import tokenize


def _ohne_kommentare(code: str) -> str:
    """Python-Quelltext ohne Kommentare.

    Register: ein Waechter, der die eigene Begruendung im Kommentar liest,
    prueft nichts – im Projekt bereits ein Dutzend Mal bezahlt."""
    try:
        raus = [t for t in tokenize.generate_tokens(io.StringIO(code).readline)
                if t.type != tokenize.COMMENT]
        return tokenize.untokenize(raus)
    except Exception:  # noqa: BLE001
        return code


M_OK = _ohne_kommentare(M)
# Positivkontrolle: die Entfernung greift wirklich (sonst waere jede Pruefung
# darunter moeglicherweise ein Treffer im Kommentar).
pruefe("WOLLTE branden" in M and "WOLLTE branden" not in M_OK,
       "Kommentare werden vor der Pruefung entfernt (Positivkontrolle)")
pruefe("class TemplateInfoTool" in M_OK, "…und der Code bleibt dabei erhalten")

TI = M_OK.split("class TemplateInfoTool")[1].split("\nclass ")[0]
pruefe("Pfad (fuer python-pptx)" in TI,
       "office_template_info gibt den PFAD der Vorlage aus")
pruefe('Presentation(\\"{pfad}\\")' in TI or 'Presentation(\\"' in TI,
       "…und die fertige Startzeile fuer python-pptx")
pruefe("NIEMALS Presentation() ohne Argument" in TI,
       "…und warnt vor Presentation() ohne Argument")

# ── Der System-Prompt ────────────────────────────────────────────────────────
AGENT = (ROOT / "backend" / "agent.py").read_text(encoding="utf-8")
_i = AGENT.index("SYSTEM_PROMPT = ")
_j = AGENT.index("SUB_AGENT_PROMPT = ")
PROMPT = AGENT[_i:_j]
pruefe(len(PROMPT) > 5000, f"System-Prompt geschnitten ({len(PROMPT)} Zeichen)")

# (a) REGEL: jede Prompt-Zeile, die python-pptx fuer eine PRAESENTATION
#     anleitet, muss die Hausvorlage nennen. Zeilen, die python-pptx nur
#     erwaehnen (Paketliste, Abgrenzung zu matplotlib), sind ausgenommen –
#     sie leiten nicht an.
anleitend = [z for z in PROMPT.splitlines()
             if "python-pptx" in z and (".pptx" in z or "Folien" in z
                                        or "Praesentation" in z)
             and ("MUSST du" in z or "verwenden" in z or "zusammen" in z)]
pruefe(len(anleitend) >= 2, f"anleitende python-pptx-Zeilen gefunden ({len(anleitend)})")
ohne_vorlage = [z[:90] for z in anleitend
                if "hausvorlage" not in z.casefold()]
pruefe(not ohne_vorlage,
       "jede anleitende python-pptx-Zeile nennt die Hausvorlage", str(ohne_vorlage))

# (b) `Presentation()` leer darf NUR als Verbot vorkommen.
roh = [m.start() for m in re.finditer(r"Presentation\(\)", PROMPT)]
unverboten = [PROMPT[max(0, p - 60):p + 20] for p in roh
              if "NIEMALS" not in PROMPT[max(0, p - 60):p]
              and "NICHT" not in PROMPT[max(0, p - 60):p]]
pruefe(roh and not unverboten,
       f"Presentation() ohne Argument kommt nur als Verbot vor ({len(roh)}x)",
       str(unverboten))

# (c) Der Formen-Absatz nennt den Weg zum Pfad – ohne ihn ist das Verbot
#     eine Sackgasse (genau die Lage vom 01.09.).
def _absatz_ab(text: str, pos: int) -> str:
    """Vom Fund bis zum Ende SEINES Listenpunkts.

    Register: eine feste Zeichenzahl ist eine Zeitbombe – der Absatz waechst,
    und der Waechter meldet einen Fehler, den es nicht gibt (beim Ergaenzen des
    'schaubild'-Hinweises genau so passiert)."""
    ende = text.find("\n    - ", pos)
    return text[pos:ende if ende > 0 else len(text)]


_k = PROMPT.index("MUSST du python-pptx")
formen = _absatz_ab(PROMPT, _k)
pruefe("office_template_info" in formen,
       "der Formen-Absatz sagt, WOHER der Vorlagenpfad kommt")
pruefe("Presentation(" in formen and "hausvorlage" in formen.casefold(),
       "…und nennt die Startzeile samt Hausvorlage")
pruefe("slide_layouts" in formen and "NAME" in formen,
       "…und dass Layouts der Vorlage ueber den Namen gewaehlt werden")
pruefe("ACCENT_1" in formen,
       "…und dass eigene Formen Theme-Farben nehmen (folgen dem Branding)")

# (d) Das pauschale Verbot ("NICHT von Hand") darf den Formen-Absatz nicht
#     ueberdecken – sonst stehen zwei Saetze gegeneinander wie am 17.08.2026.
_p = PROMPT.index("Baue eine Praesentation NICHT von Hand")
verbot = _absatz_ab(PROMPT, _p)
pruefe("hausvorlage" in verbot.casefold(),
       "das pauschale Verbot benennt die Ausnahme (kein Widerspruch mehr)")

# ── Teil B: die Anleitung wirklich befolgen ─────────────────────────────────
if HAT_PPTX:
    import tempfile
    import asyncio as _aio
    tmp7 = Path(tempfile.mkdtemp(prefix="jvorl7_"))
    _echt = VO.VORLAGEN_DIR
    VO.VORLAGEN_DIR = tmp7
    # SANDKASTEN-WAECHTER: ein Test, der die echte data/vorlagen anfasst,
    # ueberschreibt eine womoeglich von Hand hinterlegte Firmenvorlage.
    if not str(VO.VORLAGEN_DIR).startswith(str(tmp7)):
        print("ABBRUCH: Vorlagenpfad zeigt nicht in den Sandkasten")
        sys.exit(2)
    try:
        from skills.office.main import TemplateInfoTool
        ausgabe = _aio.run(TemplateInfoTool().execute())
        pruefe("Layouts (Name" in ausgabe, "das Werkzeug antwortet normal")
        m_pfad = re.search(r"Pfad \(fuer python-pptx\): (\S+)", ausgabe)
        pruefe(bool(m_pfad), "die Ausgabe nennt einen Pfad")
        if m_pfad:
            pfad = Path(m_pfad.group(1))
            pruefe(pfad.is_absolute() and pfad.exists(),
                   f"der genannte Pfad ist absolut und existiert ({pfad.name})")
            pruefe(f'Presentation("{pfad}")' in ausgabe,
                   "die Startzeile nennt GENAU diesen Pfad (kein zweiter Wert)")
            # DIE EIGENTLICHE ZUSAGE: wer der Anleitung folgt, bekommt das
            # Hausdesign. Gemessen am Theme, nicht am Wortlaut.
            aus7 = tmp7 / "nachgebaut.pptx"
            Presentation(str(pfad)).save(str(aus7))
            def _accent1(p):
                x = zipfile.ZipFile(str(p)).read("ppt/theme/theme1.xml").decode(
                    "utf-8", "replace")
                t = re.search(r"<a:accent1>.*?</a:accent1>", x, re.S).group(0)
                return re.search(r'val="([0-9A-Fa-f]{6})"', t).group(1).upper()
            soll = VO.branding_farben()["akzent"]
            pruefe(_accent1(aus7) == soll,
                   f"ein Skript auf dieser Vorlage traegt den Hausakzent ({soll})",
                   f"gemessen {_accent1(aus7)}")
            # Gegenprobe: der Weg, den der Lauf vom 01.09. genommen hat.
            leer7 = tmp7 / "wie_am_0109.pptx"
            Presentation().save(str(leer7))
            pruefe(_accent1(leer7) == "4F81BD",
                   "Gegenprobe: Presentation() leer liefert Office-Blau 4F81BD",
                   f"gemessen {_accent1(leer7)}")
            pruefe(_accent1(leer7) != _accent1(aus7),
                   "…die Messung unterscheidet die beiden Wege ueberhaupt")
    finally:
        VO.VORLAGEN_DIR = _echt
        import shutil as _sh7
        _sh7.rmtree(tmp7, ignore_errors=True)

# ═════════════════════════════════════════════════════════════════════════════
print("\n=== 8. Schaubild: Ablaufkette im Werkzeug (Folge des 01.09.) ===")

# Die Konsequenz aus Abschnitt 7: der python-pptx-Ausweichweg war noetig, WEIL
# das Werkzeug keine Kaesten mit Pfeilen konnte. Jetzt kann es das – und der
# Ausweichweg bleibt fuer das, was wirklich darueber hinausgeht.
# Gemessen wird die ERZEUGTE DATEI (Formen, Positionen, Farbtyp), nicht der
# Quelltext: ob eine Kette entsteht, sieht man nur an den Formen.

pruefe("schaubild" in M, "das Werkzeug kennt das Feld 'schaubild'")
pruefe("SCHAUBILD_MAX" in M and "SCHAUBILD_QUER_MAX" in M,
       "Deckel und Umschaltpunkt sind benannte Konstanten")
_beschr = M.split("class CreatePowerPointTool")[1].split("def parameters_schema")[0]
pruefe("schaubild" in _beschr,
       "die Werkzeug-Beschreibung nennt es (sonst findet das Modell es nie)")

if HAT_PPTX:
    import asyncio as _a8
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.enum.shapes import MSO_SHAPE
    from skills.office.main import CreatePowerPointTool, SCHAUBILD_MAX

    def _alle(seq, bedingung) -> bool:
        """all() ueber eine LEERE Menge ist True – und damit waere jede
        Eigenschaftspruefung gruen, sobald gar nichts gezeichnet wurde. Genau
        so blieb die Gegenprobe "Schaubild nicht gezeichnet" bei sieben
        Pruefungen gruen (Register: eine Gegenprobe, die nicht beisst, ist ein
        Testmangel)."""
        seq = list(seq)
        return bool(seq) and all(bedingung(x) for x in seq)

    def _formen(pfad, folie=0, art=None):
        prs = Presentation(str(pfad))
        raus = []
        for sh in prs.slides[folie].shapes:
            try:
                if sh.auto_shape_type is not None and (art is None or sh.auto_shape_type == art):
                    raus.append(sh)
            except Exception:  # noqa: BLE001
                continue
        return prs, raus

    # Die Proben landen im ECHTEN data/documents (das Werkzeug kennt kein
    # anderes Ziel). Sie werden gesammelt und am Ende IMMER entfernt – ein
    # unlink() erst hinter den Pruefungen laesst bei jedem Fehlschlag eine
    # Datei im Bestand zurueck (beim Fahren der Gegenproben genau so passiert).
    _muell8 = []

    def _erzeuge8(slides, name):
        antwort = _a8.run(CreatePowerPointTool().execute(filename=name, slides=slides))
        m = re.search(r"/api/documents/([0-9a-f]{32}__[^)\s]+)", antwort)
        pfad = (ROOT / "data" / "documents" / m.group(1)) if m else None
        if pfad:
            _muell8.append(pfad)
        return antwort, pfad

    # ── a) Der gemeldete Fall: vier Schritte, quer ───────────────────────────
    SCHRITTE = ["Auslöser", "Prüfung", "Verarbeitung", "Benachrichtigung"]
    antw, pfad = _erzeuge8([{"title": "Ablauf",
                             "schaubild": {"schritte": SCHRITTE}}], "probe_ablauf.pptx")
    pruefe(pfad is not None and pfad.exists(), "Datei mit Schaubild entsteht", antw[:120])
    if pfad and pfad.exists():
        prs8, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        _, pfeile = _formen(pfad, 0, MSO_SHAPE.RIGHT_ARROW)
        pruefe(len(kaesten) == 4, f"vier Kaesten gezeichnet ({len(kaesten)})")
        pruefe(len(pfeile) == 3, f"drei Verbindungspfeile dazwischen ({len(pfeile)})")
        texte = " | ".join(k.text_frame.text for k in kaesten)
        pruefe(all(s in texte for s in SCHRITTE), "jeder Schritt steht in einem Kasten", texte)
        # DIE KERNZUSAGE: Theme-Farbe, kein fester RGB-Wert. Genau daran ist der
        # python-pptx-Weg vom 01.09. gescheitert.
        pruefe(_alle(kaesten, lambda k: k.fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1),
               "die Kaesten nehmen die THEME-Farbe ACCENT_1 (folgt dem Branding)")
        pruefe(_alle(kaesten, lambda k: k.text_frame.paragraphs[0].runs[0].font.color.theme_color
                     == MSO_THEME_COLOR.BACKGROUND_1),
               "…und der Text die Gegenfarbe der Vorlage")
        # Satzspiegel: nichts darf aus dem Inhaltsbereich laufen.
        rand_r = prs8.slide_width - VO.RAND_LINKS
        raus = [(s.left, s.top, s.width, s.height) for s in kaesten + pfeile
                if s.left < VO.RAND_LINKS - 1 or s.left + s.width > rand_r + 1
                or s.top < VO.INHALT_Y - 1
                or s.top + s.height > VO.INHALT_Y + VO.INHALT_H + 1]
        pruefe(bool(kaesten) and not raus, "alle Formen liegen im Satzspiegel", str(raus[:2]))
        folge = sorted(kaesten, key=lambda s: s.left)
        luecken = [(a.left + a.width, b.left) for a, b in zip(folge, folge[1:])
                   if a.left + a.width > b.left + 1]
        pruefe(bool(folge) and not luecken, "die Kaesten ueberlappen sich nicht", str(luecken))
        pruefe(_alle(folge, lambda k: k.height == folge[0].height and k.width == folge[0].width),
               "…und sind gleich gross")
        # Kein Schlagschatten – und zwar an BEIDEN Stellen: das leere
        # <a:effectLst/> allein liess LibreOffice den Schatten aus dem
        # <p:style>-Verweis weiterzeichnen (am PDF gesehen).
        xmls = [k._element.xml for k in kaesten + pfeile]
        pruefe(_alle(xmls, lambda x: "<a:effectLst/>" in x),
               "die Formen tragen einen leeren Effekt-Block")
        pruefe(_alle(xmls, lambda x: 'effectRef idx="0"' in x),
               "…UND der Style-Verweis zeigt auf 'kein Effekt' (sonst Schatten im PDF)")
        # Die Datei traegt weiterhin die Hausvorlage.
        x8 = zipfile.ZipFile(str(pfad)).read("ppt/theme/theme1.xml").decode("utf-8", "replace")
        a8 = re.search(r'<a:accent1>.*?val="([0-9A-Fa-f]{6})"', x8, re.S).group(1).upper()
        pruefe(a8 == VO.branding_farben()["akzent"],
               f"das Deck steht weiterhin auf der Hausvorlage ({a8})")

    # ── b) Ab sechs Schritten untereinander ──────────────────────────────────
    antw, pfad = _erzeuge8([{"title": "Lang",
                             "schaubild": {"schritte": [f"Schritt {i}" for i in range(1, 7)]}}],
                           "probe_lang.pptx")
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        _, runter = _formen(pfad, 0, MSO_SHAPE.DOWN_ARROW)
        pruefe(len(kaesten) == 6 and len(runter) == 5,
               f"sechs Kaesten, fuenf Pfeile ({len(kaesten)}/{len(runter)})")
        pruefe(bool(kaesten) and len({k.left for k in kaesten}) == 1,
               "sie stehen untereinander (gleiche linke Kante)")
        oben = sorted(k.top for k in kaesten)
        pruefe(len(oben) > 1 and all(b > a for a, b in zip(oben, oben[1:])),
               "…und in aufsteigender Reihenfolge")
        # Nicht flacher als 10:1 – sechs Kaesten ueber die volle Spaltenbreite
        # ergaben im ersten Rendering Balken von 14:1.
        # max() ueber eine leere Menge WIRFT – die Gegenprobe brach dadurch ab
        # und sah aus wie ein nicht gelaufener Test (Register).
        v = max((k.width / max(k.height, 1) for k in kaesten), default=0)
        pruefe(bool(kaesten) and v <= 10,
               f"die Kaesten bleiben kompakt (Verhaeltnis {v:.1f}:1)")

    # ── c) Deckel und Fehlerfaelle: nie stillschweigend ──────────────────────
    antw, pfad = _erzeuge8([{"title": "Zuviel",
                             "schaubild": {"schritte": [f"S{i}" for i in range(12)]}}],
                           "probe_deckel.pptx")
    pruefe("Schritten" in antw and str(SCHAUBILD_MAX) in antw,
           "zu viele Schritte werden gekappt UND das steht im Ergebnis", antw[-200:])
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        pruefe(len(kaesten) == SCHAUBILD_MAX, f"gezeichnet werden {SCHAUBILD_MAX}")

    antw, pfad = _erzeuge8([{"title": "Leer", "schaubild": {"schritte": []}}],
                           "probe_leer.pptx")
    pruefe("NICHT gezeichnet" in antw,
           "ein Schaubild ohne Schritte sagt es – statt eine leere Folie zu liefern")
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        pruefe(not kaesten, "…und zeichnet wirklich nichts")

    # ── d) Toleranz: das Modell schickt Text statt eines Objekts ─────────────
    antw, pfad = _erzeuge8([{"title": "AlsText",
                             "schaubild": "{'schritte': ['Eins', 'Zwei', 'Drei']}"}],
                           "probe_text.pptx")
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        pruefe(len(kaesten) == 3, f"'schaubild' als Text wird gelesen ({len(kaesten)})")
    antw, pfad = _erzeuge8([{"title": "Pfeiltext",
                             "schaubild": {"schritte": "Eins → Zwei → Drei"}}],
                           "probe_pfeil.pptx")
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        pruefe(len(kaesten) == 3, f"auch 'A → B → C' als Kette ({len(kaesten)})")

    # ── e) Zwei Grafiken auf einer Folie: Vorrang und Ansage ─────────────────
    antw, pfad = _erzeuge8([{"title": "Beides",
                             "chart": {"kategorien": ["A", "B"], "werte": [1, 2]},
                             "schaubild": {"schritte": ["Eins", "Zwei"]}}],
                           "probe_beides.pptx")
    pruefe("chart" in antw and "schaubild" in antw,
           "chart UND schaubild auf einer Folie wird ausdruecklich gemeldet", antw[-160:])
    if pfad and pfad.exists():
        _, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        pruefe(not kaesten, "…und es liegen nicht zwei Grafiken uebereinander")

    # ── f) Neben Text: senkrecht in der rechten Spalte ───────────────────────
    antw, pfad = _erzeuge8([{"title": "Mit Text", "bullets": ["Punkt 1", "Punkt 2"],
                             "schaubild": {"schritte": ["Eins", "Zwei", "Drei"]}}],
                           "probe_neben.pptx")
    if pfad and pfad.exists():
        prs8, kaesten = _formen(pfad, 0, MSO_SHAPE.ROUNDED_RECTANGLE)
        mitte = VO.RAND_LINKS + VO.INHALT_B // 2
        pruefe(len(kaesten) == 3 and _alle(kaesten, lambda k: k.left >= mitte - 1),
               "neben einer Aufzaehlung steht die Kette in der rechten Spalte")
        pruefe(bool(kaesten) and len({k.left for k in kaesten}) == 1,
               "…und zwangslaeufig untereinander (quer waere sie 2 cm breit)")

    for _p8 in _muell8:
        try:
            _p8.unlink()
        except Exception:  # noqa: BLE001
            pass
    pruefe(not [q for q in _muell8 if q.exists()],
           f"alle {len(_muell8)} Probedateien sind wieder abgeraeumt")

print(f"\n{'=' * 62}\nErgebnis: {_ok}/{_ok + _fail} Pruefungen bestanden")
if _fail:
    print(f"FEHLGESCHLAGEN: {_fail}")
sys.exit(1 if _fail else 0)
